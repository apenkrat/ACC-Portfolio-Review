5#!/usr/bin/env python3
"""
{REGION_LABEL} Portfolio Audit — standalone, no Excel dependency.
Pulls live data from Org62 via sf CLI, applies all SPSM/DAF rules,
writes a dated audit report to Google Drive TMT Reports folder.

Usage:  python3 ~/Documents/claude/tmt-reports/run_acc_audit.py
"""
import subprocess, json, sys, os, re
from datetime import date, datetime, timedelta
from collections import defaultdict

TODAY = date.today()
REPORT_DATE = datetime.now().strftime('%Y-%m-%d %H:%M')

# ── Output directory — always prompted, last value saved as default ───────────
_CONFIG_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), '.tmt_config')
_DEFAULT_DIR = "/Users/apenkrat/Library/CloudStorage/GoogleDrive-apenkrat@salesforce.com/My Drive/TMT Reports"

def _load_output_dir():
    if os.path.exists(_CONFIG_FILE):
        with open(_CONFIG_FILE) as _f:
            for _line in _f:
                if _line.startswith('output_dir='):
                    return _line.strip().split('=', 1)[1].strip().strip("'\"")
    return None

def _save_output_dir(path):
    lines = []
    if os.path.exists(_CONFIG_FILE):
        with open(_CONFIG_FILE) as _f:
            lines = [l for l in _f if not l.startswith('output_dir=')]
    lines.append(f'output_dir={path}\n')
    with open(_CONFIG_FILE, 'w') as _f:
        _f.writelines(lines)

# ── Slack intelligence cache (written by Claude Code, read by the script) ─────
# Run `/slack-intel` in Claude Code to refresh before generating reports.
_SLACK_CACHE = os.path.join(os.path.dirname(os.path.abspath(__file__)), f'slack_intel_{date.today().isoformat()}.json')

import argparse as _argparse

# ── CLI argument parsing ──────────────────────────────────────────────────────
_parser = _argparse.ArgumentParser(
    prog='run_acc_audit.py',
    description='ACC Portfolio Review — pull live Org62 data and generate audit reports.',
    formatter_class=_argparse.RawTextHelpFormatter,
)
_parser.add_argument('--region', '-r', metavar='REGION',
    help='Region to audit: "tmt", "cbs", or "all" (default: interactive prompt)')
_parser.add_argument('--format', '-f', metavar='FORMAT', dest='fmt',
    help='Output format: txt, docx, pptx, html, or all (default: interactive prompt)')
_parser.add_argument('--output', '-o', metavar='DIR',
    help='Output directory path (default: last-used or Google Drive)')
_parser.add_argument('--sf-alias', metavar='ALIAS', default=None,
    help='Salesforce CLI alias to use (default: org62)')
_parser.add_argument('--data-only', action='store_true', default=False,
    help='Write JSON data file(s) only; skip HTML/TXT/DOCX/PPTX generation')
_CLI = _parser.parse_args()

# ── Output directory ──────────────────────────────────────────────────────────
_last_dir = _load_output_dir() or _DEFAULT_DIR
if _CLI.output:
    OUTPUT_DIR = _CLI.output
    _save_output_dir(OUTPUT_DIR)
elif not sys.stdin.isatty():
    OUTPUT_DIR = _last_dir
else:
    print(f"\n📁  Output directory [{_last_dir}]: ", end='', flush=True)
    _entered = input().strip()
    OUTPUT_DIR = _entered if _entered else _last_dir
    _save_output_dir(OUTPUT_DIR)
    print(f"    → Saving to: {OUTPUT_DIR}\n")

# ── Output format ─────────────────────────────────────────────────────────────
_FMT_MAP = {'txt': ['txt'], 'docx': ['docx'], 'pptx': ['pptx'], 'html': ['html'],
            'all': ['txt','docx','pptx','html']}
if _CLI.data_only:
    OUTPUT_FORMATS = ['json']
elif _CLI.fmt:
    OUTPUT_FORMATS = _FMT_MAP.get(_CLI.fmt.lower(), ['txt'])
elif not sys.stdin.isatty():
    OUTPUT_FORMATS = ['txt']
else:
    print("Select output format:")
    print("  1 — TXT  (plain text, default)")
    print("  2 — DOCX (Word document)")
    print("  3 — PPTX (PowerPoint)")
    print("  4 — HTML (interactive dashboard)")
    print("  5 — ALL  (all four formats)")
    _choice = input("Enter choice [1]: ").strip() or '1'
    OUTPUT_FORMATS = {'1':['txt'],'2':['docx'],'3':['pptx'],'4':['html'],'5':['txt','docx','pptx','html']}.get(_choice, ['txt'])
print(f"  → Generating: {', '.join(f.upper() for f in OUTPUT_FORMATS)}\n")

# ── Region registry — add new regions here ────────────────────────────────────
# Each entry covers all subregions for that region. To add a new region,
# append a dict with: key, label, subregions, pipe_like, slug, meta_file.
REGION_REGISTRY = [
    {'key': 'AMER TMT', 'label': 'AMER TMT', 'subregions': ['AMER TMT - 1', 'AMER TMT - 2'], 'pipe_like': 'AMER TMT%', 'slug': 'AMER_TMT', 'meta_file': 'acc_portfolio_metadata.md'},
    {'key': 'AMER CBS', 'label': 'AMER CBS', 'subregions': ['AMER CBS - 1', 'AMER CBS - 2', 'AMER CBS - ENTR & TTH'], 'pipe_like': 'AMER CBS%', 'slug': 'AMER_CBS', 'meta_file': 'acc_portfolio_metadata.md'},
]
_REGION_SLUG_MAP = {'tmt': 'AMER TMT', 'cbs': 'AMER CBS', 'all': None}

# ── Region selection ──────────────────────────────────────────────────────────
if _CLI.region:
    _rkey = _CLI.region.lower()
    if _rkey == 'all':
        SELECTED_REGION_KEYS = [_r['key'] for _r in REGION_REGISTRY]
    elif _rkey in _REGION_SLUG_MAP:
        SELECTED_REGION_KEYS = [_REGION_SLUG_MAP[_rkey]]
    else:
        # allow passing the full key directly e.g. "AMER TMT"
        SELECTED_REGION_KEYS = [_r['key'] for _r in REGION_REGISTRY if _r['key'].lower() == _rkey] or [REGION_REGISTRY[0]['key']]
    print(f"  → Region(s): {', '.join(SELECTED_REGION_KEYS)}\n")
elif not sys.stdin.isatty():
    SELECTED_REGION_KEYS = [REGION_REGISTRY[0]['key']]
else:
    print("Select region to audit:")
    for _i, _r in enumerate(REGION_REGISTRY, 1):
        print(f"  {_i} — {_r['label']}")
    print(f"  {len(REGION_REGISTRY)+1} — ALL regions")
    _rchoice = input("Enter choice [1]: ").strip() or '1'
    if _rchoice == str(len(REGION_REGISTRY)+1):
        SELECTED_REGION_KEYS = [_r['key'] for _r in REGION_REGISTRY]
    elif _rchoice.isdigit() and 1 <= int(_rchoice) <= len(REGION_REGISTRY):
        SELECTED_REGION_KEYS = [REGION_REGISTRY[int(_rchoice)-1]['key']]
    else:
        SELECTED_REGION_KEYS = [REGION_REGISTRY[0]['key']]
    print(f"  → Region(s): {', '.join(SELECTED_REGION_KEYS)}\n")

# ── PM-asserted SWE overrides (project name fragments, lowercase) ─────────────
SWE_PM_OVERRIDES = [
    'aws - afx salesforce rca - csg',
    'disney parks - ccn - r2',
    'nvidia <> slack',
    'nvidia — slack',
]

# ── Auth ───────────────────────────────────────────────────────────────────────
_SF_ALIAS = _CLI.sf_alias or 'org62'

def get_sf():
    try:
        from simple_salesforce import Salesforce
        import re as _re

        # Get instance URL
        display = subprocess.run(
            ['sf', 'org', 'display', '--target-org', _SF_ALIAS, '--json'],
            capture_output=True, text=True, check=True
        )
        instance_url = json.loads(display.stdout)['result']['instanceUrl']

        # sf CLI ≥2.x redacts accessToken in org display — use auth show-access-token
        token_out = subprocess.run(
            ['sf', 'org', 'auth', 'show-access-token', '--target-org', _SF_ALIAS],
            capture_output=True, text=True, input='y\n'
        )
        m = _re.search(r'Access Token\s*│\s*(\S+)', token_out.stdout)
        if not m:
            raise ValueError(f"Could not parse access token. Output: {token_out.stdout[:200]}")
        access_token = m.group(1)

        return Salesforce(instance_url=instance_url, session_id=access_token)
    except Exception as e:
        print(f"Auth failed (alias: {_SF_ALIAS}): {e}")
        sys.exit(1)

sf = get_sf()

def soql(q, label=''):
    print(f'  {label}...')
    recs, result = [], sf.query(q)
    recs.extend(result['records'])
    while not result['done']:
        result = sf.query_more(result['nextRecordsUrl'], True)
        recs.extend(result['records'])
    def flatten(obj, prefix=''):
        row = {}
        for k, v in obj.items():
            if k == 'attributes': continue
            key = f'{prefix}.{k}' if prefix else k
            row.update(flatten(v, key) if isinstance(v, dict) else {key: v})
        return row
    rows = [flatten(r) for r in recs]
    print(f'    -> {len(rows)} rows')
    return rows

def to_f(v):
    if v is None: return None
    try: return float(str(v).replace(',', '').replace('$', '').strip())
    except (ValueError, TypeError): return None

def to_d(v):
    if not v: return None
    if isinstance(v, (date, datetime)): return v.date() if isinstance(v, datetime) else v
    for fmt in ['%Y-%m-%d', '%m/%d/%Y']:
        try: return datetime.strptime(str(v)[:10], fmt).date()
        except (ValueError, TypeError): pass
    return None

# ── Pull data ─────────────────────────────────────────────────────────────────
for _RC in [_r for _r in REGION_REGISTRY if _r['key'] in SELECTED_REGION_KEYS]:
    REGION_LABEL       = _RC['label']
    REGION_SLUG        = _RC['slug']
    REGION_SUBREGIONS  = _RC['subregions']
    REGION_PIPE_LIKE   = _RC['pipe_like']
    METADATA_FILE_NAME = _RC['meta_file']
    _SR_LIST = "'" + "','".join(REGION_SUBREGIONS) + "'"
    # Filter on the PSA Region lookup (pse__Region__r.Name) — cleaner than subregion matching
    # and naturally excludes [INACTIVE - FY25] regions which won't match 'AMER TMT%'/'AMER CBS%'.
    _SR_PROJ_CLAUSE   = f"pse__Region__r.Name LIKE '{REGION_PIPE_LIKE}'"
    _SR_CHILD_CLAUSE  = f"Project__r.pse__Region__r.Name LIKE '{REGION_PIPE_LIKE}'"
    _SR_PSE_CLAUSE    = f"pse__Project__r.pse__Region__r.Name LIKE '{REGION_PIPE_LIKE}'"
    _SR_PSE2_CLAUSE   = f"PSE_Project__r.pse__Region__r.Name LIKE '{REGION_PIPE_LIKE}'"

    if len(SELECTED_REGION_KEYS) > 1:
        print(f"\n{'='*60}")
        print(f"  Running region: {REGION_LABEL}")
        print(f"{'='*60}\n")

    print(f"Pulling live data from Org62 ({REGION_LABEL})...")

# ── Parallel SOQL fetch ───────────────────────────────────────────────────────
import concurrent.futures as _cf

_last_sat = (TODAY - timedelta(days=(TODAY.weekday() + 2) % 7)).isoformat()

_queries = {
    'q1': (f"""
SELECT Id, Name, pse__Stage__c, pse__Account__c, pse__Account__r.Name,
  pse__Project_Manager__r.Name, ProjectManager2Contact__r.Name,
  pse__Practice__r.Name, pse__Billing_Type__c, pse__Start_Date__c, pse__End_Date__c,
  pse__Bookings__c, pse__Billings__c, Total_Amount_Remaining__c,
  Revenue_Recognized_Comments__c, Forecasted_Amount_Remaining_Subreason__c,
  Forecasted_Amount_Remaining_Details__c,
  Overall_Bid_Margin_new__c, Margin_at_Close_Percent__c,
  Percent_Complete__c, Revenue_Treatment__c,
  Health_Risk_Score__c, Data_Quality_Score__c,
  pse__Opportunity__r.Owner.Name,
  pse__Opportunity__r.Salesforce_Exec_Sponsor__c,
  pse__Account__r.Owner.Name,
  pse__Unscheduled_Backlog__c, T_M_Amount_Remaining__c,
  Do_Not_Survey__c, Do_Not_Survey_Reason__c, pse_survey_send_date__c,
  pse__Opportunity__r.ContractId
FROM pse__Proj__c
WHERE pse__Stage__c IN ('In Progress', 'In Progress - SWE', 'On Hold')
  AND {_SR_PROJ_CLAUSE}
  AND pse__Practice__r.Name != 'FDE'
  AND pse__Account__c NOT IN (SELECT Id FROM Account WHERE Name LIKE '%Salesforce%')
ORDER BY pse__Account__r.Name, Name
""", 'Q1 Projects'),
    'q2': (f"""
SELECT Id, Project__c, Overall_Health__c, Trend_new__c, High_Watch_Visibility__c,
  Pulse_Update_Frequency_Required__c, LastModifiedDate,
  Scope_Status__c, Schedule_Status__c, Budget_Status__c,
  Resource_Status__c, Customer_Status__c, Governance_Status__c,
  Reason_for_RY_Path_to_Green_Scope__c, Reason_for_RY_Path_to_Green_Schedule__c,
  Reason_for_RY_Path_to_Green_Budget__c, Reason_for_RY_Path_to_Green_Resource__c,
  Reason_for_RY_Path_to_Green_Customer__c,
  Summary_Status__c, Leadership_Notes__c, Action_Needed_from_Leadership__c,
  Next_Steering_Committee_Date__c, Next_Go_Live_Date__c, SWE_or_CO_anticipated__c,
  Overall_Pulse_Status__c
FROM Project_Health_Check__c
WHERE Not_Primary_Pulse_Record__c = false
  AND Project__r.pse__Stage__c IN ('In Progress', 'In Progress - SWE', 'On Hold')
  AND {_SR_CHILD_CLAUSE}
  AND Project__r.pse__Practice__r.Name != 'FDE'
  AND Project__r.pse__Account__r.Name NOT IN ('Salesforce', 'Salesforce.com')
""", 'Q2 Pulse'),
    'q3': (f"""
SELECT Project__c
FROM Project_Health_Check__c
WHERE Project__r.pse__Stage__c IN ('In Progress', 'In Progress - SWE', 'On Hold')
  AND {_SR_CHILD_CLAUSE}
  AND Project__r.pse__Practice__r.Name != 'FDE'
  AND Project__r.pse__Account__r.Name NOT IN ('Salesforce', 'Salesforce.com')
GROUP BY Project__c
LIMIT 2000
""", 'Q3 Any Pulse'),
    'q4': (f"""
SELECT pse__Project__c, COUNT(Id), SUM(pse__Request_Billable_Amount__c), MIN(pse__Start_Date__c)
FROM pse__Resource_Request__c
WHERE pse__Status__c IN ('Draft', 'Ready to Staff', 'Tentative', 'Hold')
  AND pse__Project__r.pse__Stage__c IN ('In Progress', 'In Progress - SWE', 'On Hold')
  AND {_SR_PSE_CLAUSE}
  AND pse__Project__r.pse__Practice__r.Name != 'FDE'
GROUP BY pse__Project__c
LIMIT 2000
""", 'Q4 Open RRs'),
    'q5': (f"""
SELECT AccountId, SUM(Amount) totalAmt
FROM Opportunity
WHERE RecordType.Name IN ('Services', 'Services - LOCKED')
  AND ForecastCategory != 'Omitted'
  AND StageName NOT IN (
    '06 - Project Booked', '061 - Project on Hold',
    '09 - Project in Progress', '10 - Awaiting Acceptance',
    '11 - Implemented',
    'Dead - Duplicate', 'Dead - Lost', 'Dead - No Decision', 'Dead Attrition',
    'Closed Won', 'Closed Lost'
  )
  AND pse__Region__r.Name LIKE '{REGION_PIPE_LIKE}'
  AND Amount > 0
GROUP BY AccountId
LIMIT 2000
""", 'Q5 Pipeline'),
    'q6': (f"""
SELECT pse__Project__c,
  SUM(pse__Actual_Billable_Amount__c) actualAmt,
  SUM(Estimated_Amount__c) estimatedAmt,
  SUM(pse__Actual_Hours__c) actualHrs,
  SUM(pse__Estimated_Hours__c) estimatedHrs
FROM pse__Est_vs_Actuals__c
WHERE pse__Time_Period_Type__c = 'Week'
  AND pse__End_Date__c = {_last_sat}
  AND {_SR_PSE_CLAUSE}
  AND pse__Project__r.pse__Practice__r.Name != 'FDE'
  AND pse__Project__r.pse__Account__r.Name NOT IN ('Salesforce', 'Salesforce.com')
GROUP BY pse__Project__c
LIMIT 2000
""", 'Q6 EvA'),
    'q7a': (f"""
SELECT pse__Account__c FROM pse__Proj__c
WHERE {_SR_PROJ_CLAUSE}
  AND pse__Practice__r.Name != 'FDE'
GROUP BY pse__Account__c LIMIT 2000
""", 'Q7a Account IDs'),
    'q8': (f"""
SELECT PSE_Project__c, US_Overall_Satisfaction__c, COMPLETIONTIME__c
FROM Clicktools_Survey_Results__c
WHERE PSE_Project__c != null
  AND {_SR_PSE2_CLAUSE}
  AND Survey_Status__c = 'Complete'
ORDER BY COMPLETIONTIME__c DESC
LIMIT 2000
""", 'Q8 CSAT'),
    'q9': (f"""
SELECT pse__Project__c, Resource_Region__c,
  pse__Resource__r.Name, pse__Role__c,
  pse__Planned_Hours__c, pse__Scheduled_Hours__c,
  Total_Billable_and_Credited_Hours__c, Actual_Hours_Remaining__c,
  pse__Start_Date__c, pse__End_Date__c
FROM pse__Assignment__c
WHERE pse__Status__c IN ('Tentative', 'Scheduled')
  AND pse__Project__r.pse__Stage__c IN ('In Progress', 'In Progress - SWE', 'On Hold')
  AND {_SR_PSE_CLAUSE}
  AND pse__Project__r.pse__Practice__r.Name != 'FDE'
  AND pse__Project__r.pse__Account__r.Name NOT IN ('Salesforce', 'Salesforce.com')
LIMIT 5000
""", 'Q9 Assignments'),
}

with _cf.ThreadPoolExecutor(max_workers=9) as _pool:
    _futures = {k: _pool.submit(soql, q, lbl) for k, (q, lbl) in _queries.items()}
    _results = {k: f.result() for k, f in _futures.items()}

q1_rows  = _results['q1']
q2_rows  = _results['q2']
q3_rows  = _results['q3']
q4_rows  = _results['q4']
q5_pipe  = _results['q5']
q6_eva   = _results['q6']
q7_accts  = _results['q7a']
q9_rows   = _results['q9']

# Q5b: pipeline detail scoped to portfolio account IDs (run after Q1 so IDs are known)
_q5b_acct_ids = list({r.get('pse__Account__c') for r in q1_rows if r.get('pse__Account__c')})
_Q5B_FIELDS = """SELECT Id, Name, AccountId, StageName, Amount,
  CurrencyIsoCode, convertCurrency(Amount) convertedAmount,
  CloseDate, Owner.Name, ForecastCategory,
  Manager_Forecast_Judgement__c, Description,
  Overall_Bid_Margin__c, PushCount, CreatedDate
FROM Opportunity
WHERE RecordType.Name IN ('Services', 'Services - LOCKED')
  AND ForecastCategory != 'Omitted'
  AND StageName NOT IN (
    '06 - Project Booked', '061 - Project on Hold',
    '09 - Project in Progress', '10 - Awaiting Acceptance',
    '11 - Implemented',
    'Dead - Duplicate', 'Dead - Lost', 'Dead - No Decision', 'Dead Attrition',
    'Closed Won', 'Closed Lost'
  )
  AND Amount > 0
  AND AccountId IN ({ids})
ORDER BY AccountId, CloseDate ASC"""
_chunk = 100
q5b_pipe = []
for _i in range(0, len(_q5b_acct_ids), _chunk):
    _ids_str = "'" + "','".join(_q5b_acct_ids[_i:_i+_chunk]) + "'"
    q5b_pipe.extend(soql(_Q5B_FIELDS.format(ids=_ids_str), f'Q5b Pipeline Detail chunk {_i//_chunk+1}'))

# Q_ironclad: Ironclad Workflow URLs keyed by Contract ID
# Path: pse__Proj__c → Opportunity.ContractId → ironclad__Ironclad_Workflow__c.Contract__c
_contract_ids = list({
    p.get('pse__Opportunity__r.ContractId')
    for p in q1_rows
    if p.get('pse__Opportunity__r.ContractId')
})
ironclad_map = {}  # contract_id → url
if _contract_ids:
    _ic_chunk = 200
    for _i in range(0, len(_contract_ids), _ic_chunk):
        _ids_str = "'" + "','".join(_contract_ids[_i:_i+_ic_chunk]) + "'"
        _ic_rows = soql(
            f"SELECT Contract__c, ironclad__Workflow_Link__c "
            f"FROM ironclad__Ironclad_Workflow__c "
            f"WHERE Contract__c IN ({_ids_str}) LIMIT 2000",
            f'Q_Ironclad chunk {_i//_ic_chunk+1}'
        )
        for _r in _ic_rows:
            _cid = _r.get('Contract__c')
            _url = _r.get('ironclad__Workflow_Link__c') or ''
            if _cid and _url:
                ironclad_map[_cid] = _url

# Build CSAT score map: pid → most recent US_Overall_Satisfaction__c score
q8_csat_rows = _results['q8']
q8_csat = {}
for _row in q8_csat_rows:
    _pid = _row.get('PSE_Project__c')
    if _pid and _pid not in q8_csat:  # first = most recent (query is ORDER BY DESC)
        q8_csat[_pid] = _row.get('US_Overall_Satisfaction__c')

# ── Q7b/c: Overdue invoices (depends on q7a account IDs) ─────────────────────
overdue_map = {}
if q7_accts:
    acct_ids = list({r.get('pse__Account__c') for r in q7_accts if r.get('pse__Account__c')})
    chunk = 200
    inv_rows = []
    for i in range(0, len(acct_ids), chunk):
        ids_str = "'" + "','".join(acct_ids[i:i+chunk]) + "'"
        inv_rows += soql(f"""
SELECT Name, sfbill__BalanceDue__c
FROM sfbill__Transaction__c
WHERE sfbill__TransactionType__c = 'INV'
  AND sfbill__BalanceDue__c > 0
  AND sfbill__DueDate__c < TODAY
  AND sfbill__InvoiceAge__c > 30
  AND CurrencyIsoCode = 'USD'
  AND sfbill__AccountName__c IN ({ids_str})
LIMIT 2000
""", f'Q7b Invoices chunk {i//chunk+1}')
    if inv_rows:
        be_rows = soql(f"""
SELECT pse__Invoice_Number__c, pse__Project__c
FROM pse__Billing_Event__c
WHERE {_SR_PSE_CLAUSE}
  AND pse__Invoice_Number__c != null
LIMIT 5000
""", 'Q7c Billing Events')
        inv_to_proj = {r['pse__Invoice_Number__c']: r['pse__Project__c'] for r in be_rows if r.get('pse__Invoice_Number__c')}
        for r in inv_rows:
            pid = inv_to_proj.get(r.get('Name', ''))
            if pid:
                overdue_map.setdefault(pid, {'amount': 0, 'count': 0})
                overdue_map[pid]['amount'] += to_f(r.get('sfbill__BalanceDue__c')) or 0
                overdue_map[pid]['count'] += 1

acct_pipe_map = {}
for r in q5_pipe:
    aid = r.get('AccountId') or r.get('accountid') or ''
    amt = to_f(r.get('totalAmt') or r.get('totalamt')) or 0
    if aid:
        acct_pipe_map[aid] = amt

# SF fiscal year starts Feb 1 — quarters: Q1=Feb-Apr, Q2=May-Jul, Q3=Aug-Oct, Q4=Nov-Jan
def _sf_qtr_end(d):
    """Return last day of the SF fiscal quarter containing date d."""
    import calendar
    # Map each calendar month to its SF fiscal quarter end month
    # Q1 ends Apr, Q2 ends Jul, Q3 ends Oct, Q4 ends Jan
    _qend = {2:4, 3:4, 4:4, 5:7, 6:7, 7:7, 8:10, 9:10, 10:10, 11:1, 12:1, 1:1}
    end_month = _qend[d.month]
    end_year  = d.year if end_month != 1 else (d.year + 1 if d.month != 1 else d.year)
    last_day  = calendar.monthrange(end_year, end_month)[1]
    return date(end_year, end_month, last_day)

import calendar as _cal
_this_month_end = date(TODAY.year, TODAY.month, _cal.monthrange(TODAY.year, TODAY.month)[1])
_this_qtr_end   = _sf_qtr_end(TODAY)

# Map acct_id → list of opps closing this month / this quarter
acct_pipe_closing_m = {}   # opps closing this calendar month
acct_pipe_closing_q = {}   # opps closing this SF fiscal quarter (but not this month)
for _opp in q5b_pipe:
    _aid  = _opp.get('AccountId') or ''
    _cdt  = (_opp.get('CloseDate') or '')[:10]
    if not _aid or not _cdt:
        continue
    try:
        _cd = date.fromisoformat(_cdt)
    except ValueError:
        continue
    _amt  = to_f(_opp.get('convertedAmount') or _opp.get('Amount')) or 0
    _name = _opp.get('Name', '')
    if _cd <= _this_month_end:
        acct_pipe_closing_m.setdefault(_aid, []).append((_name, _amt, _cdt))
    elif _cd <= _this_qtr_end:
        acct_pipe_closing_q.setdefault(_aid, []).append((_name, _amt, _cdt))
eva_map = {}
for r in q6_eva:
    pid = r.get('pse__Project__c','')
    if not pid: continue
    act_amt  = to_f(r.get('actualAmt')    or r.get('actualamt'))    or 0
    est_amt  = to_f(r.get('estimatedAmt') or r.get('estimatedamt')) or 0
    act_hrs  = to_f(r.get('actualHrs')    or r.get('actualhrs'))    or 0
    est_hrs  = to_f(r.get('estimatedHrs') or r.get('estimatedhrs')) or 0
    eva_amt  = act_amt - est_amt
    eva_pct  = ((act_hrs - est_hrs) / est_hrs * 100) if est_hrs else None
    eva_map[pid] = {'eva_amt': eva_amt, 'eva_pct': eva_pct,
                    'act_amt': act_amt, 'est_amt': est_amt}

# ── Build lookup maps ─────────────────────────────────────────────────────────
pulse_map = {}
for r in q2_rows:
    pid = r.get('Project__c', '')
    if pid:
        pulse_map[pid] = r

any_pulse_pids = {r.get('Project__c') for r in q3_rows if r.get('Project__c')}

rr_map = {}
for r in q4_rows:
    pid = r.get('pse__Project__c', '')
    if not pid: continue
    cnt = to_f(r.get('expr0')) or 0
    rev = to_f(r.get('expr1')) or 0
    earliest = r.get('expr2') or ''
    if earliest: earliest = str(earliest)[:10]
    rr_map[pid] = {'count': int(cnt), 'revenue': rev, 'earliest_start': earliest}

gdc_map = {}
for r in q9_rows:
    pid = r.get('pse__Project__c', '')
    if not pid: continue
    region   = (r.get('Resource_Region__c') or '').strip()
    res_name = (r.get('pse__Resource__r.Name') or '').strip()
    role     = (r.get('pse__Role__c') or '').strip()
    est_hrs  = to_f(r.get('pse__Planned_Hours__c'))
    sch_hrs  = to_f(r.get('pse__Scheduled_Hours__c'))
    act_hrs  = to_f(r.get('Total_Billable_and_Credited_Hours__c'))
    rem_hrs  = to_f(r.get('Actual_Hours_Remaining__c'))
    start    = (r.get('pse__Start_Date__c') or '')[:10]
    end      = (r.get('pse__End_Date__c') or '')[:10]
    if pid not in gdc_map:
        gdc_map[pid] = {'total': 0, 'india': 0, 'resources': []}
    gdc_map[pid]['total'] += 1
    if region == 'GDC India':
        gdc_map[pid]['india'] += 1
    gdc_map[pid]['resources'].append({
        'name': res_name, 'role': role, 'region': region,
        'est_hrs': est_hrs, 'sch_hrs': sch_hrs, 'act_hrs': act_hrs, 'rem_hrs': rem_hrs,
        'start': start, 'end': end,
    })

# ── Helpers ───────────────────────────────────────────────────────────────────
def is_swe_co(name, stage, swe_field):
    nl = name.lower()
    if any(o in nl for o in SWE_PM_OVERRIDES): return True
    if stage == 'In Progress - SWE': return True
    if swe_field and swe_field not in ('', 'No', None): return True
    if 'swe' in nl: return True
    if 'ari' in nl: return True  # ARI projects always treated as SWE — margin excluded
    return False

def is_multiyr_cto(name, end_dt, rev_treat):
    if not end_dt: return False
    keywords = ['cto', 'sela', 'seh', 'advisory', 'term', 'rta']
    return end_dt > date(2027, 1, 31) and (
        any(k in name.lower() for k in keywords) or
        (rev_treat or '') in ('Term', 'Percent Complete')
    )

def rag(v):
    if not v: return ''
    return str(v).strip().lower().capitalize()

# ── Rule group labels ─────────────────────────────────────────────────────────
RULE_GROUPS = {
    'MARGIN_RED':        'Margin',
    'MARGIN_YELLOW':     'Margin',
    'FAR_RED_NEG':       'FAR',
    'FAR_RED_UNDERUTIL': 'FAR',
    'FAR_YELLOW':        'FAR',
    'RR_RISK':           'Resource',
    'GDC_LOW':           'Resource',
    'OVERDUE_INV':       'Invoice',
    'SWE_BURNING_HOT':   'SWE Burn',
    'NO_PULSE':          'Governance',
    'NO_STEERCO':        'Governance',
    'MISSING_PTG':       'Governance',
    'END_DATE_PAST':  'End Date',
    'END_DATE_30':    'End Date',
    'END_DATE_60':    'End Date',
    'END_DATE_90':    'End Date',
    'END_DATE_120':   'End Date',
    'CSAT_OVERDUE':      'CSAT',
    'CSAT_EXEMPT':       'CSAT',
    'PIPE_CLOSE_THIS_M':  'Pipeline',
    'PIPE_CLOSE_THIS_Q':  'Pipeline',
}

def fmt_violation(v):
    code, msg = v[0], v[1]
    grp = RULE_GROUPS.get(code, '')
    prefix = f"[{grp}] " if grp else ''
    return f"{prefix}{code}: {msg}"

# ── Evaluate ──────────────────────────────────────────────────────────────────
print("Evaluating projects...")
results = []

for p in q1_rows:
    pid      = p.get('Id', '')
    name     = p.get('Name', '') or ''
    acct_id  = p.get('pse__Account__c', '') or ''
    acct     = p.get('pse__Account__r.Name', '') or ''
    pm       = p.get('pse__Project_Manager__r.Name', '') or ''
    pm2      = p.get('ProjectManager2Contact__r.Name', '') or ''
    opp_owner   = p.get('pse__Opportunity__r.Owner.Name', '') or ''
    exec_sponsor= p.get('pse__Opportunity__r.Salesforce_Exec_Sponsor__c', '') or ''
    acct_owner  = p.get('pse__Account__r.Owner.Name', '') or ''
    stage    = p.get('pse__Stage__c', '') or ''
    practice = p.get('pse__Practice__r.Name', '') or ''
    rev_treat    = p.get('Revenue_Treatment__c', '') or ''
    billing_type = p.get('pse__Billing_Type__c', '') or ''
    bookings = to_f(p.get('pse__Bookings__c'))
    billings = to_f(p.get('pse__Billings__c'))
    far      = to_f(p.get('Total_Amount_Remaining__c'))
    unsch_backlog  = to_f(p.get('pse__Unscheduled_Backlog__c'))
    actuals_rem    = to_f(p.get('T_M_Amount_Remaining__c'))
    far_reason   = p.get('Revenue_Recognized_Comments__c') or ''
    far_details  = p.get('Forecasted_Amount_Remaining_Details__c') or ''
    far_subreason= p.get('Forecasted_Amount_Remaining_Subreason__c') or ''
    bid_margin   = to_f(p.get('Overall_Bid_Margin_new__c'))
    close_margin = to_f(p.get('Margin_at_Close_Percent__c'))
    work_pct      = to_f(p.get('Percent_Complete__c'))
    health_risk_score  = to_f(p.get('Health_Risk_Score__c'))
    data_quality_score = to_f(p.get('Data_Quality_Score__c'))
    start_dt     = to_d(p.get('pse__Start_Date__c'))
    end_dt       = to_d(p.get('pse__End_Date__c'))
    do_not_survey        = bool(p.get('Do_Not_Survey__c'))
    do_not_survey_reason = p.get('Do_Not_Survey_Reason__c') or ''
    survey_send_date     = p.get('pse_survey_send_date__c') or ''
    _contract_id         = p.get('pse__Opportunity__r.ContractId') or ''
    ironclad_url         = ironclad_map.get(_contract_id, '') if _contract_id else ''
    csat_score           = q8_csat.get(pid)

    pulse = pulse_map.get(pid, {})
    pulse_id = pulse.get('Id') or ''
    health   = pulse.get('Overall_Health__c') or 'Null'
    swe_field= pulse.get('SWE_or_CO_anticipated__c') or ''
    scope_s  = rag(pulse.get('Scope_Status__c'))
    sched_s  = rag(pulse.get('Schedule_Status__c'))
    budget_s = rag(pulse.get('Budget_Status__c'))
    resource_s = rag(pulse.get('Resource_Status__c'))
    customer_s = rag(pulse.get('Customer_Status__c'))
    ptg_scope    = pulse.get('Reason_for_RY_Path_to_Green_Scope__c') or ''
    ptg_sched    = pulse.get('Reason_for_RY_Path_to_Green_Schedule__c') or ''
    ptg_budget   = pulse.get('Reason_for_RY_Path_to_Green_Budget__c') or ''
    ptg_resource = pulse.get('Reason_for_RY_Path_to_Green_Resource__c') or ''
    ptg_customer = pulse.get('Reason_for_RY_Path_to_Green_Customer__c') or ''
    action_needed     = pulse.get('Action_Needed_from_Leadership__c') or ''
    leadership_notes  = pulse.get('Leadership_Notes__c') or ''
    overall_summary   = pulse.get('Summary_Status__c') or ''
    trend             = pulse.get('Trend_new__c') or ''
    high_watch        = (pulse.get('High_Watch_Visibility__c') or '').strip().lower() == 'include'
    last_updated   = pulse.get('LastModifiedDate') or ''
    steerco_date   = to_d(pulse.get('Next_Steering_Committee_Date__c'))
    next_golive    = to_d(pulse.get('Next_Go_Live_Date__c'))

    rr = rr_map.get(pid, {})
    rr_count   = rr.get('count', 0)
    rr_revenue = rr.get('revenue', 0)
    rr_earliest= rr.get('earliest_start', '')

    gdc_data  = gdc_map.get(pid, {})
    gdc_total = gdc_data.get('total', 0)
    gdc_india = gdc_data.get('india', 0)
    gdc_other = gdc_total - gdc_india
    gdc_pct       = (gdc_india / gdc_total) if gdc_total > 0 else None
    gdc_resources = sorted(gdc_data.get('resources', []), key=lambda x: (x['region'] != 'GDC India', x['name']))

    overdue_inv = overdue_map.get(pid, {}).get('amount', 0)
    overdue_cnt = overdue_map.get(pid, {}).get('count', 0)

    open_pipe = acct_pipe_map.get(acct_id, 0) or 0
    _eva      = eva_map.get(pid, {})
    eva_amt   = _eva.get('eva_amt')
    eva_pct   = _eva.get('eva_pct')
    swe_co    = is_swe_co(name, stage, swe_field)
    cto_exempt= is_multiyr_cto(name, end_dt, rev_treat)
    has_pulse = bool(pulse) or pid in any_pulse_pids

    baseline_ry = any(s in ('Red', 'Yellow') for s in [scope_s, sched_s, budget_s, resource_s, customer_s])

    violations = []

    # 1A: Margin degradation (SWE/CO excluded)
    if has_pulse and not swe_co and bid_margin is not None and close_margin is not None:
        delta = close_margin - bid_margin
        if delta < -5:
            violations.append(('MARGIN_RED', f'Margin delta {delta:+.1f}% (bid {bid_margin:.1f}% → close {close_margin:.1f}%)'))
        elif delta < 0:
            violations.append(('MARGIN_YELLOW', f'Margin delta {delta:+.1f}% (bid {bid_margin:.1f}% → close {close_margin:.1f}%)'))

    # 1C: FAR
    def far_context():
        lns = []
        if far_reason:    lns.append(f'  FAR Reason: "{far_reason}"')
        if far_details:   lns.append(f'  FAR Details: "{far_details}"')
        if far_subreason: lns.append(f'  FAR Subreason: "{far_subreason}"')
        return '\n'.join(lns) if lns else '  (no FAR context fields populated)'

    if has_pulse and far is not None:
        bk_val = bookings or 0
        if far < -1.0:
            violations.append(('FAR_RED_NEG', f'FAR overrun: ${far:,.0f}\n{far_context()}'))
        elif far > 0 and bk_val > 0:
            pct = (far / bk_val) * 100
            if pct > 15 and not cto_exempt:
                violations.append(('FAR_RED_UNDERUTIL', f'FAR ${far:,.0f} = {pct:.1f}% of bookings\n{far_context()}'))
            elif pct > 5:
                violations.append(('FAR_YELLOW', f'FAR ${far:,.0f} = {pct:.1f}% of bookings\n{far_context()}'))

    # 1C-RR: Pending RR revenue risk
    if rr_revenue and rr_revenue > 0:
        urgency = ''
        if far and far > 0 and rr_revenue / far > 0.10:
            urgency = ' ⚠️ ELEVATED — exceeds 10% of remaining FAR'
        violations.append(('RR_RISK', f'Pending RR Revenue: ${rr_revenue:,.0f} ({rr_count} open RRs){urgency}'))

    # 1D: GDC India resourcing threshold
    if gdc_total > 0 and gdc_pct is not None and gdc_pct <= 0.65:
        violations.append(('GDC_LOW',
            f'GDC India share {gdc_pct*100:.0f}% ({gdc_india}/{gdc_total} assigned) — below 65% threshold'))

    # 1E: Overdue invoices
    if overdue_inv > 0:
        violations.append(('OVERDUE_INV', f'Overdue invoices: ${overdue_inv:,.0f} ({overdue_cnt} invoice(s))'))

    # 1F: SWE burn rate
    swe_burn_str = None
    if swe_co and start_dt and end_dt and work_pct is not None:
        total_days = (end_dt - start_dt).days
        if total_days > 0:
            elapsed = (TODAY - start_dt).days
            time_pct = (elapsed / total_days) * 100
            if time_pct > 0:
                br = work_pct / time_pct
                tag = '🟢 NORMAL'
                if br > 1.20:   tag = '🔴 HOT'
                elif br > 1.10: tag = '🟡 WATCH'
                elif br < 0.80: tag = '🔵 SLOW'
                swe_burn_str = f'{work_pct:.0f}% complete | {time_pct:.0f}% elapsed | Ratio {br:.2f} {tag}'
                if br > 1.20:
                    violations.append(('SWE_BURNING_HOT', f'SWE burn ratio {br:.2f} — HOT'))

    # 2A: No pulse (≥$150K)
    if not has_pulse and (bookings or 0) >= 150000:
        violations.append(('NO_PULSE', f'No pulse — ${(bookings or 0):,.0f} project'))

    # 2F: SteerCo Date required for projects ≥$750K
    # Exempt projects must set date to 01/01/2100; missing or null = violation
    STEERCO_EXEMPT_DATE = date(2100, 1, 1)
    if (bookings or 0) >= 750000:
        if steerco_date is None:
            violations.append(('NO_STEERCO', f'Next Steering Committee Date is required for this ${(bookings or 0):,.0f} project. ACTION: If exempt (SEH, Advisory, etc.) set date to 01/01/2100.'))
        elif steerco_date == STEERCO_EXEMPT_DATE:
            pass  # Compliant exempt project — no violation
        # else: valid future date set — compliant

    # 3A: Watermelon
    fin_reds = {'MARGIN_RED', 'FAR_RED_NEG', 'FAR_RED_UNDERUTIL', 'OVERDUE_INV', 'SWE_BURNING_HOT'}
    has_fin_red = bool(fin_reds & {v[0] for v in violations})
    is_green = health.lower() == 'green'
    is_watermelon = is_green and (baseline_ry or has_fin_red)

    # 3B: Missing PTG
    if baseline_ry and has_pulse:
        for dim, (stat, ptg) in [
            ('Scope',    (scope_s,    ptg_scope)),
            ('Schedule', (sched_s,    ptg_sched)),
            ('Budget',   (budget_s,   ptg_budget)),
            ('Resource', (resource_s, ptg_resource)),
            ('Customer', (customer_s, ptg_customer)),
        ]:
            if stat in ('Red', 'Yellow') and not ptg:
                violations.append(('MISSING_PTG', f'Missing PTG for {dim} ({stat})'))

    # End date rules
    if end_dt:
        days_to_end = (end_dt - TODAY).days
        if days_to_end < 0:
            violations.append(('END_DATE_PAST', f'End date past due ({end_dt.isoformat()})'))
        elif days_to_end <= 30:
            violations.append(('END_DATE_30', f'End date in {days_to_end} day(s) ({end_dt.isoformat()})'))
        elif days_to_end <= 60:
            violations.append(('END_DATE_60', f'End date in {days_to_end} day(s) ({end_dt.isoformat()})'))
        elif days_to_end <= 90:
            violations.append(('END_DATE_90', f'End date in {days_to_end} day(s) ({end_dt.isoformat()})'))
        elif days_to_end <= 120:
            violations.append(('END_DATE_120', f'End date in {days_to_end} day(s) ({end_dt.isoformat()})'))

    # Pipeline closing this month / this SF fiscal quarter
    _m_opps = acct_pipe_closing_m.get(acct_id, [])
    _q_opps = acct_pipe_closing_q.get(acct_id, [])
    if _m_opps:
        _m_total = sum(a for _, a, _ in _m_opps)
        _m_names = '; '.join(f'{n[:40]} ({c})' for n, _, c in _m_opps[:3])
        if len(_m_opps) > 3: _m_names += f' + {len(_m_opps)-3} more'
        violations.append(('PIPE_CLOSE_THIS_M',
            f'{len(_m_opps)} opp(s) closing this month — ${_m_total:,.0f}: {_m_names}'))
    if _q_opps:
        _q_total = sum(a for _, a, _ in _q_opps)
        _q_names = '; '.join(f'{n[:40]} ({c})' for n, _, c in _q_opps[:3])
        if len(_q_opps) > 3: _q_names += f' + {len(_q_opps)-3} more'
        violations.append(('PIPE_CLOSE_THIS_Q',
            f'{len(_q_opps)} opp(s) closing this SF quarter — ${_q_total:,.0f}: {_q_names}'))

    # 4B: CSAT rules
    if do_not_survey:
        reason_txt = do_not_survey_reason or 'No reason provided'
        violations.append(('CSAT_EXEMPT', f'CSAT exempt: {reason_txt}'))
    elif (
        (bookings or 0) > 150_000
        and stage == 'In Progress'
        and start_dt
        and (TODAY - start_dt).days >= 90
        and not survey_send_date
    ):
        violations.append(('CSAT_OVERDUE',
            f'No CSAT survey sent; project active {(TODAY - start_dt).days} days, ${(bookings or 0):,.0f} bookings'))

    # Resource concern flag (resource baseline R/Y or open RRs)
    has_resource_concern = (resource_s in ('Red', 'Yellow')) or (rr_count > 0)

    results.append({
        'pid': pid, 'name': name, 'acct': acct, 'acct_id': acct_id, 'pm': pm, 'pm2': pm2, 'opp_owner': opp_owner, 'exec_sponsor': exec_sponsor, 'acct_owner': acct_owner,
        'stage': stage, 'health': health, 'trend': trend, 'high_watch': high_watch,
        'bookings': bookings, 'billings': billings, 'far': far,
        'far_reason': far_reason, 'far_details': far_details, 'far_subreason': far_subreason,
        'bid_margin': bid_margin, 'close_margin': close_margin,
        'scope_s': scope_s, 'sched_s': sched_s, 'budget_s': budget_s,
        'resource_s': resource_s, 'customer_s': customer_s,
        'swe_co': swe_co, 'cto_exempt': cto_exempt,
        'has_pulse': has_pulse, 'pulse_id': pulse_id, 'is_watermelon': is_watermelon,
        'violations': violations, 'baseline_ry': baseline_ry,
        'swe_burn_str': swe_burn_str, 'overall_summary': overall_summary,
        'leadership_notes': leadership_notes,
        'action_needed': action_needed, 'last_updated': last_updated,
        'steerco_date': steerco_date, 'next_golive': next_golive,
        'overdue_inv': overdue_inv, 'overdue_cnt': overdue_cnt,
        'rr_count': rr_count, 'rr_revenue': rr_revenue, 'rr_earliest': rr_earliest,
        'gdc_total': gdc_total, 'gdc_india': gdc_india, 'gdc_other': gdc_other, 'gdc_pct': gdc_pct,
        'gdc_resources': gdc_resources,
        'ptg_scope': ptg_scope, 'ptg_sched': ptg_sched, 'ptg_budget': ptg_budget,
        'ptg_resource': ptg_resource, 'ptg_customer': ptg_customer,
        'work_pct': work_pct, 'start_dt': start_dt, 'end_dt': end_dt,
        'rev_treat': rev_treat, 'billing_type': billing_type, 'practice': practice,
        'open_pipe': open_pipe,
        'unsch_backlog': unsch_backlog, 'actuals_rem': actuals_rem,
        'eva_amt': eva_amt, 'eva_pct': eva_pct,
        'health_risk_score': health_risk_score,
        'data_quality_score': data_quality_score,
        'has_resource_concern': has_resource_concern,
        'csat_score': csat_score,
        'do_not_survey': do_not_survey,
        'survey_send_date': survey_send_date,
        'ironclad_url': ironclad_url,
        'region': REGION_LABEL,
    })

    # ── Load tier + portfolio owner from metadata file ──────────────────────────
    METADATA_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), METADATA_FILE_NAME)

meta_by_pid  = {}   # project record id → {tier, owner}
meta_by_name = {}   # project name      → {tier, owner}

if os.path.exists(METADATA_FILE):
    import re as _re
    _current_tier = None
    with open(METADATA_FILE, encoding='utf-8') as _mf_read:
        for _line in _mf_read:
            _m = _re.match(r'^## (Tier \d+|Unassigned)', _line)
            if _m:
                _current_tier = _m.group(1)
                continue
            _m = _re.match(r'^\|\s*([A-Za-z0-9]{15,18})\s*\|\s*(.+?)\s*\|\s*(.+?)\s*\|\s*(.+?)\s*\|', _line)
            if _m and _current_tier:
                _pid, _name, _acct, _owner = _m.group(1), _m.group(2), _m.group(3), _m.group(4)
                _entry = {'tier': int(_current_tier.split()[-1]) if _current_tier.startswith('Tier') else 0, 'owner': _owner.strip()}
                meta_by_pid[_pid.strip()]  = _entry
                meta_by_name[_name.strip()] = _entry

def _tier_lookup(r):
    entry = meta_by_pid.get(r['pid']) or meta_by_name.get(r['name'])
    return entry['tier'] if entry else None


_unmatched = [r for r in results if not (meta_by_pid.get(r['pid']) or meta_by_name.get(r['name']))]
if _unmatched and sys.stdin.isatty():
    _meta_basename = os.path.basename(METADATA_FILE)
    print(f"\n⚠️  {len(_unmatched)} project(s) not in {_meta_basename}.")
    print(f"    Tip: add rows directly to {_meta_basename} to skip this prompt.")
    print(f"    For Portfolio Owner — press Enter to use PM2 as default.\n")
    for r in _unmatched:
        _pm2_default = r.get('pm2') or r.get('pm') or 'Unassigned'
        print(f"  [{r['pid'][:18]}] {r['name'][:60]}  (PM2: {_pm2_default})")
        _t = input("    Assign Tier (1/2/3) [3]: ").strip() or '3'
        _o = input(f"    Portfolio Owner [{_pm2_default}]: ").strip() or _pm2_default
        meta_by_pid[r['pid']] = {'tier': int(_t) if _t in ('1','2','3') else 3, 'owner': _o}
        with open(METADATA_FILE, 'a', encoding='utf-8') as _mf:
            _mf.write(f"| {r['pid']} | {r['name']} | {r['acct']} | {_o} |\n")
        print(f"    → Saved to {_meta_basename} (Tier {_t}, Owner: {_o})\n")
elif _unmatched:
    for r in _unmatched:
        _pm2_default = r.get('pm2') or r.get('pm') or 'Unassigned'
        meta_by_pid[r['pid']] = {'tier': 3, 'owner': _pm2_default}

for r in results:
    entry = meta_by_pid.get(r['pid']) or meta_by_name.get(r['name']) or {'tier': 3, 'owner': 'Unassigned'}
    r['tier']  = entry['tier']
    r['owner'] = entry['owner']

# ── Sync assignments with Page Host SQLite DB ─────────────────────────────────
# DB is the source of truth after first seed. On first run (empty table) the DB
# is seeded from metadata.md. On subsequent runs the DB overrides metadata.md.
_DB_CLIENT_ID     = os.environ.get('DB_CLIENT_ID', '')
_DB_CLIENT_SECRET = os.environ.get('DB_CLIENT_SECRET', '')
_PAGE_HOST_URL    = os.environ.get('PAGE_HOST_URL', 'https://single-html-page-app-host-07cda8a7041b.herokuapp.com')
_TILE_ID_ENV      = os.environ.get('TILE_ID', '817')

if _DB_CLIENT_ID and _DB_CLIENT_SECRET:
    import base64 as _b64, urllib.request as _ureq, urllib.error as _uerr
    import ssl as _ssl
    _db_basic = _b64.b64encode(f"{_DB_CLIENT_ID}:{_DB_CLIENT_SECRET}".encode()).decode()
    _db_write_url = f"{_PAGE_HOST_URL}/api/tiles/{_TILE_ID_ENV}/db/write"
    _db_query_url = f"{_PAGE_HOST_URL}/api/tiles/{_TILE_ID_ENV}/db/query"
    _ssl_ctx = _ssl.create_default_context()
    _ssl_ctx.check_hostname = False
    _ssl_ctx.verify_mode = _ssl.CERT_NONE

    def _db_request(url, method, body_dict, auth):
        body = json.dumps(body_dict).encode()
        req = _ureq.Request(url, data=body, method=method,
                            headers={'Content-Type': 'application/json', 'Authorization': auth})
        try:
            with _ureq.urlopen(req, timeout=15, context=_ssl_ctx) as resp:
                return json.loads(resp.read())
        except (_uerr.HTTPError, _uerr.URLError, Exception) as e:
            print(f"⚠️  DB request failed: {e}")
            return None

    try:
        # Query current DB assignments (uses session-auth REST endpoint with Basic auth as fallback)
        _bearer_token = os.environ.get('PAGE_HOST_TOKEN', '')
        _db_rows = None
        if _bearer_token:
            _qresp = _db_request(_db_query_url, 'POST',
                                 {'sql': 'SELECT pid, tier, po FROM assignments'},
                                 f'Bearer {_bearer_token}')
            if _qresp and 'rows' in _qresp:
                _db_rows = _qresp['rows']

        if _db_rows:
            # DB has data — apply DB assignments, overriding metadata.md
            _db_map = {row['pid']: row for row in _db_rows}
            for r in results:
                if r['pid'] in _db_map:
                    db_entry = _db_map[r['pid']]
                    if db_entry.get('tier') is not None:
                        r['tier']  = int(db_entry['tier'])
                    if db_entry.get('po') is not None:
                        r['owner'] = db_entry['po']
            print(f"✅  DB assignments loaded: {len(_db_rows)} rows (overriding metadata.md)")
        else:
            # DB empty or unreachable — seed from current results
            _seed_rows = [{'pid': r['pid'], 'tier': r['tier'], 'po': r['owner']} for r in results]
            _seed_resp = _db_request(_db_write_url, 'POST', {
                'table': 'assignments',
                'mode': 'upsert',
                'key_column': 'pid',
                'columns': ['pid', 'tier', 'po'],
                'rows': _seed_rows,
            }, f'Basic {_db_basic}')
            if _seed_resp and _seed_resp.get('ok'):
                print(f"✅  DB seeded with {len(_seed_rows)} assignments from metadata.md")
            else:
                print("⚠️  DB seed failed — using metadata.md assignments")
    except Exception as _db_exc:
        print(f"⚠️  DB sync skipped: {_db_exc}")

# ── Slack Intelligence — load from cache written by Claude Code ───────────────
# To refresh: ask Claude Code to run the Slack enrichment before generating reports.
# Cache file: slack_intel_YYYY-MM-DD.json (auto-dated, one per day)
_slack_cache_data = {}
if os.path.exists(_SLACK_CACHE):
    try:
        with open(_SLACK_CACHE) as _f:
            _slack_cache_data = json.load(_f)
        print(f"💬  Slack intel loaded: {sum(1 for v in _slack_cache_data.values() if v)} projects with intel.")
    except Exception as _slack_exc:
        print(f"⚠️  Slack cache load failed ({_SLACK_CACHE}): {_slack_exc}")

for r in results:
    r['slack_intel'] = _slack_cache_data.get(r['pid'], '')

# ── Segment ───────────────────────────────────────────────────────────────────
watermelons = [r for r in results if r['is_watermelon']]
reds        = [r for r in results if r['health'].lower() == 'red' and not r['is_watermelon']]
yellows     = [r for r in results if r['health'].lower() == 'yellow']
clean_green = [r for r in results if r['health'].lower() == 'green' and not r['is_watermelon']]
no_pulse    = [r for r in results if not r['has_pulse']]
no_pulse_flagged = [r for r in results if not r['has_pulse'] and (r['bookings'] or 0) >= 150000]
on_hold     = [r for r in results if r['stage'] == 'On Hold']

total_bk  = sum(r['bookings'] or 0 for r in results)
_bk_sorted = sorted(r['bookings'] or 0 for r in results)
_n = len(_bk_sorted)
median_bk = (_bk_sorted[_n//2] if _n % 2 else (_bk_sorted[_n//2-1] + _bk_sorted[_n//2]) / 2) if _n else 0
total_bil = sum(r['billings'] or 0 for r in results)
total_backlog = total_bk - total_bil
total_far = sum(r['far'] or 0 for r in results if (r['far'] or 0) != 0)
total_far_overrun = sum(r['far'] for r in results if (r['far'] or 0) < 0)
total_rr_rev = sum(r['rr_revenue'] for r in results if r['rr_revenue'])
total_overdue = sum(r['overdue_inv'] for r in results)

margin_pool = [r for r in results if not r['swe_co'] and r['bookings'] and r['bookings'] > 0 and r['bid_margin'] is not None and r['close_margin'] is not None]
tbk = sum(r['bookings'] for r in margin_pool)
w_bid   = sum(r['bid_margin']   * r['bookings'] for r in margin_pool) / tbk if tbk else 0
w_close = sum(r['close_margin'] * r['bookings'] for r in margin_pool) / tbk if tbk else 0

hr_scores  = [r['health_risk_score']  for r in results if r.get('health_risk_score')  is not None]
dq_scores  = [r['data_quality_score'] for r in results if r.get('data_quality_score') is not None]
avg_hr = sum(hr_scores) / len(hr_scores) if hr_scores else None
avg_dq = sum(dq_scores) / len(dq_scores) if dq_scores else None

csat_scores_list = [r['csat_score'] for r in results if r.get('csat_score') is not None]
avg_csat = sum(csat_scores_list) / len(csat_scores_list) if csat_scores_list else None

# ── Format helpers ────────────────────────────────────────────────────────────
def bk(r):  return r['bookings'] or 0
def bks(r): return f"${bk(r)/1e3:,.0f}K"

def snip(text, n=300):
    t = ' '.join((text or '').split())
    return t[:n] + ('…' if len(t) > n else '')

def baselines_list(r):
    dims = [('Scope', r['scope_s']), ('Schedule', r['sched_s']), ('Budget', r['budget_s']),
            ('Resource', r['resource_s']), ('Customer', r['customer_s'])]
    return [(lbl, v) for lbl, v in dims if v] or [('All', 'Green')]

def baselines(r):
    items = baselines_list(r)
    return ' | '.join(f'{lbl}={v}' for lbl, v in items)

# ── Build report ──────────────────────────────────────────────────────────────
lines = []
L = lines.append

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — EXECUTIVE SCORECARD
# ══════════════════════════════════════════════════════════════════════════════
L("=" * 80)
L(f"{REGION_LABEL} PORTFOLIO AUDIT")
L(f"Audit Date: {REPORT_DATE}  |  Data: Live Org62  |  Auditor: SPSM/DAF")
L("=" * 80)
L("")
L("RULE KEY")
L("-" * 60)
L("  ── MARGIN RULES ──────────────────────────────────────────")
L("  MARGIN_RED          Margin delta < -5%  (bid vs. close margin)")
L("  MARGIN_YELLOW       Margin delta between 0% and -5%")
L("")
L("  ── FAR / FINANCIAL UTILIZATION RULES ─────────────────────")
L("  FAR_RED_NEG      FAR is negative — active contract overrun")
L("  FAR_RED_UNDERUTIL FAR > 15% of bookings — severe under-utilization")
L("  FAR_YELLOW       FAR between 5%-14% of bookings — under-utilization watch")
L("")
L("  ── RESOURCE RULES ─────────────────────────────────────────")
L("  RR_RISK      Pending RR revenue > $0 — staffing gap risk")
L("                  ELEVATED if pending RR > 10% of remaining FAR")
L("  SWE_BURNING_HOT          SWE burn ratio > 1.20 — scope exhaustion risk")
L("")
L("  ── INVOICE / BILLING RULES ────────────────────────────────")
L("  OVERDUE_INV  Overdue invoice balance outstanding > 30 days")
L("")
L("  ── PULSE / GOVERNANCE RULES ───────────────────────────────")
L("  NO_PULSE     No pulse record on project with bookings ≥ $150K")
L("  MISSING_PTG  Baseline R/Y with no Path to Green explanation")
L("")
L("  ── DATA HYGIENE RULES ─────────────────────────────────────")
L("  NO_STEERCO   Next Steering Committee Date is required for projects ≥ $750K")
L("                  ACTION: If exempt (SEH, Advisory, etc.) set date to 01/01/2100.")
L("")
L("  ── END DATE RULES ─────────────────────────────────────────")
L("  END_DATE_PAST     End date has passed — project may need extension or closure")
L("  END_DATE_120      End date within 120 days")
L("  END_DATE_90       End date within 90 days")
L("  END_DATE_60       End date within 60 days")
L("  END_DATE_30       End date within 30 days")
L("-" * 60)
L("")
L("SECTION 1 — EXECUTIVE SCORECARD")
L("-" * 60)
L(f"Total Active Projects : {len(results)}")
L(f"Portfolio Bookings    : ${total_bk/1e6:.1f}M")
L(f"Portfolio Billings    : ${total_bil/1e6:.1f}M")
L(f"Outstanding Backlog   : ${total_backlog/1e6:.1f}M")
L("")
self_reported_green = len(clean_green) + len(watermelons)
wm_pct = len(watermelons) * 100 // max(self_reported_green, 1)
red_accts = ', '.join(sorted({r['acct'][:12].upper() for r in reds}))
L(f"  🟢 Green                 : {self_reported_green}")
L(f"     ├─ Genuine            : {len(clean_green)}")
L(f"     └─ 🍉 Watermelon      : {len(watermelons)}  ({wm_pct}% misrepresented)")
L(f"  🟡 Yellow                : {len(yellows)}")
L(f"  🔴 Red                   : {len(reds)}  {red_accts}")
L(f"  ⚫ No Pulse (Null)       : {len(no_pulse)}")
L(f"  ⏸  On Hold               : {len(on_hold)}")
L("")
L(f"Weighted Bid Margin   : {w_bid:.1f}%")
L(f"Delivered Margin      : {w_close-w_bid:+.1f}%")
L(f"Margin at Close       : {w_close:.1f}%  ({len(margin_pool)} projects / ${tbk/1e6:.1f}M bookings)")
L(f"SWE/CO Projects       : {sum(1 for r in results if r['swe_co'])} projects (excluded from margin calc)")
L("")
L("  ── FINANCIAL EXPOSURE ─────────────────────────────────────")
L(f"  Total FAR (net)       : ${total_far:,.0f}")
L(f"  Total FAR Overruns    : ${total_far_overrun:,.0f}  (negative FAR = active contract overrun)")
if total_overdue > 0:
    L(f"  Total Overdue Invoices: ${total_overdue:,.0f}")
if total_rr_rev > 0:
    L(f"  Total Revenue @ Risk  : ${total_rr_rev:,.0f}  ({sum(r['rr_count'] for r in results if r['rr_revenue'])} open RRs — staffing gaps)")
L("")

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — PORTFOLIO LEADER ACTION PLAN
# ══════════════════════════════════════════════════════════════════════════════
L("=" * 80)
L("SECTION 2 — PORTFOLIO LEADER ACTION PLAN")
L("=" * 80)
L("  Top 15 priority escalations — all watermelons + reds, ranked by bookings")
L("")
L("  Top 15 Priority Escalations:")
top = sorted([r for r in results if r['is_watermelon'] or r['health'].lower() == 'red'],
             key=lambda x: -bk(x))[:15]
for r in top:
    icon  = '🍉' if r['is_watermelon'] else '🔴'
    codes = ', '.join(v[0] for v in r['violations'][:3])
    L(f"  {icon} {r['name'][:55]}  [{r['pm']}]  {bks(r)}  {codes}")

# RR Revenue at Risk table
rr_at_risk = sorted([r for r in results if r['rr_revenue'] > 0], key=lambda x: -x['rr_revenue'])
if rr_at_risk:
    L("")
    L("  Revenue at Risk — Unassigned Resource Requests")
    L("  " + "-" * 100)
    L(f"  {'Account':<18} {'Project':<42} {'PM':<22} {'RRs':>4} {'Pending Rev':>12} {'Earliest Start':<18} {'Flag'}")
    L("  " + "-" * 100)
    for r in rr_at_risk:
        today_str = TODAY.isoformat()
        elevated = '⚠️ Elevated' if r['rr_revenue'] > (r['far'] or 0) * 0.10 and (r['far'] or 0) > 0 else ''
        start = r['rr_earliest'] or 'N/A'
        overdue = start < today_str if (start and start != 'N/A') else False
        start_label = f"{start} (overdue)" if overdue else start
        acct_short = (r['acct'] or '')[:18]
        L(f"  {acct_short:<18} {r['name'][:42]:<42} {r['pm'][:22]:<22} {r['rr_count']:>4} ${r['rr_revenue']:>11,.0f} {start_label:<18} {elevated}")
    L("  " + "-" * 100)
    L(f"  {'TOTAL':<18} {'':42} {'':22} {sum(r['rr_count'] for r in rr_at_risk):>4} ${sum(r['rr_revenue'] for r in rr_at_risk):>11,.0f}")

# Negative FAR table (overruns > $1 to exclude rounding)
neg_far_list = sorted([r for r in results if (r['far'] or 0) < -1], key=lambda x: x['far'] or 0)
if neg_far_list:
    L("")
    L("  Contract Overruns — Projects with Negative FAR")
    L("  " + "-" * 110)
    L(f"  {'Account':<20} {'Project':<42} {'PM':<22} {'Tier':>4} {'FAR':>12}  {'FAR Reason'}")
    L("  " + "-" * 110)
    for r in neg_far_list:
        tier_lbl = f"T{r['tier']}"
        reason = (r['far_reason'] or r['far_details'] or '—')[:40]
        L(f"  {r['acct'][:20]:<20} {r['name'][:42]:<42} {r['pm'][:22]:<22} {tier_lbl:>4} ${r['far']:>11,.0f}  {reason}")
    L("  " + "-" * 110)
    L(f"  Total overrun: ${sum(r['far'] for r in neg_far_list):,.0f}  ({len(neg_far_list)} projects)")

# Resource concerns
resource_concern_list = sorted(
    [r for r in results if r['has_resource_concern']],
    key=lambda x: (-x['rr_revenue'], -(x['bookings'] or 0))
)
if resource_concern_list:
    L("")
    L("  Resource Concerns — Open RRs and/or Red/Yellow Resource Baseline")
    L("  " + "-" * 80)
    L(f"  {'Project':<45} {'PM':<20} {'Res Status':<12} {'Open RRs':>9} {'Pending Rev':>12}")
    L("  " + "-" * 80)
    for r in resource_concern_list:
        res_status = r['resource_s'] or 'Green'
        flag = ' ⚠️' if res_status in ('Red', 'Yellow') else ''
        L(f"  {r['name'][:45]:<45} {r['pm'][:20]:<20} {res_status+flag:<12} {r['rr_count']:>9}  ${r['rr_revenue']:>11,.0f}")
    L("  " + "-" * 80)
    L(f"  Projects with resource concerns: {len(resource_concern_list)}  |  "
      f"Red/Yellow resource baseline: {sum(1 for r in resource_concern_list if r['resource_s'] in ('Red','Yellow'))}  |  "
      f"Total open RRs: {sum(r['rr_count'] for r in resource_concern_list)}")

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — PROJECT DETAIL BY TIER
# ══════════════════════════════════════════════════════════════════════════════
L("")
L("=" * 80)
L("SECTION 3 — PROJECT DETAIL BY TIER")
L("=" * 80)

def project_snapshot(r):
    parts = []
    ov = snip(r.get('overall_summary') or '', 400)
    if ov: parts.append(ov)
    ln = snip(r.get('leadership_notes') or '', 300)
    if ln and ln != ov: parts.append(f"Leadership: {ln}")
    if r['action_needed']:
        parts.append(f"Action: {snip(r['action_needed'], 200)}")
    # Path to green for R/Y
    ptg_parts = [r.get(f'ptg_{d}','') or '' for d in ['scope','sched','budget','resource','customer']]
    ptg = ' '.join(p for p in ptg_parts if p)
    if ptg: parts.append(f"Path to Green: {snip(ptg, 300)}")
    if not parts: parts.append('No status summary on file.')
    return '  |  '.join(parts)

def health_icon(r):
    if r.get('high_watch') and r['health'].lower() not in ('red','yellow'):
        return '⚑'
    return {'red':'🔴','yellow':'🟡','green':'🟢'}.get(r['health'].lower(), '⚫')

def pulse_indicator(r):
    """Single-char pulse presence marker for TXT/DOCX/PPTX."""
    return '📋' if r.get('has_pulse') else '○'

for tier_num in [1, 2, 3]:
    tier_label = {1:'TIER 1 — Strategic Accounts (≥$7M)',
                  2:'TIER 2 — Growth Accounts ($750K–$7M)',
                  3:'TIER 3 — Volume Accounts (<$750K)'}[tier_num]
    tier_projs = [r for r in results if r['tier'] == tier_num]
    t_hw  = [r for r in tier_projs if r.get('high_watch')]
    t_red = [r for r in tier_projs if r['health'].lower()=='red' and not r['is_watermelon']]
    t_yel = [r for r in tier_projs if r['health'].lower()=='yellow']
    t_wm  = [r for r in tier_projs if r['is_watermelon']]
    t_grn = [r for r in tier_projs if r['health'].lower()=='green' and not r['is_watermelon']]
    t_np  = [r for r in tier_projs if not r['has_pulse']]
    t_bk  = sum(bk(r) for r in tier_projs)
    t_bil = sum(r['billings'] or 0 for r in tier_projs)

    L("")
    L("=" * 80)
    L(f"  {tier_label}")
    L("=" * 80)
    L(f"  Projects: {len(tier_projs)}  |  Bookings: ${t_bk/1e6:.1f}M  |  Billings: ${t_bil/1e6:.1f}M")
    L(f"  ⚑  High Watch: {len(t_hw)}  |  🔴 Red: {len(t_red)}  |  🟡 Yellow: {len(t_yel)}  "
      f"|  🍉 Watermelon: {len(t_wm)}  |  🟢 Green: {len(t_grn)}  |  ⚫ No Status: {len(t_np)}")
    L("")

    # Order: High Watch → Red → Yellow → Watermelon → Green → No Status
    groups = [
        ('⚑  HIGH WATCH', t_hw),
        ('🔴 RED',         t_red),
        ('🟡 YELLOW',      t_yel),
        ('🍉 WATERMELON',  t_wm),
        ('🟢 GREEN',       t_grn),
        ('⚫ NO STATUS',   t_np),
    ]
    for group_label, group_projs in groups:
        if not group_projs: continue
        L(f"  ── {group_label} ({len(group_projs)}) ──────────────────────────────────────")
        L(f"  {'H':<4} {'P':<3} {'Project':<45} {'Account':<22} {'PM':<22} {'PM2':<18} {'Bookings':>10}  {'Rules':<35}  {'Baselines'}")
        L("  " + "-" * 200)
        for r in sorted(group_projs, key=lambda x: -bk(x)):
            icon  = health_icon(r)
            pulse_ind = pulse_indicator(r)
            pm2   = r.get('pm2','') or ''
            codes = ', '.join(fmt_violation(v) for v in r['violations'][:4])
            bl    = baselines(r)
            def _score_icon(v): return '🟢' if v>=70 else ('🟡' if v>=30 else '🔴') if v is not None else ''
            hr_dq = ''
            if r.get('health_risk_score') is not None: hr_dq += f"  {_score_icon(r['health_risk_score'])}H&R:{r['health_risk_score']:.0f}"
            if r.get('data_quality_score') is not None: hr_dq += f"  {_score_icon(r['data_quality_score'])}DQ:{r['data_quality_score']:.0f}"
            proj_url = f"https://org62.lightning.force.com/lightning/r/pse__Proj__c/{r['pid']}/view"
            _dates = f"{r['start_dt']} → {r['end_dt']}" if r.get('start_dt') or r.get('end_dt') else ''
            _stage_prac = '  '.join(x for x in [r.get('stage',''), r.get('practice','')] if x)
            L(f"  {icon:<4} {pulse_ind:<3} {r['name'][:45]:<45} {r['acct'][:22]:<22} {r['pm'][:22]:<22} {pm2[:18]:<18} {bks(r):>10}  {codes:<35}  {bl}{hr_dq}")
            L(f"         Link: {proj_url}  {_dates}  {_stage_prac}")
            # Financials sub-line
            def _m(v): return f'${v/1e6:.2f}M' if v and abs(v)>=1e6 else (f'${v:,.0f}' if v else '—')
            fin_parts = [f"Type:{r['rev_treat']}" if r.get('rev_treat') else None,
                         f"FAR:{_m(r['far'])}", f"BidM:{r['bid_margin']:.1f}%" if r['bid_margin'] is not None else None,
                         f"Margin@Close:{r['close_margin']:.1f}%" if r['close_margin'] is not None else None,
                         f"UnschBL:{_m(r['unsch_backlog'])}" if r.get('unsch_backlog') else None,
                         f"ActRem:{_m(r['actuals_rem'])}" if r.get('actuals_rem') else None,
                         f"EvA$:{r['eva_amt']:+,.0f}" if r.get('eva_amt') is not None else None,
                         f"EvA%:{r['eva_pct']:+.1f}%" if r.get('eva_pct') is not None else None,
                         f"Pipe:{_m(r['open_pipe'])}" if r.get('open_pipe') else None]
            fin_str = '  '.join(p for p in fin_parts if p)
            if fin_str: L(f"         Financials: {fin_str}")
            snap = project_snapshot(r)
            L(f"         Snapshot: {snap}")
            if r.get('slack_intel'):
                L(f"         💬 Slack Intel: {r['slack_intel']}")
            L("")

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 4 — DATA GAPS
# ══════════════════════════════════════════════════════════════════════════════
L("")
L("=" * 80)
L("SECTION 4 — DATA GAPS")
L("=" * 80)
L("  ARMV (Rule 1D): Milestone at-risk query not in this audit run — add Q10 to next version")
L(f"  CSAT (Rule 4B): {len(csat_scores_list)} projects with scores | CSAT_OVERDUE: {sum(1 for r in results if any(v[0]=='CSAT_OVERDUE' for v in r['violations']))} | CSAT_EXEMPT: {sum(1 for r in results if any(v[0]=='CSAT_EXEMPT' for v in r['violations']))}")
L("  SteerCo (Rule 4B): Queried from pulse — check Next_Steering_Committee_Date__c field coverage")
L("  Pulse Staleness (2B): LastModifiedDate captured in pulse; calculate staleness from REPORT_DATE")
L(f"  No-Pulse total  : {len(no_pulse)} projects ({len(no_pulse_flagged)} flagged ≥$150K)")

L("")
L("=" * 80)
L(f"Audit complete — {len(results)} projects evaluated — {REPORT_DATE}")
L("=" * 80)

# ── Write report ──────────────────────────────────────────────────────────────
report_text = '\n'.join(lines)
print(report_text)

os.makedirs(OUTPUT_DIR, exist_ok=True)
FILE_STAMP = datetime.now().strftime('%Y-%m-%d_%H%M')
base_path = f"{OUTPUT_DIR}/{REGION_SLUG}_Audit_{FILE_STAMP}"

def write_txt():
    path = base_path + '.txt'
    with open(path, 'w', encoding='utf-8') as f:
        f.write(report_text)
    print(f"✅  TXT saved:  {path}")

def write_docx():
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    C_BLUE   = RGBColor(0x1F, 0x49, 0x7D)
    C_RED    = RGBColor(0xC0, 0x00, 0x00)
    C_YELLOW = RGBColor(0x7F, 0x60, 0x00)
    C_GREEN  = RGBColor(0x37, 0x86, 0x30)
    C_WM     = RGBColor(0x2E, 0x75, 0xB6)
    C_GREY   = RGBColor(0x40, 0x40, 0x40)

    HEALTH_COLOR = {
        'red':       C_RED,
        'yellow':    C_YELLOW,
        'green':     C_GREEN,
        'watermelon':C_WM,
    }

    doc = Document()
    doc.core_properties.title = f"{REGION_LABEL} Portfolio Audit {REPORT_DATE}"

    # Page layout — landscape
    for sec in doc.sections:
        sec.orientation   = 1  # WD_ORIENT.LANDSCAPE
        sec.page_width    = Inches(11)
        sec.page_height   = Inches(8.5)
        sec.left_margin   = Inches(0.75)
        sec.right_margin  = Inches(0.75)
        sec.top_margin    = Inches(0.75)
        sec.bottom_margin = Inches(0.75)

    normal = doc.styles['Normal']
    normal.font.name = 'Calibri'
    normal.font.size = Pt(9)

    # ── helpers ──────────────────────────────────────────────────────────────
    def h1(text):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = C_BLUE
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after  = Pt(2)

    def h2(text):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = C_GREY
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after  = Pt(1)

    def body(text, bold=False, color=None, indent=0):
        p = doc.add_paragraph()
        p.paragraph_format.left_indent  = Inches(indent * 0.2)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(1)
        run = p.add_run(text)
        run.bold = bold
        if color: run.font.color.rgb = color

    def spacer():
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(0)

    def shade_row(row, hex_fill):
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement('w:shd')
            shd.set(qn('w:val'), 'clear')
            shd.set(qn('w:color'), 'auto')
            shd.set(qn('w:fill'), hex_fill)
            tcPr.append(shd)

    def add_hyperlink_run(paragraph, url, text, bold=False, color=None, size=9):
        part = paragraph.part
        r_id = part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
        hyperlink = OxmlElement('w:hyperlink')
        hyperlink.set(qn('r:id'), r_id)
        run_el = OxmlElement('w:r')
        rPr = OxmlElement('w:rPr')
        # Underline + blue hyperlink style
        u = OxmlElement('w:u'); u.set(qn('w:val'), 'single'); rPr.append(u)
        c = OxmlElement('w:color')
        hex_col = '{:02X}{:02X}{:02X}'.format(*(tuple(color) if color else (0x44, 0x72, 0xC4)))
        c.set(qn('w:val'), hex_col); rPr.append(c)
        if bold:
            b = OxmlElement('w:b'); rPr.append(b)
        sz = OxmlElement('w:sz'); sz.set(qn('w:val'), str(int(size * 2))); rPr.append(sz)
        fnt = OxmlElement('w:rFonts'); fnt.set(qn('w:ascii'), 'Calibri'); fnt.set(qn('w:hAnsi'), 'Calibri'); rPr.append(fnt)
        run_el.append(rPr)
        t = OxmlElement('w:t'); t.text = text
        t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
        run_el.append(t)
        hyperlink.append(run_el)
        paragraph._p.append(hyperlink)
        return hyperlink

    def score_rag_color(v):
        if v is None: return C_GREY
        if v >= 70: return C_GREEN
        if v >= 30: return C_YELLOW
        return C_RED

    def cell_text(cell, text, bold=False, color=None, size=9, align=None):
        p = cell.paragraphs[0]
        p.clear()
        run = p.add_run(text)
        run.bold = bold
        run.font.size = Pt(size)
        run.font.name = 'Calibri'
        if color: run.font.color.rgb = color
        if align:
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            p.alignment = {'right': WD_ALIGN_PARAGRAPH.RIGHT,
                           'center': WD_ALIGN_PARAGRAPH.CENTER}.get(align)

    RAG_COLOR = {'Red': C_RED, 'Yellow': RGBColor(0xBF,0x8F,0x00), 'Green': C_GREEN}

    def cell_plain_bullets(cell, items, size=8):
        # items: list of strings — each rendered as a plain bullet line
        first = True
        for line in items:
            p = cell.paragraphs[0] if first else cell.add_paragraph()
            first = False
            p.clear()
            p.paragraph_format.space_before = Pt(0)
            p.paragraph_format.space_after  = Pt(0)
            p.paragraph_format.left_indent  = Inches(0.05)
            run = p.add_run(f'• {line}')
            run.font.size = Pt(size)
            run.font.name = 'Calibri'

    def fmt_financials(r, size=8):
        # Write financials as bulleted lines directly into a cell
        def _fmt_m(v):
            if v is None: return '—'
            return f'${v/1e6:.2f}M' if abs(v) >= 1_000_000 else f'${v:,.0f}'
        def _fmt_pct(v):
            return f'{v:.1f}%' if v is not None else '—'
        delivered = (r['close_margin'] or 0) - (r['bid_margin'] or 0) if (r['bid_margin'] is not None and r['close_margin'] is not None) else None
        lines = []
        if r.get('rev_treat'): lines.append(f"Type: {r['rev_treat']}")
        lines += [
            f"Bookings: {_fmt_m(r['bookings'])}",
            f"Billings: {_fmt_m(r['billings'])}",
            f"FAR: {_fmt_m(r['far'])}",
            f"Bid Margin: {_fmt_pct(r['bid_margin'])}",
            f"Delivered Margin: {(f'{delivered:+.1f}%') if delivered is not None else '—'}",
            f"Margin at Close: {_fmt_pct(r['close_margin'])}",
        ]
        if r.get('unsch_backlog'): lines.append(f"Unsch Backlog: {_fmt_m(r['unsch_backlog'])}")
        if r.get('actuals_rem'):   lines.append(f"Actuals Rem: {_fmt_m(r['actuals_rem'])}")
        if r.get('eva_amt') is not None: lines.append(f"EvA $: {r['eva_amt']:+,.0f}")
        if r.get('eva_pct') is not None: lines.append(f"EvA %: {r['eva_pct']:+.1f}%")
        if r.get('open_pipe'): lines.append(f"Open Pipe: {_fmt_m(r['open_pipe'])}")
        return lines

    def cell_bullets(cell, items, size=8):
        # items: list of (label, value) tuples — each rendered as a bullet line
        first = True
        for lbl, val in items:
            p = cell.paragraphs[0] if first else cell.add_paragraph()
            first = False
            p.clear()
            p.paragraph_format.space_before = Pt(0)
            p.paragraph_format.space_after  = Pt(0)
            p.paragraph_format.left_indent  = Inches(0.05)
            run = p.add_run(f'• {lbl}=')
            run.font.size = Pt(size)
            run.font.name = 'Calibri'
            run.bold = False
            val_run = p.add_run(val)
            val_run.font.size = Pt(size)
            val_run.font.name = 'Calibri'
            val_run.bold = True
            val_run.font.color.rgb = RAG_COLOR.get(val, C_GREY)

    def add_table(headers, col_widths, rows_data, header_fill='1F497D', stripe='EBF3FB', repeat_header=False):
        n_cols = len(headers)
        tbl = doc.add_table(rows=1, cols=n_cols)
        tbl.style = 'Table Grid'

        # Header row
        hdr = tbl.rows[0]
        shade_row(hdr, header_fill)
        for i, (h, w) in enumerate(zip(headers, col_widths)):
            hdr.cells[i].width = Inches(w)
            cell_text(hdr.cells[i], h, bold=True,
                      color=RGBColor(0xFF,0xFF,0xFF), size=9)

        # Repeat header on every page
        if repeat_header:
            trPr = hdr._tr.get_or_add_trPr()
            tblHeader = OxmlElement('w:tblHeader')
            trPr.append(tblHeader)

        # Data rows
        for idx, row_vals in enumerate(rows_data):
            row = tbl.add_row()
            if idx % 2 == 1 and stripe:
                shade_row(row, stripe)
            for i, val in enumerate(row_vals):
                if isinstance(val, dict) and 'plain_bullets' in val:
                    cell_plain_bullets(row.cells[i], val['plain_bullets'])
                elif isinstance(val, dict) and 'bullets' in val:
                    cell_bullets(row.cells[i], val['bullets'])
                elif isinstance(val, dict) and 'scores' in val:
                    # {'text':..., 'color':..., 'bold':..., 'scores': [...], 'url': optional}
                    p = row.cells[i].paragraphs[0]; p.clear()
                    if val.get('url'):
                        add_hyperlink_run(p, val['url'], val.get('text',''),
                                          bold=val.get('bold', False),
                                          color=val.get('color'), size=9)
                    else:
                        r0 = p.add_run(val.get('text',''))
                        r0.bold = val.get('bold', False)
                        r0.font.size = Pt(9); r0.font.name = 'Calibri'
                        if val.get('color'): r0.font.color.rgb = val['color']
                    for slbl, sval in val['scores']:
                        if sval is None: continue
                        sep = p.add_run('  ')
                        sep.font.size = Pt(7); sep.font.name = 'Calibri'
                        chip = p.add_run(f'{slbl}:{sval:.0f}')
                        chip.bold = True; chip.font.size = Pt(7); chip.font.name = 'Calibri'
                        chip.font.color.rgb = score_rag_color(sval)
                    if val.get('subtitle'):
                        sub_p = row.cells[i].add_paragraph()
                        sub_r = sub_p.add_run(val['subtitle'])
                        sub_r.font.size = Pt(7); sub_r.font.name = 'Calibri'
                        sub_r.font.color.rgb = C_GREY
                elif isinstance(val, dict):
                    cell_text(row.cells[i], val.get('text',''),
                              bold=val.get('bold', False),
                              color=val.get('color'),
                              align=val.get('align'))
                else:
                    cell_text(row.cells[i], str(val))
        spacer()
        return tbl

    def health_color(r):
        if r['is_watermelon']: return C_WM
        h = r['health'].lower()
        return HEALTH_COLOR.get(h, C_GREY)

    def health_label(r):
        if r['is_watermelon']: return '🍉 Watermelon'
        return {'red':'🔴 Red','yellow':'🟡 Yellow','green':'🟢 Green'}.get(r['health'].lower(), r['health'])

    # ── Cover / title ─────────────────────────────────────────────────────────
    p = doc.add_paragraph()
    run = p.add_run(f"{REGION_LABEL} PORTFOLIO AUDIT")
    run.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = C_BLUE
    p.paragraph_format.space_after = Pt(2)

    body(f"Audit Date: {REPORT_DATE}  |  Data: Live Org62  |  Auditor: SPSM/DAF", color=C_GREY)
    spacer()

    # ── Section 1: Executive Scorecard ───────────────────────────────────────
    h1("SECTION 1 — EXECUTIVE SCORECARD")

    h2("Portfolio Summary")
    add_table(
        headers    = ['Metric', 'Value'],
        col_widths = [3.0, 2.0],
        rows_data  = [
            ['Total Active Projects', str(len(results))],
            ['Portfolio Bookings',    f'${total_bk/1e6:.1f}M'],
            ['Portfolio Billings',    f'${total_bil/1e6:.1f}M'],
            ['Outstanding Backlog',   f'${total_backlog/1e6:.1f}M'],
            ['SWE/CO Projects',       f'{sum(1 for r in results if r["swe_co"])} (excluded from margin calc)'],
            ['Weighted Bid Margin',   f'{w_bid:.1f}%'],
            ['Delivered Margin',      f'{w_close-w_bid:+.1f}%'],
            ['Margin at Close',       f'{w_close:.1f}%'],
        ],
        stripe=None,
    )

    h2("Health Breakdown")
    add_table(
        headers    = ['Status', 'Count', 'Notes'],
        col_widths = [1.8, 0.7, 3.5],
        rows_data  = [
            [{'text':'🟢 Green',            'color':C_GREEN,  'bold':True}, str(self_reported_green), ''],
            [{'text':'🟢 Green (genuine)',  'color':C_GREEN,  'bold':False}, str(len(clean_green)),   ''],
            [{'text':'🍉 Watermelon',       'color':C_WM,     'bold':True}, str(len(watermelons)),   f'{wm_pct}% of self-reported greens'],
            [{'text':'🟡 Yellow',           'color':C_YELLOW, 'bold':True}, str(len(yellows)),       ''],
            [{'text':'🔴 Red',              'color':C_RED,    'bold':True}, str(len(reds)),          red_accts],
            [{'text':'⚫ No Pulse (Null)',  'bold':True},                   str(len(no_pulse)),      f'{len(no_pulse_flagged)} flagged ≥$150K'],
            [{'text':'⏸ On Hold',          'bold':True},                   str(len(on_hold)),       ', '.join(r['acct'][:20] for r in on_hold)],
        ],
        stripe=None,
    )

    h2("Financial Exposure")
    expo_rows = [
        ['Total FAR (net)',       f'${total_far:,.0f}'],
        ['Total FAR Overruns',    f'${total_far_overrun:,.0f}',],
    ]
    if total_overdue > 0:
        expo_rows.append(['Total Overdue Invoices', f'${total_overdue:,.0f}'])
    if total_rr_rev > 0:
        n_rr = sum(r['rr_count'] for r in results if r['rr_revenue'])
        expo_rows.append(['Total Revenue @ Risk', f'${total_rr_rev:,.0f}  ({n_rr} open RRs)'])
    add_table(
        headers    = ['Exposure Item', 'Amount'],
        col_widths = [3.0, 2.5],
        rows_data  = expo_rows,
        stripe=None,
    )

    # ── Section 2: Portfolio Leader Action Plan ───────────────────────────────
    h1("SECTION 2 — PORTFOLIO LEADER ACTION PLAN")

    h2("Top Priority Escalations (Reds → Yellows → Watermelons, by Bookings)")
    def esc_sort(x):
        h = x['health'].lower()
        if h == 'red' and not x['is_watermelon']:    return (0, -bk(x))
        if h == 'yellow' and not x['is_watermelon']: return (1, -bk(x))
        return (2, -bk(x))
    top10 = sorted([r for r in results if r['is_watermelon'] or r['health'].lower() in ('red', 'yellow')],
                   key=esc_sort)[:15]
    def esc_people(r):
        p = [f"PM: {r['pm']}"]
        if r.get('pm2'): p.append(f"PM2: {r['pm2']}")
        if r.get('opp_owner'): p.append(f"AP: {r['opp_owner']}")
        if r.get('owner') and r.get('owner') != 'Unassigned': p.append(f"PO: {r['owner']}")
        return p
    add_table(
        headers    = ['Project', 'Team Leadership', 'Financials', 'Status', 'Rule Codes'],
        col_widths = [2.2, 1.6, 1.5, 0.9, 1.5],
        rows_data  = [
            [
                {'text': r['name'][:55], 'color': health_color(r), 'bold': True},
                {'plain_bullets': esc_people(r)},
                {'plain_bullets': fmt_financials(r)},
                health_label(r),
                '\n'.join(fmt_violation(v) for v in r['violations'][:3]),
            ]
            for r in top10
        ],
        repeat_header=True,
    )

    if rr_at_risk:
        h2("Revenue at Risk — Unassigned Resource Requests")
        rr_rows = [
            [
                r['name'][:50],
                r['pm'],
                {'text': f"${r['rr_revenue']:,.0f}", 'align': 'right'},
                {'text': str(r['rr_count']),          'align': 'right'},
                r['rr_earliest'] or 'N/A',
                '⚠️' if r['rr_revenue'] > (r['far'] or 0) * 0.10 and (r['far'] or 0) > 0 else '',
            ]
            for r in rr_at_risk
        ]
        rr_rows.append([
            {'text': 'TOTAL', 'bold': True},
            '',
            {'text': f"${sum(r['rr_revenue'] for r in rr_at_risk):,.0f}", 'align': 'right', 'bold': True},
            {'text': str(sum(r['rr_count'] for r in rr_at_risk)),         'align': 'right', 'bold': True},
            '', '',
        ])
        add_table(
            headers    = ['Project', 'PM', 'Pending Rev', 'RRs', 'Earliest Start', 'Flag'],
            col_widths = [2.4, 1.5, 1.0, 0.4, 1.0, 0.4],
        rows_data  = rr_rows,
        )

    if neg_far_list:
        h2("Contract Overruns — Projects with Negative FAR")
        neg_far_rows = [
            [
                r['acct'][:25],
                r['name'][:50],
                r['pm'],
                f"T{r['tier']}",
                {'text': f"${r['far']:,.0f}", 'align': 'right', 'color': C_RED, 'bold': True},
                (r['far_reason'] or r['far_details'] or '—')[:50],
            ]
            for r in neg_far_list
        ]
        neg_far_rows.append([
            {'text': 'TOTAL', 'bold': True}, '', '', '',
            {'text': f"${sum(r['far'] for r in neg_far_list):,.0f}", 'align': 'right', 'bold': True, 'color': C_RED},
            '',
        ])
        add_table(
            headers    = ['Account', 'Project', 'PM', 'Tier', 'FAR', 'FAR Reason'],
            col_widths = [1.4, 2.3, 1.4, 0.4, 1.0, 2.2],
            rows_data  = neg_far_rows,
        )


    # ── Section 3: Project Detail by Tier ────────────────────────────────────
    h1("SECTION 3 — PROJECT DETAIL BY TIER")

    GROUP_COLOR = {
        'HIGH WATCH': C_BLUE,
        'RED':        C_RED,
        'YELLOW':     C_YELLOW,
        'WATERMELON': C_WM,
        'GREEN':      C_GREEN,
        'NO STATUS':  C_GREY,
    }

    for tier_num in [1, 2, 3]:
        tier_label = {1:'TIER 1 — Strategic Accounts (≥$7M)',
                      2:'TIER 2 — Growth Accounts ($1.5M–$7M)',
                      3:'TIER 3 — Volume Accounts (<$1.5M)'}[tier_num]
        tier_projs = [r for r in results if r['tier'] == tier_num]
        t_hw  = [r for r in tier_projs if r.get('high_watch')]
        t_red = [r for r in tier_projs if r['health'].lower()=='red' and not r['is_watermelon']]
        t_yel = [r for r in tier_projs if r['health'].lower()=='yellow']
        t_wm  = [r for r in tier_projs if r['is_watermelon']]
        t_grn = [r for r in tier_projs if r['health'].lower()=='green' and not r['is_watermelon']]
        t_np  = [r for r in tier_projs if not r['has_pulse']]
        t_bk  = sum(bk(r) for r in tier_projs)
        t_bil = sum(r['billings'] or 0 for r in tier_projs)

        h2(tier_label)
        add_table(
            headers    = ['Metric', 'Value'],
            col_widths = [2.5, 1.5],
            rows_data  = [
                ['Projects',    str(len(tier_projs))],
                ['Bookings',    f'${t_bk/1e6:.1f}M'],
                ['Billings',    f'${t_bil/1e6:.1f}M'],
                ['⚑ High Watch', str(len(t_hw))],
                ['🔴 Red',      str(len(t_red))],
                ['🟡 Yellow',   str(len(t_yel))],
                ['🍉 Watermelon',str(len(t_wm))],
                ['🟢 Green',    str(len(t_grn))],
                ['⚫ No Status',str(len(t_np))],
            ],
            stripe=None,
        )

        for group_key, group_label, group_projs in [
            ('RED',        '🔴 Red',         t_red),
            ('YELLOW',     '🟡 Yellow',      t_yel),
            ('WATERMELON', '🍉 Watermelon',  t_wm),
            ('GREEN',      '🟢 Green',       t_grn),
            ('NO STATUS',  '⚫ No Status',   t_np),
        ]:
            if not group_projs: continue
            gc = GROUP_COLOR[group_key]
            h2(f"{group_label} ({len(group_projs)})")
            rows = []
            for r in sorted(group_projs, key=lambda x: (-bool(x.get('high_watch')), -bk(x))):
                codes = '\n'.join(fmt_violation(v) for v in r['violations'][:4])
                snap  = project_snapshot(r)
                hw_prefix = '⚑ ' if r.get('high_watch') else ''
                people = [f"PM: {r['pm']}"]
                if r.get('pm2'): people.append(f"PM2: {r['pm2']}")
                if r.get('opp_owner'): people.append(f"AP: {r['opp_owner']}")
                if r.get('owner') and r.get('owner') != 'Unassigned': people.append(f"PO: {r['owner']}")
                rows.append([
                    {'text': f"{health_icon(r)} {pulse_indicator(r)}", 'color': gc, 'bold': True},
                    {'text': f"{hw_prefix}{r['name'][:48]}", 'color': gc, 'bold': True,
                     'url': f"https://org62.lightning.force.com/lightning/r/pse__Proj__c/{r['pid']}/view",
                     'scores': [('H&R', r.get('health_risk_score')), ('DQ', r.get('data_quality_score'))],
                     'subtitle': '  '.join(x for x in [r['acct'][:25], f"{r.get('start_dt','?')} → {r.get('end_dt','?')}", r.get('stage',''), r.get('practice','')] if x)},
                    {'plain_bullets': people},
                    {'plain_bullets': fmt_financials(r)},
                    codes,
                    {'bullets': baselines_list(r)},
                    snap + (f'\n💬 {r["slack_intel"]}' if r.get('slack_intel') else ''),
                ])
            add_table(
                headers    = ['H', 'Project', 'Team Leadership', 'Financials', 'Rules', 'Baselines', 'Summary'],
                col_widths = [0.2, 2.5, 1.4, 1.4, 1.0, 1.0, 3.0],
                rows_data  = rows,
                repeat_header = True,
            )

    # ── Section 3: Data Gaps (end) ────────────────────────────────────────────
    h1("SECTION 4 — DATA GAPS")
    add_table(
        headers    = ['Item', 'Status'],
        col_widths = [2.5, 5.0],
        rows_data  = [
            ['ARMV (Rule 1D)',         'Milestone at-risk query not in this audit run'],
            ['CSAT (Rule 4B)',         f"{len(csat_scores_list)} scores | CSAT_OVERDUE: {sum(1 for r in results if any(v[0]=='CSAT_OVERDUE' for v in r['violations']))} | No-Exempt: {sum(1 for r in results if any(v[0]=='CSAT_EXEMPT' for v in r['violations']))}"],
            ['SteerCo (Rule 4B)',      'Queried from pulse — check Next_Steering_Committee_Date__c coverage'],
            ['Pulse Staleness (2B)',   'LastModifiedDate captured; staleness calc not yet implemented'],
            ['No-Pulse total',         f'{len(no_pulse)} projects  ({len(no_pulse_flagged)} flagged ≥$150K)'],
        ],
        stripe=None,
    )

    path = base_path + '.docx'
    doc.save(path)
    print(f"✅  DOCX saved: {path}")

def write_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    # ── Salesforce Brand Colors ───────────────────────────────────────────────
    SF_BLUE    = RGBColor(0x00, 0x70, 0xD2)   # Salesforce primary blue
    SF_NAVY    = RGBColor(0x03, 0x2D, 0x60)   # Salesforce dark navy
    SF_CLOUD   = RGBColor(0xF4, 0xF6, 0xF9)   # Salesforce cloud grey bg
    SF_RAIN    = RGBColor(0x16, 0x32, 0x5B)   # Salesforce midnight
    C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    C_BODY     = RGBColor(0x18, 0x18, 0x18)
    C_SUBTEXT  = RGBColor(0x54, 0x65, 0x7E)
    C_RED      = RGBColor(0xBA, 0x00, 0x21)
    C_YELLOW   = RGBColor(0xB8, 0x76, 0x00)
    C_GREEN    = RGBColor(0x2E, 0x7D, 0x32)
    C_WM       = RGBColor(0x00, 0x70, 0xD2)
    C_GREY     = RGBColor(0x54, 0x65, 0x7E)
    C_LINK     = RGBColor(0x00, 0x70, 0xD2)
    FONT       = 'Salesforce Sans'   # falls back to Calibri if not installed

    HEALTH_COL = {
        'red':       C_RED,
        'yellow':    C_YELLOW,
        'green':     C_GREEN,
        'watermelon':C_WM,
    }
    STATUS_ICONS = {'red':'🔴','yellow':'🟡','green':'🟢','watermelon':'🍉','no pulse':'⚫','on hold':'⏸'}

    SLIDE_W, SLIDE_H = Inches(13.33), Inches(7.5)
    HEADER_H = Inches(0.58)
    FOOTER_H = Inches(0.28)

    prs2 = Presentation()
    prs2.slide_width  = SLIDE_W
    prs2.slide_height = SLIDE_H
    blank = prs2.slide_layouts[6]

    def _run(p, text, size=9, bold=False, color=None, italic=False, font=FONT):
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        if color: run.font.color.rgb = color
        return run

    def _para(tf, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(0); p.space_after = Pt(0)
        return p

    def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
        from pptx.util import Pt as _Pt
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.line.fill.background()
        if fill_color:
            shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
            if line_width: shape.line.width = _Pt(line_width)
        else:
            shape.line.fill.background()
        return shape

    def add_textbox(slide, left, top, width, height, word_wrap=True):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tb.text_frame.word_wrap = word_wrap
        return tb

    def add_slide_header(title_text, subtitle_text='', accent_color=SF_BLUE):
        slide = prs2.slides.add_slide(blank)
        # Background
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = SF_CLOUD
        # Top accent bar
        add_rect(slide, 0, 0, SLIDE_W, HEADER_H, fill_color=accent_color)
        # SF logo dot (circle) — left of title
        add_rect(slide, Inches(0.18), Inches(0.13), Inches(0.32), Inches(0.32), fill_color=C_WHITE)
        # Title text
        tb = add_textbox(slide, Inches(0.6), Inches(0.08), SLIDE_W - Inches(0.8), HEADER_H, word_wrap=False)
        tf = tb.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        _run(p, title_text, size=16, bold=True, color=C_WHITE)
        if subtitle_text:
            _run(p, f'  ·  {subtitle_text}', size=11, bold=False, color=RGBColor(0xCC, 0xE0, 0xFF))
        # Bottom accent line
        add_rect(slide, 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H, fill_color=SF_NAVY)
        # Footer text
        tb2 = add_textbox(slide, Inches(0.3), SLIDE_H - FOOTER_H + Inches(0.04), SLIDE_W - Inches(0.6), FOOTER_H)
        tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
        _run(p2, f'Salesforce Professional Services  ·  {REGION_LABEL} Portfolio Review  ·  {REPORT_DATE}', size=7, color=RGBColor(0xAA, 0xBB, 0xCC))
        return slide

    def body_tb(slide, left=Inches(0.3), top=HEADER_H + Inches(0.1),
                width=None, height=None):
        w = width  or (SLIDE_W - Inches(0.6))
        h = height or (SLIDE_H - HEADER_H - FOOTER_H - Inches(0.15))
        tb = add_textbox(slide, left, top, w, h)
        return tb.text_frame

    def score_col(v):
        if v is None: return C_GREY
        return C_GREEN if v >= 70 else (C_YELLOW if v >= 30 else C_RED)

    def fmt_m(v):
        if v is None: return '—'
        return f'${v/1e6:.2f}M' if abs(v) >= 1_000_000 else f'${v:,.0f}'

    def pptx_table(slide, headers, col_widths_in, rows,
                   top=None, left=Inches(0.3), height=None,
                   hdr_fill=SF_NAVY, hdr_color=C_WHITE,
                   stripe=RGBColor(0xEB, 0xF3, 0xFB),
                   row_ht=Inches(0.26), font_size=8, hdr_size=8):
        from pptx.util import Inches as _In
        from pptx.oxml.ns import qn as _qn
        from lxml import etree as _et
        n_cols = len(headers)
        n_rows = len(rows) + 1
        tbl_w  = sum(col_widths_in)
        _top   = top if top is not None else HEADER_H + Inches(0.15)
        _h     = height or (SLIDE_H - _top - FOOTER_H - Inches(0.1))
        tbl = slide.shapes.add_table(n_rows, n_cols,
                                     left, _top, _In(tbl_w), _h).table
        tbl.first_row  = True
        tbl.horz_banding = True
        # Column widths
        for ci, cw in enumerate(col_widths_in):
            tbl.columns[ci].width = _In(cw)
        # Row heights  — set uniform (can't do per-row cleanly without XML)
        for ri in range(n_rows):
            tbl.rows[ri].height = row_ht

        def _cell_text(cell, text, size=font_size, bold=False,
                       color=C_BODY, align=PP_ALIGN.LEFT, italic=False):
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            p.space_before = Pt(0); p.space_after = Pt(0)
            run = p.add_run()
            run.text = str(text) if text is not None else ''
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = FONT
            run.font.color.rgb = color

        def _cell_fill(cell, rgb):
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = _et.SubElement(tcPr, _qn('a:solidFill'))
            srgbClr   = _et.SubElement(solidFill, _qn('a:srgbClr'))
            srgbClr.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')

        # Header row
        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            _cell_fill(cell, (hdr_fill[0], hdr_fill[1], hdr_fill[2]))
            align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            _cell_text(cell, hdr, size=hdr_size, bold=True,
                       color=hdr_color, align=align)

        # Data rows
        for ri, row_spec in enumerate(rows):
            is_odd = ri % 2 == 1
            bg_rgb = (stripe[0], stripe[1], stripe[2]) if is_odd else (0xFF, 0xFF, 0xFF)
            for ci, cell_spec in enumerate(row_spec):
                cell = tbl.cell(ri + 1, ci)
                if is_odd:
                    _cell_fill(cell, bg_rgb)
                if isinstance(cell_spec, dict):
                    text  = cell_spec.get('text', '')
                    bold  = cell_spec.get('bold', False)
                    color = cell_spec.get('color', C_BODY)
                    align = cell_spec.get('align', PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
                    italic = cell_spec.get('italic', False)
                    if cell_spec.get('fill'):
                        fr = cell_spec['fill']
                        _cell_fill(cell, (fr[0], fr[1], fr[2]))
                    _cell_text(cell, text, size=font_size, bold=bold,
                               color=color, align=align, italic=italic)
                else:
                    align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                    _cell_text(cell, cell_spec, size=font_size,
                               color=C_BODY, align=align)
        return tbl

    # ── SLIDE 1: Cover ────────────────────────────────────────────────────────
    slide = add_slide_header(f"{REGION_LABEL} Portfolio Review", REPORT_DATE, SF_NAVY)

    # Big stat block — left column
    col_w = (SLIDE_W - Inches(1.0)) / 2
    tb = add_textbox(slide, Inches(0.5), HEADER_H + Inches(0.3), col_w, Inches(5.8))
    tf = tb.text_frame; first = True

    def stat_line(label, value, color=C_BODY, size=11, bold=False):
        nonlocal first
        p = _para(tf, first); first = False; p.space_after = Pt(4)
        _run(p, f'{label}  ', size=size, color=C_SUBTEXT)
        _run(p, value, size=size, bold=bold, color=color)

    _para(tf, True); first_p = tf.paragraphs[0]
    _run(first_p, f'{REGION_LABEL}', size=26, bold=True, color=SF_NAVY)
    p2 = _para(tf); p2.space_after = Pt(12)
    _run(p2, 'Portfolio Audit Report', size=14, color=SF_BLUE)

    stat_line('Total Projects',  str(len(results)),                         SF_NAVY,   13, True)
    stat_line('Bookings',        f'${total_bk/1e6:.1f}M',                  C_BODY,    11)
    stat_line('Billings',        f'${total_bil/1e6:.1f}M',                 C_BODY,    11)
    stat_line('Backlog',         f'${total_backlog/1e6:.1f}M',             SF_BLUE,   11)
    _bid_col = C_GREEN if w_bid > 13 else (C_YELLOW if w_bid >= 5 else C_RED)
    stat_line('Weighted Bid Margin',  f'{w_bid:.1f}%',                     _bid_col,  11)
    _del_col = C_GREEN if w_close - w_bid >= 0 else (C_YELLOW if w_close - w_bid >= -5 else C_RED)
    stat_line('Delivered Margin',     f'{w_close-w_bid:+.1f}%',            _del_col,  11)
    _clo_col = C_RED if w_close < 0 else (C_RED if w_close < w_bid - 5 else C_GREEN)
    stat_line('Margin at Close',      f'{w_close:.1f}%',                   _clo_col,  11)
    p3 = _para(tf); p3.space_before = Pt(8)
    _run(p3, f'Total FAR: {fmt_m(total_far)}  ·  Overruns: {fmt_m(total_far_overrun)}',
         size=10, color=C_RED if (total_far_overrun or 0) < -100000 else C_BODY)
    p4 = _para(tf)
    _run(p4, f'Rev @ Risk: {fmt_m(total_rr_rev)}  ·  Overdue Inv: {fmt_m(total_overdue)}',
         size=10, color=C_YELLOW)

    # Status pills — right column
    pill_top = HEADER_H + Inches(0.3)
    pill_x   = Inches(0.5) + col_w + Inches(0.2)
    pill_data = [
        ('🟢 Green',      len(clean_green),  C_GREEN, RGBColor(0xE8,0xF5,0xE9)),
        ('🍉 Watermelon', len(watermelons),  C_WM,    RGBColor(0xE3,0xF2,0xFD)),
        ('🟡 Yellow',     len(yellows),      C_YELLOW,RGBColor(0xFF,0xF8,0xE1)),
        ('🔴 Red',        len(reds),         C_RED,   RGBColor(0xFF,0xEB,0xEE)),
        ('⚫ No Pulse',   len(no_pulse),     C_GREY,  RGBColor(0xF5,0xF5,0xF5)),
    ]
    for label, count, fg, bg in pill_data:
        add_rect(slide, pill_x, pill_top, col_w - Inches(0.3), Inches(0.62), fill_color=bg,
                 line_color=fg, line_width=0.75)
        tb = add_textbox(slide, pill_x + Inches(0.1), pill_top + Inches(0.08),
                         col_w - Inches(0.5), Inches(0.46))
        tf2 = tb.text_frame; p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        _run(p, f'{label}   ', size=11, color=C_BODY)
        _run(p, str(count), size=18, bold=True, color=fg)
        pill_top += Inches(0.72)

    # ── SLIDE 2: Executive Scorecard ─────────────────────────────────────────
    slide = add_slide_header('Executive Scorecard', f'{len(results)} Projects', SF_BLUE)

    sc_headers   = ['Tier', 'Projects', 'Bookings', '🔴 Red', '🍉 WM', '🟡 Yellow', '🟢 Green', '⚫ No Pulse', 'Bid M%', 'Margin at Close%']
    sc_col_w     = [2.1, 0.8, 1.1, 0.65, 0.65, 0.85, 0.85, 1.0, 0.8, 1.25]
    sc_rows      = []
    for tn in [1, 2, 3]:
        tp     = [r for r in results if r['tier'] == tn]
        tname  = {1:'T1 Strategic', 2:'T2 Growth', 3:'T3 Volume'}[tn]
        t_bk   = sum(r['bookings'] or 0 for r in tp)
        t_red  = sum(1 for r in tp if r['health'].lower()=='red' and not r['is_watermelon'])
        t_wm   = sum(1 for r in tp if r['is_watermelon'])
        t_yel  = sum(1 for r in tp if r['health'].lower()=='yellow')
        t_grn  = sum(1 for r in tp if r['health'].lower()=='green' and not r['is_watermelon'])
        t_np   = sum(1 for r in tp if not r['has_pulse'])
        bid_pool = [(r['bid_margin'], r['bookings'] or 0) for r in tp if r['bid_margin'] is not None and not r['swe_co']]
        clo_pool = [(r['close_margin'], r['bookings'] or 0) for r in tp if r['close_margin'] is not None and not r['swe_co']]
        _bp_bk = sum(bk for _, bk in bid_pool)
        _cp_bk = sum(bk for _, bk in clo_pool)
        t_bid_v = sum(m * bk for m, bk in bid_pool) / _bp_bk if _bp_bk else (sum(m for m, _ in bid_pool) / len(bid_pool) if bid_pool else None)
        t_clo_v = sum(m * bk for m, bk in clo_pool) / _cp_bk if _cp_bk else (sum(m for m, _ in clo_pool) / len(clo_pool) if clo_pool else None)
        t_bid_s = f"{t_bid_v:.1f}%" if t_bid_v is not None else '—'
        t_clo_s = f"{t_clo_v:.1f}%" if t_clo_v is not None else '—'
        t_bid_c = (C_GREEN if t_bid_v > 13 else C_YELLOW if t_bid_v >= 5 else C_RED) if t_bid_v is not None else C_BODY
        t_clo_c = C_RED if (t_clo_v is not None and t_clo_v < 0) else C_BODY
        sc_rows.append([
            {'text': tname, 'bold': True, 'color': SF_NAVY, 'align': PP_ALIGN.LEFT},
            str(len(tp)), fmt_m(t_bk),
            {'text': str(t_red),  'color': C_RED    if t_red  else C_BODY},
            {'text': str(t_wm),   'color': C_WM     if t_wm   else C_BODY},
            {'text': str(t_yel),  'color': C_YELLOW if t_yel  else C_BODY},
            {'text': str(t_grn),  'color': C_GREEN  if t_grn  else C_BODY},
            {'text': str(t_np),   'color': C_GREY   if t_np   else C_BODY},
            {'text': t_bid_s, 'color': t_bid_c},
            {'text': t_clo_s, 'color': t_clo_c},
        ])
    # Totals row
    sc_rows.append([
        {'text': 'TOTAL', 'bold': True, 'color': SF_NAVY, 'align': PP_ALIGN.LEFT,
         'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        {'text': str(len(results)), 'bold': True, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        {'text': fmt_m(total_bk),   'bold': True, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        {'text': str(len(reds)),        'bold': True, 'color': C_RED    if reds        else C_BODY, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        {'text': str(len(watermelons)), 'bold': True, 'color': C_WM     if watermelons else C_BODY, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        {'text': str(len(yellows)),     'bold': True, 'color': C_YELLOW if yellows     else C_BODY, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        {'text': str(len(clean_green)), 'bold': True, 'color': C_GREEN  if clean_green else C_BODY, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        {'text': str(len(no_pulse)),    'bold': True, 'color': C_GREY   if no_pulse    else C_BODY, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        {'text': f"{w_bid:.1f}%",   'bold': True, 'fill': RGBColor(0xD9, 0xE8, 0xF7),
         'color': C_GREEN if w_bid > 13 else C_YELLOW if w_bid >= 5 else C_RED},
        {'text': f"{w_close:.1f}%", 'bold': True, 'fill': RGBColor(0xD9, 0xE8, 0xF7),
         'color': C_RED if w_close < 0 else C_BODY},
    ])
    pptx_table(slide, sc_headers, sc_col_w, sc_rows,
               top=HEADER_H + Inches(0.2), row_ht=Inches(0.48), font_size=10, hdr_size=10)

    # ── SLIDES 3+: Per-tier project tables ────────────────────────────────────
    TIER_LABELS = {1:'Tier 1 — Strategic (≥$7M)', 2:'Tier 2 — Growth ($750K–$7M)', 3:'Tier 3 — Volume (<$750K)'}
    ROWS_PER_SLIDE = 12

    for tn in [1, 2, 3]:
        tier_projs = sorted(
            [r for r in results if r['tier'] == tn],
            key=lambda x: (not x.get('high_watch'),
                           {'red':0,'yellow':1,'watermelon':2,'green':3}.get(x['health'].lower(), 4),
                           -bk(x))
        )
        if not tier_projs: continue

        pt_headers  = ['St', 'Project', 'Account', 'PM / PO', 'Dates', 'Bookings', 'FAR', 'Bid%', 'Margin at Close%', 'H&R', 'DQ', 'Rules', 'Summary']
        pt_col_w    = [0.28, 2.5, 1.55, 1.45, 1.0, 0.75, 0.75, 0.55, 0.6, 0.38, 0.38, 1.4, 2.2]

        for chunk_i, chunk_start in enumerate(range(0, len(tier_projs), ROWS_PER_SLIDE)):
            chunk = tier_projs[chunk_start:chunk_start + ROWS_PER_SLIDE]
            pg_label = f'Page {chunk_i+1}' if len(tier_projs) > ROWS_PER_SLIDE else ''
            slide = add_slide_header(TIER_LABELS[tn], pg_label, SF_NAVY)

            pt_rows = []
            for r in chunk:
                h       = r['health'].lower()
                hcol    = HEALTH_COL.get('watermelon' if r['is_watermelon'] else h, C_GREY)
                icon    = STATUS_ICONS.get('watermelon' if r['is_watermelon'] else h, '⚫')
                hw_pfx  = '⚑ ' if r.get('high_watch') else ''
                pulse   = '✓' if r.get('has_pulse') else '—'
                pm      = (r['pm'] or '')
                po      = (r.get('owner') or '')
                pm_po   = f"{pm}\n{po}" if po else pm
                dates   = f"Start: {r.get('start_dt','?')}\nEnd: {r.get('end_dt','?')}"
                far_val = r.get('far')
                far_col = C_RED if (far_val or 0) < 0 else C_BODY
                bid_v   = r['bid_margin']
                clo_v   = r['close_margin']
                bid_s   = f"{bid_v:.1f}%" if bid_v is not None else '—'
                clo_s   = f"{clo_v:.1f}%" if clo_v is not None else '—'
                bid_col = (C_GREEN if bid_v > 13 else C_YELLOW if bid_v >= 5 else C_RED) if bid_v is not None else C_BODY
                clo_col = C_RED if (clo_v is not None and clo_v < 0) else C_BODY
                hr_val  = r.get('health_risk_score')
                dq_val  = r.get('data_quality_score')
                codes   = '\n'.join(fmt_violation(v) for v in r['violations'][:3])
                ov      = snip(r.get('overall_summary',''), 140)
                slack_s = r.get('slack_intel','')
                summary = ov + (f'\n💬 {slack_s[:80]}' if slack_s else '')
                pt_rows.append([
                    {'text': icon, 'color': hcol, 'bold': True},
                    {'text': hw_pfx + r['name'], 'color': hcol, 'bold': bool(hw_pfx), 'align': PP_ALIGN.LEFT},
                    {'text': r['acct'], 'color': C_SUBTEXT, 'align': PP_ALIGN.LEFT},
                    {'text': pm_po, 'color': C_BODY, 'align': PP_ALIGN.LEFT},
                    {'text': dates, 'color': C_SUBTEXT, 'align': PP_ALIGN.LEFT},
                    {'text': bks(r), 'align': PP_ALIGN.RIGHT, 'color': C_BODY},
                    {'text': fmt_m(far_val), 'color': far_col, 'align': PP_ALIGN.RIGHT},
                    {'text': bid_s, 'color': bid_col},
                    {'text': clo_s, 'color': clo_col},
                    {'text': f'{hr_val:.0f}' if hr_val is not None else '—', 'bold': True, 'color': score_col(hr_val)},
                    {'text': f'{dq_val:.0f}' if dq_val is not None else '—', 'bold': True, 'color': score_col(dq_val)},
                    {'text': codes, 'color': C_RED if codes else C_BODY, 'align': PP_ALIGN.LEFT},
                    {'text': summary, 'color': C_SUBTEXT, 'align': PP_ALIGN.LEFT},
                ])
            row_h = min(Inches(0.52), Inches((SLIDE_H - HEADER_H - FOOTER_H - Inches(0.35)) / (len(pt_rows) + 1)))
            pptx_table(slide, pt_headers, pt_col_w, pt_rows,
                       top=HEADER_H + Inches(0.15), row_ht=row_h, font_size=7, hdr_size=7)

    # ── SLIDE: Top Escalations ────────────────────────────────────────────────
    top_esc = sorted(
        [r for r in results if r['is_watermelon'] or r['health'].lower() in ('red','yellow')],
        key=lambda x: (-bool(x.get('high_watch')),
                       {'red':0,'watermelon':1,'yellow':2}.get(x['health'].lower() if not x['is_watermelon'] else 'watermelon', 3),
                       -bk(x))
    )[:15]
    slide = add_slide_header('Priority Escalations', f'Top {len(top_esc)} by severity & bookings', C_RED)

    esc_headers = ['St', 'Project', 'Account', 'PM', 'Bookings', 'FAR', 'Rules', 'H&R', 'DQ', 'Summary']
    esc_col_w   = [0.28, 2.8, 1.6, 1.5, 0.85, 0.85, 1.7, 0.38, 0.38, 3.29]
    esc_rows    = []
    for r in top_esc:
        h      = r['health'].lower()
        hcol   = HEALTH_COL.get('watermelon' if r['is_watermelon'] else h, C_GREY)
        icon   = STATUS_ICONS.get('watermelon' if r['is_watermelon'] else h, '⚫')
        hw_pfx = '⚑ ' if r.get('high_watch') else ''
        codes  = '\n'.join(fmt_violation(v) for v in r['violations'][:4])
        hr_val = r.get('health_risk_score'); dq_val = r.get('data_quality_score')
        ov     = snip(r.get('overall_summary',''), 160)
        slack_s = r.get('slack_intel','')
        summary = ov + (f'\n💬 {slack_s[:120]}' if slack_s else '')
        far_val = r.get('far')
        esc_rows.append([
            {'text': icon, 'color': hcol, 'bold': True},
            {'text': hw_pfx + r['name'], 'color': hcol, 'bold': True, 'align': PP_ALIGN.LEFT},
            {'text': r['acct'], 'color': C_SUBTEXT, 'align': PP_ALIGN.LEFT},
            {'text': r['pm'] or '', 'color': C_BODY, 'align': PP_ALIGN.LEFT},
            {'text': bks(r), 'align': PP_ALIGN.RIGHT, 'color': C_BODY},
            {'text': fmt_m(far_val), 'align': PP_ALIGN.RIGHT,
             'color': C_RED if (far_val or 0) < 0 else C_BODY},
            {'text': codes, 'color': C_RED if codes else C_BODY, 'align': PP_ALIGN.LEFT},
            {'text': f'{hr_val:.0f}' if hr_val is not None else '—', 'bold': True, 'color': score_col(hr_val)},
            {'text': f'{dq_val:.0f}' if dq_val is not None else '—', 'bold': True, 'color': score_col(dq_val)},
            {'text': summary, 'color': C_SUBTEXT, 'align': PP_ALIGN.LEFT, 'italic': bool(slack_s)},
        ])
    row_h = min(Inches(0.55), Inches((SLIDE_H - HEADER_H - FOOTER_H - Inches(0.35)) / (len(esc_rows) + 1)))
    pptx_table(slide, esc_headers, esc_col_w, esc_rows,
               top=HEADER_H + Inches(0.15), hdr_fill=C_RED,
               row_ht=row_h, font_size=8, hdr_size=8)

    # ── SLIDE: Revenue at Risk ────────────────────────────────────────────────
    if rr_at_risk:
        slide = add_slide_header('Revenue at Risk — Open Resource Requests', '', C_YELLOW)

        rr_headers = ['Project', 'PM', 'Pending Revenue', 'RRs', 'Earliest RR', 'Resource Status']
        rr_col_w   = [3.8, 1.8, 1.6, 0.5, 1.1, 1.5]
        rr_rows    = []
        for r in rr_at_risk:
            elev  = ' ⚠' if r['rr_revenue'] > (r['far'] or 0) * 0.10 and (r['far'] or 0) > 0 else ''
            res   = r['resource_s'] or ''
            rflag = ' ⚠' if res in ('Red','Yellow') else ''
            res_col = C_RED if res=='Red' else (C_YELLOW if res=='Yellow' else C_BODY)
            rr_rows.append([
                {'text': r['name'], 'color': C_BODY, 'align': PP_ALIGN.LEFT},
                {'text': r['pm'] or '', 'color': C_SUBTEXT, 'align': PP_ALIGN.LEFT},
                {'text': f"${r['rr_revenue']:,.0f}{elev}", 'bold': True, 'color': C_RED, 'align': PP_ALIGN.RIGHT},
                {'text': str(r['rr_count']), 'align': PP_ALIGN.CENTER},
                {'text': str(r['rr_earliest'] or 'N/A'), 'color': C_BODY},
                {'text': f'{res}{rflag}', 'color': res_col, 'bold': res in ('Red','Yellow')},
            ])
        # Totals row
        rr_total = sum(r['rr_revenue'] for r in rr_at_risk)
        rr_cnt   = sum(r['rr_count']   for r in rr_at_risk)
        rr_rows.append([
            {'text': 'TOTAL', 'bold': True, 'color': SF_NAVY, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
            {'text': '', 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
            {'text': f'${rr_total:,.0f}', 'bold': True, 'color': C_RED,   'align': PP_ALIGN.RIGHT, 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
            {'text': str(rr_cnt), 'bold': True, 'align': PP_ALIGN.CENTER,  'fill': RGBColor(0xD9, 0xE8, 0xF7)},
            {'text': '', 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
            {'text': '', 'fill': RGBColor(0xD9, 0xE8, 0xF7)},
        ])
        row_h = min(Inches(0.42), Inches((SLIDE_H - HEADER_H - FOOTER_H - Inches(0.35)) / (len(rr_rows) + 1)))
        pptx_table(slide, rr_headers, rr_col_w, rr_rows,
                   top=HEADER_H + Inches(0.15), hdr_fill=RGBColor(0xB8, 0x76, 0x00),
                   row_ht=row_h, font_size=9, hdr_size=9)

    path = base_path + '.pptx'
    prs2.save(path)
    print(f"✅  PPTX saved: {path}")

# ── Dispatch ──────────────────────────────────────────────────────────────────
def write_html():
    import json as _json
    import html as _html
    html_version = os.environ.get('NEXT_HTML_VERSION', '')

    def status_label(r):
        if r['stage'] == 'On Hold': return 'On Hold'
        if r['is_watermelon']: return 'Watermelon'
        h = r['health'].lower()
        if h == 'red':    return 'Red'
        if h == 'yellow': return 'Yellow'
        if h == 'green':  return 'Green'
        return 'No Pulse'

    def fmt_m(v):
        if v is None: return ''
        return f'${v/1e6:.2f}M' if abs(v) >= 1_000_000 else f'${v:,.0f}'

    def fmt_pct(v):
        return f'{v:.1f}%' if v is not None else ''

    rows_data = []
    for r in results:
        dd = (r['close_margin'] - r['bid_margin']) if (r['bid_margin'] is not None and r['close_margin'] is not None) else None
        viol_codes = ', '.join(v[0] for v in r['violations'])
        bl = ' | '.join(f"{lbl}={v}" for lbl, v in baselines_list(r))
        team = []
        if r['pm']:        team.append(f"PM: {r['pm']}")
        if r.get('pm2'):   team.append(f"PM2: {r['pm2']}")
        if r.get('opp_owner'): team.append(f"AP: {r['opp_owner']}")
        if r.get('exec_sponsor'): team.append(f"ES: {r['exec_sponsor']}")
        if r.get('acct_owner'): team.append(f"AO: {r['acct_owner']}")
        if r.get('owner') and r.get('owner') != 'Unassigned': team.append(f"PO: {r['owner']}")
        po = r.get('owner','') or ''
        rows_data.append({
            'name':          r['name'],
            'acct':          r['acct'],
            'acct_id':       r.get('acct_id') or '',
            'status':        status_label(r),
            'tier':          r['tier'],
            'team':          ' · '.join(team),
            'po':            po,
            'pid':           r['pid'],
            'region':        r.get('region', REGION_LABEL),
            'url':           f"https://org62.lightning.force.com/lightning/r/pse__Proj__c/{r['pid']}/view",
            'rev_treat':     r.get('rev_treat') or '',
            'billing_type':  r.get('billing_type') or '',
            'fmt_pipe':      fmt_m(r.get('open_pipe')) if r.get('open_pipe') else '',
            'health_risk':   r.get('health_risk_score'),
            'data_quality':  r.get('data_quality_score'),
            'bookings':      r['bookings'] or 0,
            'billings':      r['billings'] or 0,
            'far':           r['far'] or 0,
            'overdue_inv':   r['overdue_inv'] or 0,
            'rr_revenue':    r['rr_revenue'] or 0,
            'bid_margin_raw':   r['bid_margin'],
            'close_margin_raw': r['close_margin'],
            'rules':         viol_codes,
            'baselines':     bl,
            'pulse_id':      r.get('pulse_id', ''),
            'summary':       ' '.join((r.get('overall_summary') or '').split()),
            'leadership':    ' '.join((r.get('leadership_notes') or '').split())[:400],
            'swe_co':        r['swe_co'],
            'high_watch':    bool(r.get('high_watch')),
            'has_pulse':     r.get('has_pulse', False),
            'pulse_trend':   r.get('trend') or '',
            'pulse_updated': (r.get('last_updated') or '')[:10],
            'pulse_scope':   r.get('scope_s') or '',
            'pulse_sched':   r.get('sched_s') or '',
            'pulse_budget':  r.get('budget_s') or '',
            'pulse_resource':r.get('resource_s') or '',
            'pulse_customer':r.get('customer_s') or '',
            'pulse_action':  ' '.join((r.get('action_needed') or '').split())[:300],
            'pulse_steerco': r['steerco_date'].isoformat() if r.get('steerco_date') else '',
            'pulse_golive':  r['next_golive'].isoformat() if r.get('next_golive') else '',
            'slack_intel':   r.get('slack_intel') or '',
            'start_dt':      r['start_dt'].isoformat() if r.get('start_dt') else '',
            'end_dt':        r['end_dt'].isoformat() if r.get('end_dt') else '',
            'stage':         r.get('stage') or '',
            'practice':      r.get('practice') or '',
            'unsch_backlog': r.get('unsch_backlog') or 0,
            'actuals_rem':   r.get('actuals_rem') or 0,
            'eva_amt':       r.get('eva_amt'),
            'eva_pct':       r.get('eva_pct'),
            'csat_score':    r.get('csat_score'),
            'fmt_unsch':     fmt_m(r.get('unsch_backlog')) if r.get('unsch_backlog') else '',
            'fmt_actuals':   fmt_m(r.get('actuals_rem')) if r.get('actuals_rem') else '',
            'fmt_eva_amt':   (f"+${r['eva_amt']:,.0f}" if r.get('eva_amt') is not None and r['eva_amt'] >= 0 else f"-${abs(r['eva_amt']):,.0f}" if r.get('eva_amt') is not None else ''),
            'fmt_eva_pct':   (f"{r['eva_pct']:+.1f}%" if r.get('eva_pct') is not None else ''),
            'far_reason':    r.get('far_reason') or '',
            'far_subreason': r.get('far_subreason') or '',
            'ptg_budget':    r.get('ptg_budget') or '',
            'ptg_scope':     r.get('ptg_scope') or '',
            'ptg_sched':     r.get('ptg_sched') or '',
            'ptg_resource':  r.get('ptg_resource') or '',
            'ptg_customer':  r.get('ptg_customer') or '',
            'fmt_bookings':  fmt_m(r['bookings']),
            'fmt_billings':  fmt_m(r['billings']),
            'fmt_bid':       fmt_pct(r['bid_margin']),
            'fmt_close':     fmt_pct(r['close_margin']),
            'fmt_delivered': (f'{dd:+.1f}%' if dd is not None else ''),
            'fmt_far':       fmt_m(r['far']),
            'fmt_overdue':   fmt_m(r['overdue_inv']) if r['overdue_inv'] else '',
            'fmt_rr':        fmt_m(r['rr_revenue']) if r['rr_revenue'] else '',
            'gdc_total':     r.get('gdc_total') or 0,
            'gdc_india':     r.get('gdc_india') or 0,
            'gdc_pct':       round(r['gdc_pct'] * 100, 1) if r.get('gdc_pct') is not None else None,
            'gdc_resources': r.get('gdc_resources') or [],
            'ironclad_url':  r.get('ironclad_url') or '',
        })

    # Unique portfolio owners for dropdown
    po_list = sorted(set(r['po'] for r in rows_data if r['po'] and r['po'] != 'Unassigned'))
    po_options = '\n'.join(f'<option value="{_html.escape(p)}">{_html.escape(p)}</option>' for p in po_list)
    # POs grouped by region for the assignment modal
    _po_by_region_dict = {}
    for _r in rows_data:
        if _r.get('po') and _r['po'] != 'Unassigned':
            _po_by_region_dict.setdefault(_r['region'], set()).add(_r['po'])
    po_by_region_js = _json.dumps({k: sorted(v) for k, v in _po_by_region_dict.items()}, ensure_ascii=False)

    # Build pipeline-by-account dict for inline embedding (no browser fetch needed)
    _pipe_by_acct_inline = {}
    for _opp in q5b_pipe:
        _aid = _opp.get('AccountId') or ''
        if not _aid:
            continue
        _pipe_by_acct_inline.setdefault(_aid, []).append({
            'id':           _opp.get('Id', ''),
            'name':         _opp.get('Name', ''),
            'stage':        _opp.get('StageName', ''),
            'amount':       _opp.get('Amount') or 0,
            'converted':    _opp.get('convertedAmount') or _opp.get('Amount') or 0,
            'currency':     _opp.get('CurrencyIsoCode', 'USD'),
            'close':        (_opp.get('CloseDate') or '')[:10],
            'created':      (_opp.get('CreatedDate') or '')[:10],
            'owner':        _opp.get('Owner.Name', '') or '',
            'forecast':     _opp.get('ForecastCategory', ''),
            'mgr_forecast': _opp.get('Manager_Forecast_Judgement__c', ''),
            'description':  _opp.get('Description', '') or '',
            'bid_margin':   _opp.get('Overall_Bid_Margin__c'),
            'push_count':   _opp.get('PushCount') or 0,
        })

    rows_json = _json.dumps({'generated': REPORT_DATE, 'region': REGION_LABEL, 'rows': rows_data, 'pipe': _pipe_by_acct_inline}, ensure_ascii=False)

    # DB credentials for self-serve assignments (embedded at publish time)
    import base64 as _b64_html
    _db_cid    = os.environ.get('DB_CLIENT_ID', '')
    _db_csec   = os.environ.get('DB_CLIENT_SECRET', '')
    _tile_id_h = os.environ.get('TILE_ID', '817')
    _db_b64    = _b64_html.b64encode(f"{_db_cid}:{_db_csec}".encode()).decode() if (_db_cid and _db_csec) else ''

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>ACC Delivery Portfolio — {REPORT_DATE}</title>
<script src="/assets/tile-db.js"></script>
<style>
  :root {{
    --navy:    #0D2B4E;
    --blue:    #1F497D;
    --lblue:   #4472C4;
    --red:     #C0392B;
    --yellow:  #D4AC0D;
    --green:   #1E8449;
    --wm:      #E74C3C;
    --grey:    #6B7280;
    --bg:      #F8F9FA;
    --card:    #FFFFFF;
    --border:  #E5E7EB;
    --text:    #111827;
    --subtext: #6B7280;
    --stripe:  #F1F5FB;
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  html, body {{ height: 100%; overflow: hidden; }}
  body {{ font-family: 'Segoe UI', system-ui, sans-serif; background: var(--bg); color: var(--text); font-size: 13px; display: flex; flex-direction: column; }}

  /* ── Header ── */
  .page-header {{ background: var(--blue); color: #fff; padding: 7px 18px; display: flex; align-items: center; justify-content: space-between; }}
  .page-header h1 {{ font-size: 14px; font-weight: 700; letter-spacing: .3px; }}
  .page-header .meta {{ font-size: 10px; opacity: .75; margin-top: 1px; }}

  /* ── Scorecard ── */
  .scorecard {{ display: flex; flex-wrap: wrap; gap: 6px; padding: 8px 18px; background: var(--card); border-bottom: 1px solid var(--border); transition: padding .2s; }}
  .scorecard.collapsed {{ padding-top: 0; padding-bottom: 0; overflow: hidden; max-height: 0; }}
  .scorecard-bar {{ display: flex; align-items: center; justify-content: space-between; padding: 4px 18px; background: var(--card); border-bottom: 1px solid var(--border); cursor: pointer; user-select: none; }}
  .scorecard-bar:hover {{ background: #f0f4ff; }}
  .scorecard-bar .sc-toggle {{ font-size: 11px; color: var(--subtext); display: flex; align-items: center; gap: 5px; }}
  .scorecard-bar .sc-summary {{ font-size: 11px; color: var(--subtext); }}
  .stat {{ background: var(--bg); border: 1px solid var(--border); border-radius: 6px; padding: 6px 10px; min-width: 90px; display:flex; flex-direction:column; }}
  .stat .val {{ font-size: 16px; font-weight: 700; color: var(--blue); }}
  .stat .lbl {{ font-size: 9px; text-transform: uppercase; letter-spacing: .5px; color: var(--subtext); margin-bottom: 1px; }}
  .stat.red    .val {{ color: var(--red); }}
  .stat.yellow .val {{ color: var(--yellow); }}
  .stat.green  .val {{ color: var(--green); }}
  .stat.wm     .val {{ color: var(--wm); }}

  /* ── Region tabs ── */
  .region-tabs {{ display: flex; gap: 0; background: var(--navy); padding: 0 18px; }}
  .rtab {{ background: none; border: none; border-bottom: 3px solid transparent; color: rgba(255,255,255,.65); font-size: 12px; font-weight: 600; padding: 7px 18px; cursor: pointer; letter-spacing: .3px; transition: color .15s, border-color .15s; }}
  .rtab:hover {{ color: #fff; }}
  .rtab.active {{ color: #fff; border-bottom-color: #4FC3F7; }}

  /* ── Controls ── */
  .controls {{ padding: 0; background: var(--card); border-bottom: 1px solid var(--border); }}
  .controls-row {{ display: flex; flex-wrap: nowrap; gap: 6px; align-items: center; padding: 6px 18px; }}
  .controls-row-1 {{ border-bottom: 1px solid var(--border); }}
  .controls-row-spacer {{ flex: 1; }}
  @media (max-width: 900px) {{ .controls-row {{ flex-wrap: wrap; }} }}
  .controls input[type=search] {{
    border: 1px solid var(--border); border-radius: 6px; padding: 4px 10px;
    font-size: 12px; width: 240px; outline: none;
  }}
  .controls input[type=search]:focus {{ border-color: var(--lblue); box-shadow: 0 0 0 2px #4472c420; }}
  .filter-group {{ display: flex; gap: 4px; align-items: center; }}
  .pill {{
    border: 1px solid var(--border); border-radius: 20px; padding: 2px 9px;
    font-size: 11px; cursor: pointer; background: var(--bg); color: var(--subtext);
    transition: all .15s; user-select: none;
  }}
  .pill:hover {{ border-color: var(--lblue); color: var(--lblue); }}
  .pill.active {{ background: var(--blue); color: #fff; border-color: var(--blue); font-weight: 600; }}
  .pill.red-pill.active    {{ background: var(--red);    border-color: var(--red); }}
  .pill.yellow-pill.active {{ background: var(--yellow); border-color: var(--yellow); }}
  .pill.green-pill.active  {{ background: var(--green);  border-color: var(--green); }}
  .pill.wm-pill.active     {{ background: var(--wm);     border-color: var(--wm); }}
  .sep {{ width: 1px; height: 18px; background: var(--border); align-self: center; }}
  .count-badge {{ font-size: 10px; color: var(--subtext); margin-left: 3px; }}
  .po-select {{
    border: 1px solid var(--border); border-radius: 6px; padding: 3px 8px;
    font-size: 11px; background: var(--bg); color: var(--text); cursor: pointer;
    outline: none; min-width: 150px;
  }}
  .po-select:focus {{ border-color: var(--lblue); }}

  /* ── Multi-select combobox ── */
  .ms-wrap {{ position: relative; display: inline-block; }}
  .ms-btn {{
    border: 1px solid var(--border); border-radius: 6px; padding: 3px 24px 3px 8px;
    font-size: 11px; background: var(--bg); color: var(--text); cursor: pointer;
    outline: none; min-width: 120px; text-align: left; white-space: nowrap;
    user-select: none; position: relative;
  }}
  .ms-btn::after {{ content: '▾'; position: absolute; right: 7px; top: 50%; transform: translateY(-50%); font-size: 10px; color: var(--subtext); }}
  .ms-btn.ms-active {{ border-color: var(--blue); color: var(--blue); background: #eef2ff; font-weight: 600; }}
  .ms-dropdown {{
    display: none; position: absolute; top: calc(100% + 3px); left: 0; z-index: 9000;
    background: var(--card); border: 1px solid var(--border); border-radius: 8px;
    box-shadow: 0 4px 16px rgba(0,0,0,.12); min-width: 180px; max-height: 260px;
    overflow-y: auto; padding: 4px 0;
  }}
  .ms-dropdown.open {{ display: block; }}
  .ms-item {{
    display: flex; align-items: center; gap: 7px; padding: 5px 12px;
    font-size: 11px; cursor: pointer; color: var(--text);
  }}
  .ms-item:hover {{ background: #f3f6ff; }}
  .ms-item input[type=checkbox] {{ accent-color: var(--blue); cursor: pointer; margin: 0; }}
  .ms-clear {{ display:block; padding: 4px 12px 5px; font-size:10px; color:var(--blue); cursor:pointer; border-top:1px solid var(--border); }}
  .ms-clear:hover {{ text-decoration: underline; }}

  /* ── Table ── */
  .table-wrap {{ overflow: auto; padding: 0 24px 32px; flex: 1; }}
  table {{ width: 100%; border-collapse: collapse; margin-top: 12px; }}
  thead th {{
    background: var(--blue); color: #fff; font-size: 11px; font-weight: 600;
    text-transform: uppercase; letter-spacing: .4px; padding: 9px 10px;
    text-align: left; position: sticky; top: 0; z-index: 10; cursor: pointer;
    white-space: nowrap; user-select: none;
  }}
  thead th:hover {{ background: var(--lblue); }}
  thead th .sort-arrow {{ margin-left: 4px; opacity: .4; font-size: 10px; }}
  thead th.sorted-asc  .sort-arrow::after {{ content: '▲'; opacity: 1; }}
  thead th.sorted-desc .sort-arrow::after {{ content: '▼'; opacity: 1; }}
  thead th:not(.sorted-asc):not(.sorted-desc) .sort-arrow::after {{ content: '⇅'; }}

  tbody tr {{ border-bottom: 1px solid var(--border); transition: background .1s; }}
  tbody tr:nth-child(even) {{ background: var(--stripe); }}
  tbody tr:hover {{ background: #EBF3FB; }}
  tbody td {{ padding: 8px 10px; vertical-align: top; }}

  /* status badge */
  .badge {{ display: inline-block; padding: 2px 8px; border-radius: 12px; font-size: 11px; font-weight: 600; white-space: nowrap; }}
  .badge-Red        {{ background: #FDEDEC; color: var(--red); }}
  .badge-Yellow     {{ background: #FEF9E7; color: #9A7D0A; }}
  .badge-Green      {{ background: #EAFAF1; color: var(--green); }}
  .badge-Watermelon {{ background: #FDEDEC; color: var(--wm); }}
  .badge-No-Pulse   {{ background: #F3F4F6; color: var(--grey); }}
  .badge-On-Hold    {{ background: #EDE7F6; color: #6A1B9A; }}

  .hw-flag {{ display: inline-block; font-size: 10px; background: #EBF5FB; color: var(--lblue); border-radius: 4px; padding: 1px 5px; margin-right: 4px; font-weight: 600; }}
  .score-chip      {{ display: inline-block; font-size: 10px; border-radius: 4px; padding: 1px 5px; font-weight: 700; border: 1px solid; text-decoration: none; cursor: pointer; }}
  a.score-chip:hover {{ opacity: .8; text-decoration: underline; }}
  .score-green {{ background: #EAFAF1; color: #1E8449; border-color: #A9DFBF; }}
  .score-yellow {{ background: #FEF9E7; color: #9A7D0A; border-color: #F9E79F; }}
  .score-red   {{ background: #FDEDEC; color: #C0392B; border-color: #FADBD8; }}

  /* tier badge */
  .tier {{ display: inline-block; font-size: 10px; font-weight: 700; padding: 1px 6px; border-radius: 4px; background: #E8EAF6; color: #3949AB; }}

  /* team list */
  .team-list {{ line-height: 1.7; }}
  .team-list span {{ display: block; font-size: 11px; color: var(--subtext); }}
  .team-list span b {{ color: var(--text); font-weight: 600; }}
  .role-abbr {{ position: relative; display: inline-block; }}
  .role-abbr .role-tip {{
    display: none; position: absolute; bottom: calc(100% + 4px); left: 0;
    background: #1e293b; color: #fff; font-size: 10px; font-weight: 400;
    white-space: nowrap; padding: 3px 7px; border-radius: 4px;
    box-shadow: 0 2px 6px rgba(0,0,0,.25); pointer-events: none; z-index: 9999;
  }}
  .role-abbr:hover .role-tip {{ display: block; }}
  .res-hover-wrap {{ display: inline-block; cursor: pointer; }}
  .popup-drag-handle {{ cursor: move; user-select: none; }}
  #res-modal {{
    display: none; position: fixed; z-index: 9999;
    top: 50%; left: 50%; transform: translate(-50%, -50%);
    width: min(980px, 97vw); max-height: 80vh;
    min-width: 400px; min-height: 200px;
    background: #fff; border: 1px solid #ddd; border-radius: 10px;
    box-shadow: 0 12px 40px rgba(0,0,0,.28);
    font-size: 13px; color: var(--text);
    overflow: hidden; pointer-events: auto;
    flex-direction: column; resize: both;
  }}
  #res-modal.open {{ display: flex; }}
  #res-modal-header {{
    display: flex; align-items: center; justify-content: space-between;
    padding: 12px 16px 10px; border-bottom: 1px solid rgba(255,255,255,.15);
    flex-shrink: 0; background: var(--navy);
  }}
  #res-modal-header strong {{ font-size: 14px; color: #fff; }}
  #res-modal-close {{
    background: none; border: none; cursor: pointer; font-size: 16px;
    color: rgba(255,255,255,.8); line-height: 1; padding: 0 0 0 10px;
  }}
  #res-modal-body {{ overflow-y: auto; flex: 1; padding: 12px 16px 16px; }}
  #res-modal-body table {{ border-collapse: collapse; font-size: 11px; width: 100%; }}
  #res-modal-body th {{ font-size: 10px; color: #fff; font-weight: 600; padding: 5px 10px; text-align: left; text-transform: uppercase; letter-spacing: .3px; background: var(--navy); white-space: nowrap; cursor: pointer; user-select: none; }}
  #res-modal-body td {{ padding: 4px 10px; vertical-align: middle; color: var(--text); white-space: nowrap; border-bottom: 1px solid var(--border); }}
  .assign-btn {{ background: none; border: none; cursor: pointer; padding: 0 2px; font-size: 11px; color: var(--subtext); line-height: 1; vertical-align: middle; }}
  .assign-btn:hover {{ color: var(--blue); }}
  /* Assignment modal */
  #assign-modal {{ display: none; position: fixed; inset: 0; z-index: 9999; align-items: center; justify-content: center; }}
  #assign-modal.open {{ display: flex; }}
  #assign-backdrop {{ position: absolute; inset: 0; background: rgba(0,0,0,.35); }}
  #assign-dialog {{ position: relative; background: #fff; border-radius: 10px; padding: 20px 24px 16px; min-width: 260px; box-shadow: 0 8px 32px rgba(0,0,0,.22); z-index: 1; }}
  #assign-dialog h3 {{ margin: 0 0 4px; font-size: 13px; color: var(--navy); }}
  #assign-dialog .assign-acct {{ font-size: 11px; color: var(--subtext); margin-bottom: 14px; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; max-width: 220px; }}
  #assign-dialog label {{ display: block; font-size: 11px; font-weight: 600; color: var(--text); margin-bottom: 3px; }}
  #assign-dialog select {{ width: 100%; font-size: 12px; padding: 5px 6px; border: 1px solid var(--border); border-radius: 5px; background: var(--bg); margin-bottom: 12px; }}
  #assign-dialog .assign-actions {{ display: flex; gap: 8px; justify-content: flex-end; margin-top: 4px; }}
  #assign-dialog .btn-save {{ background: var(--blue); color: #fff; border: none; border-radius: 5px; padding: 5px 16px; font-size: 12px; font-weight: 600; cursor: pointer; }}
  #assign-dialog .btn-save:hover {{ background: var(--navy); }}
  #assign-dialog .btn-cancel {{ background: none; border: 1px solid var(--border); border-radius: 5px; padding: 5px 12px; font-size: 12px; cursor: pointer; color: var(--subtext); }}
  #assign-dialog .btn-cancel:hover {{ border-color: var(--text); color: var(--text); }}
  #assign-dialog {{ min-width: 300px; }}
  .assign-roster-toggle {{ background: none; border: none; cursor: pointer; font-size: 11px; color: var(--blue); padding: 0; margin-bottom: 10px; text-decoration: underline; }}
  #assign-roster-panel {{ display: none; border: 1px solid var(--border); border-radius: 6px; padding: 8px 10px; margin-bottom: 12px; background: var(--bg); }}
  #assign-roster-panel.open {{ display: block; }}
  .roster-region-label {{ font-size: 10px; font-weight: 700; color: var(--subtext); text-transform: uppercase; letter-spacing: .05em; margin-bottom: 6px; }}
  .roster-list {{ max-height: 140px; overflow-y: auto; margin-bottom: 8px; }}
  .roster-item {{ display: flex; align-items: center; justify-content: space-between; padding: 2px 0; font-size: 12px; gap: 6px; }}
  .roster-item span {{ flex: 1; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }}
  .roster-del {{ background: none; border: none; cursor: pointer; color: var(--subtext); font-size: 13px; padding: 0 2px; line-height: 1; flex-shrink: 0; }}
  .roster-del:hover {{ color: var(--red); }}
  .roster-del:disabled {{ opacity: .35; cursor: default; color: var(--subtext) !important; }}
  .roster-add-row {{ display: flex; gap: 6px; margin-top: 4px; }}
  .roster-add-row input {{ flex: 1; font-size: 12px; padding: 4px 7px; border: 1px solid var(--border); border-radius: 4px; background: #fff; }}
  .roster-add-row button {{ font-size: 12px; padding: 4px 10px; background: var(--blue); color: #fff; border: none; border-radius: 4px; cursor: pointer; white-space: nowrap; }}
  .roster-add-row button:hover {{ background: var(--navy); }}

  /* financials */
  .fin-2col  {{ display: flex; gap: 12px; }}
  .fin-col   {{ flex: 1; min-width: 0; display: flex; flex-direction: column; gap: 1px; font-size: 11px; }}
  .fin-item  {{ display: flex; justify-content: space-between; align-items: baseline; gap: 6px; }}
  .fin-item .fk {{ color: var(--text); font-weight: 600; white-space: nowrap; flex-shrink: 0; }}
  .fin-item .fv {{ white-space: nowrap; text-align: right; min-width: 0; overflow: hidden; text-overflow: ellipsis; color: var(--subtext); }}
  .fin-item .fv.neg {{ color: var(--red); }}
  .fin-item .fv.pos {{ color: var(--green); }}
  .fin-sep   {{ border: none; border-top: 1px solid var(--border); margin: 2px 0; }}
  .fin-extras {{ margin-top: 3px; padding-top: 3px; border-top: 1px solid var(--border); display: flex; flex-direction: column; gap: 1px; font-size: 11px; }}

  /* baselines */
  .bl-list {{ font-size: 11px; line-height: 1.8; }}
  .bl-lbl  {{ color: var(--lblue); font-weight: 700; margin-right: 3px; }}
  .bl-dot  {{ display: inline-block; width: 9px; height: 9px; border-radius: 50%; vertical-align: middle; }}
  .bl-dot-Red    {{ background: var(--red); }}
  .bl-dot-Yellow {{ background: var(--yellow); }}
  .bl-dot-Green  {{ background: var(--green); }}
  .bl-dot-None   {{ background: var(--border); }}

  /* rules */
  .rule-code {{
    display: inline-block; font-size: 10px; background: #FFF3CD; color: #856404;
    border-radius: 3px; padding: 1px 5px; margin: 1px 2px 1px 0; border: 1px solid #FFEAA7;
    cursor: help; position: relative;
  }}
  .rule-code.neg  {{ background: #FDEDEC; color: var(--red); border-color: #FADBD8; }}
  .rule-code.warn {{ background: #FFF3CD; color: #856404; border-color: #FFEAA7; }}
  .rule-code .rule-tip {{
    display: none; position: absolute; bottom: calc(100% + 4px); left: 50%;
    transform: translateX(-50%); background: #1F2937; color: #fff;
    font-size: 11px; border-radius: 5px; padding: 6px 10px; white-space: nowrap;
    z-index: 100; pointer-events: none; box-shadow: 0 2px 8px rgba(0,0,0,.25);
    font-weight: normal; letter-spacing: 0;
  }}
  .rule-code .rule-tip::after {{
    content: ''; position: absolute; top: 100%; left: 50%; transform: translateX(-50%);
    border: 5px solid transparent; border-top-color: #1F2937;
  }}
  .rule-code:hover .rule-tip {{ display: block; }}

  /* summary expandable */
  .summary-cell {{ max-width: 320px; }}
  .summary-short {{ color: var(--text); line-height: 1.5; font-size: 11px; }}
  .summary-full {{ display: none; color: var(--subtext); font-size: 11px; margin-top: 4px; line-height: 1.5; border-top: 1px dashed var(--border); padding-top: 4px; }}
  .expand-btn {{ font-size: 10px; color: var(--lblue); cursor: pointer; margin-top: 3px; display: inline-block; }}
  .expand-btn:hover {{ text-decoration: underline; }}
  tr.expanded .summary-full {{ display: block; }}

  /* project name */
  .proj-name {{ font-weight: 600; color: var(--text); font-size: 12px; }}
  .proj-acct {{ font-size: 12px; font-weight: 600; color: var(--blue); margin-top: 2px; }}
  .proj-link {{ color: inherit; text-decoration: none; }}
  .proj-link:hover {{ color: var(--lblue); text-decoration: underline; }}
  .pulse-icon {{ display: inline-block; cursor: pointer; font-size: 13px; vertical-align: middle; margin-left: 4px; }}
  #summary-modal {{
    display: none; position: fixed; z-index: 9999;
    top: 50%; left: 50%; transform: translate(-50%, -50%);
    width: min(660px, 93vw); max-height: 82vh;
    min-width: 340px; min-height: 200px;
    background: #fff; border: 1px solid #ddd; border-radius: 10px;
    box-shadow: 0 12px 40px rgba(0,0,0,.28);
    font-size: 13px; line-height: 1.7; color: var(--text);
    overflow: hidden; pointer-events: auto; flex-direction: column; resize: both;
  }}
  #summary-modal.open {{ display: flex; }}
  #summary-modal-header {{
    display: flex; align-items: center; justify-content: space-between;
    padding: 12px 16px 10px; border-bottom: 1px solid var(--border);
    flex-shrink: 0; background: var(--navy);
  }}
  #summary-modal-header strong {{ font-size: 13px; color: #fff; }}
  #summary-modal-close {{
    background: none; border: none; cursor: pointer; font-size: 16px;
    color: rgba(255,255,255,.8); line-height: 1; padding: 0 0 0 10px;
  }}
  #summary-modal-body {{ padding: 14px 16px 16px; overflow-y: auto; flex: 1; }}
  #summary-modal-footer {{
    flex-shrink: 0; padding: 10px 16px; border-top: 1px solid var(--border);
    background: var(--bg); display: flex; flex-wrap: wrap; gap: 10px;
    font-size: 11px; color: var(--subtext);
  }}
  #summary-modal-footer span strong {{ color: var(--text); }}
  #summary-modal-footer a {{ font-size: 11px; color: var(--blue); text-decoration: none; font-weight: 500; }}
  .sum-section {{ margin-bottom: 14px; }}
  .sum-section-title {{ font-size: 10px; font-weight: 700; text-transform: uppercase;
    letter-spacing: .6px; color: var(--subtext); margin-bottom: 4px; }}
  .sum-rag-row {{ display: flex; flex-wrap: wrap; gap: 6px; margin-bottom: 10px; }}
  .sum-rag-chip {{ display: inline-flex; align-items: center; gap: 4px;
    font-size: 11px; font-weight: 600; padding: 2px 8px; border-radius: 20px; }}
  .sum-rag-chip.Red    {{ background: #fde8e8; color: var(--red); }}
  .sum-rag-chip.Yellow {{ background: #fff8e1; color: #b45309; }}
  .sum-rag-chip.Green  {{ background: #e6f4ea; color: var(--green); }}
  .sum-rag-chip.Unknown {{ background: var(--bg); color: var(--subtext); }}
  .sum-block {{ background: var(--bg); border-radius: 6px; padding: 10px 12px;
    font-size: 12px; line-height: 1.6; white-space: pre-wrap; word-break: break-word; }}
  #pulse-tooltip {{
    display: none; position: fixed; z-index: 9999;
    top: 50%; left: 50%; transform: translate(-50%, -50%);
    width: min(700px, 92vw); max-height: 80vh;
    min-width: 320px; min-height: 180px;
    background: #fff; border: 1px solid #ddd; border-radius: 10px;
    box-shadow: 0 12px 40px rgba(0,0,0,.28);
    font-size: 13px; line-height: 1.7; color: var(--text);
    overflow: hidden; pointer-events: auto;
    flex-direction: column; resize: both;
  }}
  #pulse-tooltip.open {{ display: flex; }}
  #pulse-tooltip-header {{
    display: flex; align-items: center; justify-content: space-between;
    padding: 12px 16px 10px; border-bottom: 1px solid var(--border);
    flex-shrink: 0; background: #fff;
  }}
  #pulse-tooltip-header strong {{ font-size: 14px; color: var(--navy); }}
  #pulse-tooltip-header a {{ font-size: 11px; color: var(--blue); text-decoration: none; }}
  #pulse-tooltip-header a:hover {{ text-decoration: underline; }}
  #pulse-tooltip-close {{
    background: none; border: none; cursor: pointer; font-size: 16px;
    color: var(--subtext); line-height: 1; padding: 0 0 0 10px;
  }}
  #pulse-tooltip-body {{ padding: 12px 16px 16px; overflow-y: auto; flex: 1; }}
  #far-tooltip {{
    display: none; position: fixed; z-index: 9999;
    top: 50%; left: 50%; transform: translate(-50%, -50%);
    width: min(480px, 92vw); max-height: 80vh;
    min-width: 280px; min-height: 160px;
    background: #fff; border: 1px solid #ddd; border-radius: 10px;
    box-shadow: 0 12px 40px rgba(0,0,0,.28);
    font-size: 13px; line-height: 1.7; color: var(--text);
    overflow: hidden; pointer-events: auto;
    flex-direction: column; resize: both;
  }}
  #far-tooltip.open {{ display: flex; }}
  #far-tooltip-header {{
    display: flex; align-items: center; justify-content: space-between;
    padding: 12px 16px 10px; border-bottom: 1px solid var(--border);
    flex-shrink: 0; background: #fff;
  }}
  #far-tooltip-header strong {{ font-size: 14px; color: var(--navy); }}
  #far-tooltip-close {{
    background: none; border: none; cursor: pointer; font-size: 16px;
    color: var(--subtext); line-height: 1; padding: 0 0 0 10px;
  }}
  #far-tooltip-body {{ padding: 12px 16px 16px; overflow-y: auto; flex: 1; }}
  #p2g-modal {{
    display: none; position: fixed; z-index: 9999;
    top: 50%; left: 50%; transform: translate(-50%, -50%);
    width: min(520px, 93vw); max-height: 80vh;
    min-width: 280px; min-height: 160px;
    background: #fff; border: 1px solid #ddd; border-radius: 10px;
    box-shadow: 0 12px 40px rgba(0,0,0,.28);
    font-size: 13px; line-height: 1.7; color: var(--text);
    overflow: hidden; pointer-events: auto;
    flex-direction: column; resize: both;
  }}
  #p2g-modal.open {{ display: flex; }}
  #p2g-modal-header {{
    display: flex; align-items: center; justify-content: space-between;
    padding: 12px 16px 10px; border-bottom: 1px solid var(--border);
    flex-shrink: 0; background: var(--navy);
  }}
  #p2g-modal-header strong {{ font-size: 14px; color: #fff; }}
  #p2g-modal-close {{
    background: none; border: none; cursor: pointer; font-size: 16px;
    color: rgba(255,255,255,.8); line-height: 1; padding: 0 0 0 10px;
  }}
  #p2g-modal-body {{ padding: 12px 16px 16px; overflow-y: auto; flex: 1; }}
  #pipe-modal {{
    display: none; position: fixed; z-index: 9999;
    top: 50%; left: 50%; transform: translate(-50%, -50%);
    width: min(1200px, 97vw); max-height: 85vh;
    min-width: 400px; min-height: 200px;
    background: #fff; border: 1px solid #ddd; border-radius: 10px;
    box-shadow: 0 12px 40px rgba(0,0,0,.28);
    font-size: 13px; color: var(--text);
    overflow: hidden; pointer-events: auto;
    flex-direction: column; resize: both;
  }}
  #pipe-modal.open {{ display: flex; }}
  #pipe-header {{
    display: flex; align-items: center; justify-content: space-between;
    padding: 12px 16px 10px; border-bottom: 1px solid var(--border);
    flex-shrink: 0; background: #fff;
  }}
  #pipe-title {{ font-size: 14px; color: var(--navy); font-weight: 700; }}
  #pipe-close {{
    background: none; border: none; cursor: pointer; font-size: 16px;
    color: var(--subtext); line-height: 1; padding: 0 0 0 10px;
  }}
  #pipe-body {{ overflow-y: auto; flex: 1; padding: 12px 16px 16px; font-size: 12px; }}
  .pulse-dim {{ display: flex; gap: 4px; flex-wrap: wrap; margin: 5px 0 3px; }}
  .pdim {{ font-size: 10px; border-radius: 3px; padding: 1px 5px; font-weight: 700; }}
  .pdim-green  {{ background:#EAFAF1; color:#1a6b3c; }}
  .pdim-yellow {{ background:#FEF9E7; color:#7d6608; }}
  .pdim-red    {{ background:#FDEDEC; color:#922b21; }}
  .pdim-grey   {{ background:#f0f0f0; color:#666; }}

  /* no-results */
  .no-results {{ text-align: center; padding: 48px; color: var(--subtext); font-size: 14px; }}

  /* ── Group rows ── */
  .group-row td {{
    background: #EBF3FB; border-top: 2px solid var(--lblue); border-bottom: 1px solid var(--border);
    padding: 7px 10px; font-size: 11px; font-weight: 700; color: var(--navy);
    letter-spacing: .3px; cursor: pointer; user-select: none;
  }}
  .group-tier-header td {{
    background: var(--navy) !important; color: #fff !important;
    border-top: 3px solid var(--navy) !important; border-bottom: 2px solid #0a2240 !important;
    padding: 9px 12px !important; font-size: 12px !important; letter-spacing: .4px;
  }}
  .group-tier-header td .grp-toggle {{ color: #fff !important; opacity: .8; }}
  .group-tier-header td .grp-meta {{ color: rgba(255,255,255,.75) !important; }}
  .group-tier-header td .grp-pills .grp-pill {{ opacity: .92; }}
  .group-acct-header td {{ padding-left: 28px !important; border-left: 4px solid var(--lblue); }}
  .grp-btn {{ font-size: 11px; padding: 3px 9px; border-radius: 5px; border: 1px solid var(--border); background: var(--bg); color: var(--text); cursor: pointer; }}
  .grp-btn:hover {{ background: var(--blue); color: #fff; border-color: var(--blue); }}
  .group-row td .grp-toggle {{ margin-right: 6px; font-size: 10px; opacity: .6; }}
  .group-row td .grp-meta {{ font-weight: 400; color: var(--subtext); margin-left: 10px; font-size: 11px; }}
  .group-row td .grp-pills {{ display: inline-flex; gap: 4px; margin-left: 10px; }}
  .group-row td .grp-pill {{ font-size: 10px; border-radius: 10px; padding: 1px 7px; font-weight: 700; }}
  .grp-pill-nopulse {{ background: #e8e8e8; color: #555; }}
  .grp-pill-red    {{ background:#FDEDEC; color:var(--red); }}
  .grp-pill-yellow {{ background:#FEF9E7; color:#9A7D0A; }}
  .grp-pill-wm     {{ background:#FDEDEC; color:var(--wm); }}
  .grp-pill-green  {{ background:#EAFAF1; color:var(--green); }}
  tr.grp-hidden {{ display: none; }}
  .grp-order-pill {{ font-size: 10px; border-radius: 10px; padding: 2px 8px; font-weight: 600; background: var(--blue); color: #fff; white-space: nowrap; }}

  /* column widths */
  .col-status  {{ width: 90px; }}
  .col-tier    {{ width: 40px; text-align: center; }}
  .col-project {{ width: 260px; overflow-wrap: break-word; word-break: break-word; }}
  .col-team     {{ width: 180px; }}
  .col-resource {{ width: 90px; }}
  .col-fin      {{ width: 230px; }}
  .col-rules   {{ width: 130px; }}
  .col-bl      {{ width: 160px; }}
  .col-summary {{ min-width: 200px; }}

  /* ── GM Business Overview ─────────────────────────────────────────── */
  #gm-overview-btn {{
    background: var(--blue); color: #fff; border: none; border-radius: 5px;
    padding: 5px 12px; font-size: 11px; font-weight: 600; cursor: pointer;
    white-space: nowrap; margin-left: 8px;
  }}
  #gm-overview-btn:hover {{ opacity: .88; }}
  /* ── Help modal ───────────────────────────────────────────────────── */
  #help-btn {{
    background: none; color: var(--subtext); border: 1px solid var(--border);
    border-radius: 5px; padding: 5px 10px; font-size: 11px; font-weight: 600;
    cursor: pointer; white-space: nowrap; margin-left: 6px;
  }}
  #help-btn:hover {{ background: var(--hover); color: var(--text); }}
  #help-modal {{ display: none; position: fixed; inset: 0; z-index: 10001; align-items: center; justify-content: center; }}
  #help-modal.open {{ display: flex; }}
  #help-backdrop {{ position: absolute; inset: 0; background: rgba(0,0,0,.45); backdrop-filter: blur(2px); }}
  #help-dialog {{
    position: relative; z-index: 1; background: var(--card); border-radius: 10px;
    width: min(780px, 94vw); max-height: 88vh; display: flex; flex-direction: column;
    box-shadow: 0 8px 40px rgba(0,0,0,.28);
  }}
  #help-header {{
    display: flex; align-items: center; justify-content: space-between;
    background: var(--blue); color: #fff; border-radius: 10px 10px 0 0;
    padding: 10px 16px; flex-shrink: 0;
  }}
  #help-header span {{ font-size: 14px; font-weight: 700; }}
  #help-tabs {{
    display: flex; gap: 0; border-bottom: 1px solid var(--border); flex-shrink: 0;
    padding: 0 16px; background: var(--card);
  }}
  .help-tab {{
    padding: 9px 16px; font-size: 12px; font-weight: 600; cursor: pointer;
    border-bottom: 2px solid transparent; color: var(--subtext); margin-bottom: -1px;
  }}
  .help-tab.active {{ border-bottom-color: var(--blue); color: var(--blue); }}
  #help-body {{ flex: 1; overflow-y: auto; padding: 18px 20px; font-size: 13px; line-height: 1.65; }}
  #help-body h3 {{
    font-size: 11px; font-weight: 700; text-transform: uppercase; letter-spacing: .5px;
    color: var(--blue); border-bottom: 1px solid var(--border); padding-bottom: 4px;
    margin: 18px 0 8px;
  }}
  #help-body h3:first-child {{ margin-top: 0; }}
  #help-body ul {{ margin: 4px 0 10px 0; padding-left: 18px; }}
  #help-body li {{ margin-bottom: 5px; }}
  #help-body strong {{ color: var(--text); }}
  .prompt-card {{
    border: 1px solid var(--border); border-radius: 7px; margin-bottom: 16px;
    overflow: hidden;
  }}
  .prompt-card-header {{
    display: flex; align-items: center; justify-content: space-between;
    padding: 8px 12px; background: var(--row-alt); border-bottom: 1px solid var(--border);
  }}
  .prompt-card-title {{ font-size: 12px; font-weight: 700; color: var(--text); }}
  .prompt-card-audience {{ font-size: 11px; color: var(--subtext); }}
  .prompt-copy-btn {{
    background: var(--blue); color: #fff; border: none; border-radius: 4px;
    padding: 3px 10px; font-size: 11px; font-weight: 600; cursor: pointer; flex-shrink: 0;
  }}
  .prompt-copy-btn:hover {{ opacity: .85; }}
  .prompt-copy-btn.copied {{ background: var(--green); }}
  .prompt-text {{
    padding: 10px 12px; font-size: 11.5px; line-height: 1.6; color: var(--subtext);
    white-space: pre-wrap; font-family: inherit; max-height: 140px; overflow-y: auto;
    background: #fff;
  }}
  #gm-modal {{ display: none; position: fixed; inset: 0; z-index: 10000; align-items: center; justify-content: center; }}
  #gm-modal.open {{ display: flex; }}
  #gm-backdrop {{ position: absolute; inset: 0; background: rgba(0,0,0,.45); backdrop-filter: blur(2px); }}
  #gm-dialog {{
    position: relative; z-index: 1; background: var(--card); border-radius: 10px;
    width: min(740px, 93vw); min-height: 420px; max-height: 93vh; display: flex; flex-direction: column;
    box-shadow: 0 8px 40px rgba(0,0,0,.28); transition: max-height .15s ease;
  }}
  #gm-header {{
    display: flex; align-items: center; justify-content: space-between;
    background: var(--blue); color: #fff; border-radius: 10px 10px 0 0;
    padding: 10px 14px; gap: 8px; flex-shrink: 0;
  }}
  #gm-scope-label {{ font-size: 13px; font-weight: 700; flex: 1; }}
  #gm-regen-btn {{
    background: rgba(255,255,255,.18); color: #fff; border: 1px solid rgba(255,255,255,.4);
    border-radius: 4px; padding: 3px 9px; font-size: 11px; cursor: pointer;
  }}
  #gm-regen-btn:hover {{ background: rgba(255,255,255,.3); }}
  #gm-regen-btn:disabled {{ opacity: .4; cursor: default; }}
  /* GM modal tabs */
  #gm-tabs {{
    display: flex; gap: 0; border-bottom: 1px solid var(--border); flex-shrink: 0;
    padding: 0 14px; background: var(--card);
  }}
  .gm-tab {{
    padding: 8px 14px; font-size: 12px; font-weight: 600; cursor: pointer;
    border-bottom: 2px solid transparent; color: var(--subtext); margin-bottom: -1px;
    white-space: nowrap;
  }}
  .gm-tab.active {{ border-bottom-color: var(--blue); color: var(--blue); }}
  .gm-tab-panel {{ display: none; flex: 1; flex-direction: column; min-height: 0; }}
  .gm-tab-panel.active {{ display: flex; }}
  #gm-body {{
    flex: 1; overflow-y: auto; padding: 16px 20px; font-size: 13px; line-height: 1.65;
    color: var(--text);
  }}
  #gm-body h3 {{
    font-size: 12px; font-weight: 700; text-transform: uppercase; letter-spacing: .5px;
    color: var(--blue); border-bottom: 1px solid var(--border); padding-bottom: 4px;
    margin: 18px 0 8px;
  }}
  #gm-body h3:first-child {{ margin-top: 0; }}
  #gm-body ul {{ margin: 4px 0 8px 0; padding-left: 18px; }}
  #gm-body li {{ margin-bottom: 4px; }}
  #gm-body p {{ margin: 0 0 8px; }}
  #gm-footer {{
    display: flex; justify-content: space-between; padding: 6px 14px;
    border-top: 1px solid var(--border); font-size: 9px; color: var(--subtext);
    flex-shrink: 0;
  }}
  /* GM Prompts tab */
  #gm-prompts-panel {{ overflow-y: auto; padding: 14px 16px; gap: 0; }}
  .gmp-card {{
    border: 1px solid var(--border); border-radius: 7px; margin-bottom: 12px; overflow: hidden;
  }}
  .gmp-card-header {{
    display: flex; align-items: center; justify-content: space-between; gap: 8px;
    padding: 8px 12px; background: var(--row-alt); border-bottom: 1px solid var(--border);
  }}
  .gmp-card-title {{ font-size: 12px; font-weight: 700; color: var(--text); flex: 1; }}
  .gmp-card-audience {{ font-size: 11px; color: var(--subtext); }}
  .gmp-run-btn {{
    background: var(--blue); color: #fff; border: none; border-radius: 4px;
    padding: 4px 11px; font-size: 11px; font-weight: 600; cursor: pointer; flex-shrink: 0;
  }}
  .gmp-run-btn:hover {{ opacity: .85; }}
  .gmp-run-btn:disabled {{ opacity: .4; cursor: default; }}
  .gmp-chat-btn {{ font-size: 11px; padding: 3px 9px; border-radius: 5px; border: 1px solid var(--blue); background: transparent; color: var(--blue); cursor: pointer; white-space: nowrap; }}
  .gmp-chat-btn:hover {{ background: var(--blue); color: #fff; }}
  .gmp-preview {{
    padding: 8px 12px; font-size: 11px; line-height: 1.55; color: var(--subtext);
    white-space: pre-wrap; font-family: inherit; max-height: 80px; overflow: hidden;
    position: relative; background: #fff;
  }}
  .gmp-preview::after {{
    content: ''; position: absolute; bottom: 0; left: 0; right: 0; height: 28px;
    background: linear-gradient(transparent, #fff);
  }}
  /* GM Chat tab */
  #gm-chat-panel {{ flex-direction: column; }}
  #gm-chat-thread {{
    flex: 1; min-height: 80px; overflow-y: auto; padding: 12px 16px; display: flex; flex-direction: column; gap: 10px;
  }}
  .chat-msg {{ display: flex; gap: 8px; align-items: flex-start; }}
  .chat-msg.user {{ flex-direction: row-reverse; }}
  .chat-bubble {{
    max-width: 82%; padding: 8px 12px; border-radius: 10px; font-size: 12.5px; line-height: 1.6;
    white-space: pre-wrap;
  }}
  .chat-msg.user .chat-bubble {{ background: var(--blue); color: #fff; border-radius: 10px 2px 10px 10px; }}
  .chat-msg.ai .chat-bubble {{
    background: var(--row-alt); color: var(--text); border: 1px solid var(--border);
    border-radius: 2px 10px 10px 10px;
  }}
  .chat-msg.ai .chat-bubble h3 {{
    font-size: 11px; font-weight: 700; text-transform: uppercase; color: var(--blue);
    border-bottom: 1px solid var(--border); padding-bottom: 3px; margin: 12px 0 6px;
  }}
  .chat-msg.ai .chat-bubble h3:first-child {{ margin-top: 0; }}
  .chat-msg.ai .chat-bubble ul {{ margin: 4px 0 6px; padding-left: 16px; }}
  .chat-msg.ai .chat-bubble li {{ margin-bottom: 3px; }}
  .chat-msg.ai .chat-bubble p {{ margin: 0 0 6px; }}
  .chat-avatar {{
    width: 26px; height: 26px; border-radius: 50%; flex-shrink: 0;
    display: flex; align-items: center; justify-content: center;
    font-size: 12px; font-weight: 700;
  }}
  .chat-msg.user .chat-avatar {{ background: var(--blue); color: #fff; }}
  .chat-msg.ai .chat-avatar {{ background: var(--border); color: var(--subtext); }}
  .chat-thinking {{
    display: flex; gap: 4px; padding: 10px 14px; align-items: center;
  }}
  .chat-thinking span {{
    width: 7px; height: 7px; border-radius: 50%; background: var(--subtext);
    animation: gm-shimmer 1.2s infinite;
  }}
  .chat-thinking span:nth-child(2) {{ animation-delay: .2s; }}
  .chat-thinking span:nth-child(3) {{ animation-delay: .4s; }}
  #gm-chat-input-row {{
    display: flex; gap: 8px; padding: 10px 14px; border-top: 1px solid var(--border);
    flex-shrink: 0; align-items: flex-end;
  }}
  #gm-chat-input {{
    flex: 1; border: 1px solid var(--border); border-radius: 8px;
    padding: 8px 12px; font-size: 13px; font-family: inherit; resize: none;
    min-height: 68px; max-height: 40vh; overflow-y: auto; line-height: 1.5;
    outline: none; transition: max-height .1s ease;
  }}
  #gm-chat-input:focus {{ border-color: var(--blue); box-shadow: 0 0 0 2px rgba(1,118,211,.12); }}
  #gm-chat-send {{
    background: var(--blue); color: #fff; border: none; border-radius: 8px;
    padding: 8px 14px; font-size: 13px; font-weight: 600; cursor: pointer; flex-shrink: 0;
  }}
  #gm-chat-send:hover {{ opacity: .88; }}
  #gm-chat-send:disabled {{ opacity: .4; cursor: default; }}
  #gm-chat-clear {{
    background: none; color: var(--subtext); border: 1px solid var(--border); border-radius: 8px;
    padding: 8px 10px; font-size: 11px; cursor: pointer; flex-shrink: 0;
  }}
  #gm-chat-clear:hover {{ color: var(--text); border-color: var(--text); }}
  /* GM Library tab */
  #gm-library-panel {{ flex-direction: column; min-height: 0; }}
  #gm-library-filter-bar {{
    display: flex; gap: 6px; padding: 10px 14px; border-bottom: 1px solid var(--border);
    flex-shrink: 0; flex-wrap: wrap; align-items: center;
  }}
  .gml-chip {{
    font-size: 11px; font-weight: 600; padding: 3px 10px; border-radius: 12px;
    border: 1px solid var(--border); background: transparent; color: var(--subtext);
    cursor: pointer; transition: background .12s, color .12s;
  }}
  .gml-chip.active {{ background: var(--blue); color: #fff; border-color: var(--blue); }}
  #gm-library-list {{ flex: 1; overflow-y: auto; padding: 12px 14px; }}
  #gm-library-empty {{
    text-align: center; padding: 40px 20px; color: var(--subtext); font-size: 13px;
  }}
  .gml-card {{
    border: 1px solid var(--border); border-radius: 7px; margin-bottom: 10px; overflow: hidden;
  }}
  .gml-card-header {{
    display: flex; align-items: center; justify-content: space-between; gap: 8px;
    padding: 8px 12px; background: var(--row-alt); border-bottom: 1px solid var(--border);
  }}
  .gml-card-title {{ font-size: 12px; font-weight: 700; color: var(--text); flex: 1; min-width: 0; }}
  .gml-audience-badge {{
    font-size: 10px; font-weight: 600; padding: 2px 7px; border-radius: 10px;
    white-space: nowrap; flex-shrink: 0;
  }}
  .gml-audience-badge.portfolio_owner {{ background: #e8f4fd; color: #0176d3; }}
  .gml-audience-badge.dam_lead {{ background: #fef3cd; color: #916600; }}
  .gml-audience-badge.gm {{ background: #e8f5e9; color: #2e7d32; }}
  .gml-audience-badge.evp {{ background: #fce4ec; color: #c62828; }}
  .gml-card-actions {{ display: flex; gap: 5px; align-items: center; flex-shrink: 0; }}
  .gml-btn {{
    font-size: 11px; padding: 3px 9px; border-radius: 5px; cursor: pointer; border: none; font-weight: 600;
  }}
  .gml-btn-run {{ background: var(--blue); color: #fff; }}
  .gml-btn-run:hover {{ opacity: .85; }}
  .gml-btn-edit {{ background: transparent; border: 1px solid var(--border) !important; color: var(--subtext); }}
  .gml-btn-edit:hover {{ border-color: var(--blue) !important; color: var(--blue); }}
  .gml-btn-del {{ background: transparent; border: 1px solid var(--border) !important; color: var(--subtext); }}
  .gml-btn-del:hover {{ border-color: #d32f2f !important; color: #d32f2f; }}
  .gml-preview {{
    padding: 8px 12px; font-size: 11px; line-height: 1.55; color: var(--subtext);
    white-space: pre-wrap; font-family: inherit; max-height: 70px; overflow: hidden;
    position: relative; background: #fff;
  }}
  .gml-preview::after {{
    content: ''; position: absolute; bottom: 0; left: 0; right: 0; height: 24px;
    background: linear-gradient(transparent, #fff);
  }}
  /* Save-prompt dialog */
  #save-prompt-modal {{ display: none; position: fixed; inset: 0; z-index: 10001; align-items: center; justify-content: center; }}
  #save-prompt-modal.open {{ display: flex; }}
  #save-prompt-backdrop {{ position: absolute; inset: 0; background: rgba(0,0,0,.45); }}
  #save-prompt-dialog {{
    position: relative; z-index: 1; background: var(--card); border-radius: 9px;
    width: min(480px, 94vw); box-shadow: 0 6px 32px rgba(0,0,0,.28);
    display: flex; flex-direction: column;
  }}
  #save-prompt-header {{
    background: var(--blue); color: #fff; border-radius: 9px 9px 0 0;
    padding: 10px 14px; font-size: 13px; font-weight: 700; display: flex;
    align-items: center; justify-content: space-between;
  }}
  #save-prompt-body {{ padding: 16px; display: flex; flex-direction: column; gap: 12px; }}
  .spdlg-label {{ font-size: 12px; font-weight: 600; color: var(--text); margin-bottom: 4px; }}
  #save-prompt-title-inp {{
    width: 100%; border: 1px solid var(--border); border-radius: 6px;
    padding: 7px 10px; font-size: 13px; font-family: inherit; outline: none; box-sizing: border-box;
  }}
  #save-prompt-title-inp:focus {{ border-color: var(--blue); box-shadow: 0 0 0 2px rgba(1,118,211,.1); }}
  #save-prompt-audience {{
    width: 100%; border: 1px solid var(--border); border-radius: 6px;
    padding: 7px 10px; font-size: 13px; font-family: inherit; outline: none;
    background: #fff; cursor: pointer; box-sizing: border-box;
  }}
  #save-prompt-audience:focus {{ border-color: var(--blue); }}
  #save-prompt-footer {{
    display: flex; gap: 8px; justify-content: flex-end;
    padding: 10px 16px; border-top: 1px solid var(--border);
  }}
  /* Save button in chat row */
  #gm-chat-save {{
    background: none; color: var(--subtext); border: 1px solid var(--border); border-radius: 8px;
    padding: 8px 10px; font-size: 11px; cursor: pointer; flex-shrink: 0;
  }}
  #gm-chat-save:hover {{ color: var(--blue); border-color: var(--blue); }}
  .gm-slack-bar {{
    padding: 7px 14px; border-top: 1px solid var(--border); background: var(--row-alt);
    flex-shrink: 0;
  }}
  .gm-slack-chk {{
    display: flex; align-items: center; gap: 7px; cursor: pointer;
    font-size: 12px; font-weight: 500; color: var(--text); user-select: none;
  }}
  .gm-slack-chk input {{ cursor: pointer; margin: 0; }}
  .gm-slack-note {{ font-size: 10px; color: var(--subtext); font-weight: 400; }}
  .gm-slack-status {{
    font-size: 11px; color: var(--subtext); padding: 4px 14px 6px;
    border-bottom: 1px solid var(--border); background: var(--row-alt);
    flex-shrink: 0; font-style: italic;
  }}
  .gm-loading {{ animation: gm-shimmer 1.4s infinite; }}
  @keyframes gm-shimmer {{
    0%,100% {{ opacity: .35; }} 50% {{ opacity: .7; }}
  }}
  .gm-skeleton-line {{
    height: 11px; background: var(--border); border-radius: 4px;
    margin-bottom: 8px;
  }}
  .gm-spinner-wrap {{
    display: flex; flex-direction: column; align-items: center; justify-content: center;
    gap: 14px; padding: 48px 20px; flex: 1;
  }}
  .gm-spinner {{
    width: 38px; height: 38px; border-radius: 50%;
    border: 3px solid var(--border); border-top-color: var(--blue);
    animation: gm-spin .75s linear infinite;
  }}
  @keyframes gm-spin {{ to {{ transform: rotate(360deg); }} }}
  .gm-spinner-label {{
    font-size: 13px; color: var(--subtext); font-weight: 500;
  }}
</style>
</head>
<body>
<!-- Assignment modal -->
<div id="assign-modal">
  <div id="assign-backdrop" onclick="closeAssignModal()"></div>
  <div id="assign-dialog">
    <h3>Edit Assignment</h3>
    <div class="assign-acct" id="assign-acct"></div>
    <label for="assign-tier">Tier</label>
    <select id="assign-tier">
      <option value="1">T1 — Tier 1</option>
      <option value="2">T2 — Tier 2</option>
      <option value="3">T3 — Tier 3</option>
    </select>
    <label for="assign-po">Portfolio Owner</label>
    <select id="assign-po"></select>
    <button class="assign-roster-toggle" onclick="_toggleRoster()">✎ Manage PO list for this region</button>
    <div id="assign-roster-panel">
      <div class="roster-region-label" id="roster-region-label"></div>
      <div class="roster-list" id="roster-list"></div>
      <div class="roster-add-row">
        <input id="roster-add-input" type="text" placeholder="Full name…" onkeydown="if(event.key==='Enter')_addPoToRoster()">
        <button onclick="_addPoToRoster()">+ Add</button>
      </div>
    </div>
    <div class="assign-actions">
      <button class="btn-cancel" onclick="closeAssignModal()">Cancel</button>
      <button class="btn-save" onclick="commitAssign()">Save</button>
    </div>
  </div>
</div>
<div id="pulse-tooltip">
  <div id="pulse-tooltip-header" class="popup-drag-handle">
    <strong id="pulse-tooltip-title"></strong>
    <span style="display:flex;align-items:center;gap:8px">
      <a id="pulse-tooltip-link" href="#" target="_blank" style="display:none">↗ Open Pulse in Salesforce</a>
      <button id="pulse-tooltip-close" onclick="_hidePulse(true)">✕</button>
    </span>
  </div>
  <div id="pulse-tooltip-body"></div>
</div>
<div id="far-tooltip">
  <div id="far-tooltip-header" class="popup-drag-handle">
    <strong id="far-tooltip-title">FAR Details</strong>
    <button id="far-tooltip-close" onclick="_hideFar(true)">✕</button>
  </div>
  <div id="far-tooltip-body"></div>
</div>
<div id="p2g-modal">
  <div id="p2g-modal-header" class="popup-drag-handle">
    <strong id="p2g-modal-title">Path to Green</strong>
    <button id="p2g-modal-close" onclick="_hideP2g(true)">✕</button>
  </div>
  <div id="p2g-modal-body"></div>
</div>
<div id="res-modal">
  <div id="res-modal-header" class="popup-drag-handle">
    <strong id="res-modal-title">Assigned Resources</strong>
    <button id="res-modal-close" onclick="_hideRes(true)">✕</button>
  </div>
  <div id="res-modal-body"></div>
</div>
<div id="pipe-modal">
  <div id="pipe-header" class="popup-drag-handle">
    <strong id="pipe-title"></strong>
    <button id="pipe-close" onclick="_hidePipe(true)">✕</button>
  </div>
  <div id="pipe-body"></div>
</div>

<div class="page-header">
  <div>
    <h1>ACC Delivery Portfolio</h1>
    <div class="meta">Data Refresh: <span id="data-refresh-date">loading…</span> ET{'&nbsp;|&nbsp;v' + html_version if html_version else ''}</div>
  </div>
  <div style="text-align:right;font-size:12px;opacity:.8" id="header-summary">
    {len(results)} projects &nbsp;|&nbsp; {fmt_m(total_bk)} bookings
  </div>
</div>

<div class="scorecard-bar" onclick="toggleScorecard()">
  <span class="sc-summary" id="sc-bar-summary">Loading…</span>
  <span class="sc-toggle"><span id="sc-chevron">▲</span> Scorecard</span>
</div>
<div class="scorecard" id="scorecard">
  <div class="stat"><div class="lbl">Projects</div><div class="val" id="sc-total">{len(results)}</div></div>
  <div class="stat"><div class="lbl">Bookings</div><div class="val" id="sc-bookings">{fmt_m(total_bk)}</div></div>
  <div class="stat">
    <div class="lbl">Avg Bookings</div>
    <div class="val" id="sc-avg-bk">{fmt_m(total_bk / len(results) if results else 0)}</div>
    <div style="margin-top:4px;font-size:11px;color:var(--subtext)">Median: <span id="sc-median-bk" style="font-weight:600;color:var(--text)">{fmt_m(median_bk)}</span></div>
  </div>
  <div class="stat"><div class="lbl">Billings</div><div class="val" id="sc-billings">{fmt_m(total_bil)}</div></div>
  <div class="stat"><div class="lbl">Backlog</div><div class="val" id="sc-backlog">{fmt_m(total_bk - total_bil)}</div></div>
  <div class="stat" style="min-width:160px">
    <div class="lbl" style="margin-bottom:6px">Project Status</div>
    <div style="font-size:12px;line-height:2">
      <div>🔴 Red: <strong id="sc-red">{len(reds)}</strong></div>
      <div>🟡 Yellow: <strong id="sc-yellow">{len(yellows)}</strong></div>
      <div>🟢 Green: <strong id="sc-green-total">{len(clean_green) + len(watermelons)}</strong> <span style="font-size:11px;color:var(--subtext)">(<span id="sc-wm">{len(watermelons)}</span> 🍉, <span id="sc-green">{len(clean_green)}</span> ✅)</span></div>
      <div>⚫ No Pulse: <strong id="sc-nopulse">{len(no_pulse)}</strong></div>
    </div>
  </div>
  <div class="stat"><div class="lbl">Total FAR</div><div class="val" id="sc-far">{fmt_m(total_far)}</div></div>
  <div class="stat red"><div class="lbl">Overdue Inv</div><div class="val" id="sc-overdue">{fmt_m(total_overdue)}</div></div>
  <div class="stat yellow"><div class="lbl">RR Rev Risk</div><div class="val" id="sc-rr">{fmt_m(total_rr_rev)}</div></div>
  <div class="stat"><div class="lbl">Wtd Bid Margin</div><div class="val" id="sc-bid" style="color:{'var(--green)' if w_bid > 13 else 'var(--yellow)' if w_bid >= 5 else 'var(--red)'}">{w_bid:.1f}%</div></div>
  <div class="stat"><div class="lbl">Wtd Margin at Close</div><div class="val" id="sc-close" style="color:{'var(--red)' if w_close < 0 else ''}">{w_close:.1f}%</div></div>
  <div class="stat" title="Avg Margin Efficiency = Weighted Margin @ Close − Weighted Bid Margin"><div class="lbl">Margin Delta</div><div class="val" id="sc-delta">{w_close - w_bid:+.1f}%</div></div>
  <div class="stat {'red' if avg_hr is not None and avg_hr < 60 else 'yellow' if avg_hr is not None and avg_hr < 75 else 'green' if avg_hr is not None else ''}"><div class="lbl">Avg Health Score</div><div class="val" id="sc-avg-hr">{f'{avg_hr:.0f}' if avg_hr is not None else '—'}</div></div>
  <div class="stat {'red' if avg_dq is not None and avg_dq < 60 else 'yellow' if avg_dq is not None and avg_dq < 75 else 'green' if avg_dq is not None else ''}"><div class="lbl">Avg Data Quality</div><div class="val" id="sc-avg-dq">{f'{avg_dq:.0f}' if avg_dq is not None else '—'}</div></div>
  <div class="stat {'red' if avg_csat is not None and avg_csat < 3.5 else 'yellow' if avg_csat is not None and avg_csat < 4.0 else 'green' if avg_csat is not None else ''}"><div class="lbl">Avg CSAT</div><div class="val" id="sc-avg-csat">{f'{avg_csat:.1f}' if avg_csat is not None else '—'}</div></div>
</div>

<div class="region-tabs">
  <button class="rtab active" onclick="setRegion(null,this)">ACC (All)</button>
  <button class="rtab" onclick="setRegion('AMER TMT',this)">TMT</button>
  <button class="rtab" onclick="setRegion('AMER CBS',this)">CBS</button>
</div>

<div class="controls">
  <!-- Row 1: search + filters + clear -->
  <div class="controls-row controls-row-1">
    <input type="search" id="search" placeholder="🔍  Search project, account, PM, rules…" oninput="applyFilters()" onsearch="applyFilters()">
    <div class="sep"></div>
    <div class="filter-group">
      <label style="font-size:12px;color:var(--subtext)">Status:</label>
      <div class="ms-wrap" id="ms-status-wrap">
        <button class="ms-btn" id="ms-status-btn" onclick="toggleMs('status')">All Statuses</button>
        <div class="ms-dropdown" id="ms-status-dd">
          <label class="ms-item"><input type="checkbox" value="Red" onchange="msChange('status')"> 🔴 Red</label>
          <label class="ms-item"><input type="checkbox" value="Yellow" onchange="msChange('status')"> 🟡 Yellow</label>
          <label class="ms-item"><input type="checkbox" value="Watermelon" onchange="msChange('status')"> 🍉 Watermelon</label>
          <label class="ms-item"><input type="checkbox" value="Green" onchange="msChange('status')"> 🟢 Green</label>
          <label class="ms-item"><input type="checkbox" value="No Pulse" onchange="msChange('status')"> ⚫ No Pulse</label>
          <label class="ms-item"><input type="checkbox" value="On Hold" onchange="msChange('status')"> ⏸ On Hold</label>
          <span class="ms-clear" onclick="msClear('status')">Clear</span>
        </div>
      </div>
    </div>
    <div class="sep"></div>
    <div class="filter-group">
      <label style="font-size:12px;color:var(--subtext)">Type:</label>
      <div class="ms-wrap" id="ms-type-wrap">
        <button class="ms-btn" id="ms-type-btn" onclick="toggleMs('type')">All Types</button>
        <div class="ms-dropdown" id="ms-type-dd">
          <label class="ms-item"><input type="checkbox" value="hw" onchange="msChange('type')"> ⚑ High Watch</label>
          <label class="ms-item"><input type="checkbox" value="swe" onchange="msChange('type')"> 🔧 SWE</label>
          <label class="ms-item"><input type="checkbox" value="ari" onchange="msChange('type')"> 🏷 ARI</label>
          <label class="ms-item"><input type="checkbox" value="not_started" onchange="msChange('type')"> 🗓 Not Started</label>
          <span class="ms-clear" onclick="msClear('type')">Clear</span>
        </div>
      </div>
    </div>
    <div class="sep"></div>
    <div class="filter-group">
      <label style="font-size:12px;color:var(--subtext)">Billing:</label>
      <div class="ms-wrap" id="ms-billing-wrap">
        <button class="ms-btn" id="ms-billing-btn" onclick="toggleMs('billing')">All Billing</button>
        <div class="ms-dropdown" id="ms-billing-dd">
          <label class="ms-item"><input type="checkbox" value="Fixed Price + Expenses" onchange="msChange('billing')"> Fixed Price + Exp</label>
          <label class="ms-item"><input type="checkbox" value="Fixed Price POC + Expenses" onchange="msChange('billing')"> Fixed Price POC</label>
          <label class="ms-item"><input type="checkbox" value="Time and Materials" onchange="msChange('billing')"> T&amp;M</label>
          <label class="ms-item"><input type="checkbox" value="Time and Materials + Holdback" onchange="msChange('billing')"> T&amp;M + Holdback</label>
          <span class="ms-clear" onclick="msClear('billing')">Clear</span>
        </div>
      </div>
    </div>
    <div class="sep"></div>
    <div class="filter-group">
      <label style="font-size:12px;color:var(--subtext)">Tier:</label>
      <div class="ms-wrap" id="ms-tier-wrap">
        <button class="ms-btn" id="ms-tier-btn" onclick="toggleMs('tier')">All Tiers</button>
        <div class="ms-dropdown" id="ms-tier-dd">
          <label class="ms-item"><input type="checkbox" value="1" onchange="msChange('tier')"> T1 Strategic</label>
          <label class="ms-item"><input type="checkbox" value="2" onchange="msChange('tier')"> T2 Growth</label>
          <label class="ms-item"><input type="checkbox" value="3" onchange="msChange('tier')"> T3 Volume</label>
          <span class="ms-clear" onclick="msClear('tier')">Clear</span>
        </div>
      </div>
    </div>
    <div class="sep"></div>
    <div class="filter-group">
      <label style="font-size:12px;color:var(--subtext)">PO:</label>
      <div class="ms-wrap" id="ms-po-wrap">
        <button class="ms-btn" id="ms-po-btn" onclick="toggleMs('po')">All Owners</button>
        <div class="ms-dropdown" id="ms-po-dd">
          {chr(10).join(f'<label class="ms-item"><input type="checkbox" value="{_html.escape(p)}" onchange="msChange(\'po\')"> {_html.escape(p)}</label>' for p in po_list)}
          <span class="ms-clear" onclick="msClear('po')">Clear</span>
        </div>
      </div>
    </div>
    <div class="sep"></div>
    <div class="filter-group">
      <label style="font-size:12px;color:var(--subtext)">Rules:</label>
      <div class="ms-wrap" id="ms-rule-wrap">
        <button class="ms-btn" id="ms-rule-btn" onclick="toggleMs('rule')">All Rules</button>
        <div class="ms-dropdown" id="ms-rule-dd">
          <label class="ms-item"><input type="checkbox" value="MARGIN_RED" onchange="msChange('rule')"> MARGIN_RED</label>
          <label class="ms-item"><input type="checkbox" value="MARGIN_YELLOW" onchange="msChange('rule')"> MARGIN_YELLOW</label>
          <label class="ms-item"><input type="checkbox" value="FAR_RED_NEG" onchange="msChange('rule')"> FAR_RED_NEG</label>
          <label class="ms-item"><input type="checkbox" value="FAR_RED_UNDERUTIL" onchange="msChange('rule')"> FAR_RED_UNDERUTIL</label>
          <label class="ms-item"><input type="checkbox" value="FAR_YELLOW" onchange="msChange('rule')"> FAR_YELLOW</label>
          <label class="ms-item"><input type="checkbox" value="RR_RISK" onchange="msChange('rule')"> RR_RISK</label>
          <label class="ms-item"><input type="checkbox" value="GDC_LOW" onchange="msChange('rule')"> GDC_LOW</label>
          <label class="ms-item"><input type="checkbox" value="OVERDUE_INV" onchange="msChange('rule')"> OVERDUE_INV</label>
          <label class="ms-item"><input type="checkbox" value="SWE_BURNING_HOT" onchange="msChange('rule')"> SWE_BURNING_HOT</label>
          <label class="ms-item"><input type="checkbox" value="NO_PULSE" onchange="msChange('rule')"> NO_PULSE</label>
          <label class="ms-item"><input type="checkbox" value="NO_STEERCO" onchange="msChange('rule')"> NO_STEERCO</label>
          <label class="ms-item"><input type="checkbox" value="MISSING_PTG" onchange="msChange('rule')"> MISSING_PTG</label>
          <label class="ms-item"><input type="checkbox" value="END_DATE_PAST" onchange="msChange('rule')"> END_DATE_PAST</label>
          <label class="ms-item"><input type="checkbox" value="END_DATE_30" onchange="msChange('rule')"> END_DATE_30 (≤30 days)</label>
          <label class="ms-item"><input type="checkbox" value="END_DATE_60" onchange="msChange('rule')"> END_DATE_60 (31–60 days)</label>
          <label class="ms-item"><input type="checkbox" value="END_DATE_90" onchange="msChange('rule')"> END_DATE_90 (61–90 days)</label>
          <label class="ms-item"><input type="checkbox" value="END_DATE_120" onchange="msChange('rule')"> END_DATE_120 (91–120 days)</label>
          <label class="ms-item"><input type="checkbox" value="CSAT_OVERDUE" onchange="msChange('rule')"> CSAT_OVERDUE</label>
          <label class="ms-item"><input type="checkbox" value="CSAT_EXEMPT" onchange="msChange('rule')"> CSAT_EXEMPT</label>
          <span class="ms-clear" onclick="msClear('rule')">Clear</span>
        </div>
      </div>
    </div>
    <div class="controls-row-spacer"></div>
    <button id="clear-filters-btn" onclick="clearAllFilters()" title="Clear all filters and search" style="display:none;background:none;border:1px solid var(--border);border-radius:5px;padding:4px 10px;font-size:11px;font-weight:600;color:var(--subtext);cursor:pointer;white-space:nowrap">✕ Clear Filters</button>
  </div>
  <!-- Row 2: group by + collapse/expand + action buttons -->
  <div class="controls-row">
    <div class="filter-group" style="gap:6px">
      <label style="font-size:12px;color:var(--subtext);white-space:nowrap">Group by:</label>
      <div class="ms-wrap" id="ms-grp-wrap">
        <button class="ms-btn" id="ms-grp-btn" onclick="toggleGrpDd()">None</button>
        <div class="ms-dropdown" id="ms-grp-dd" style="min-width:180px">
          <div style="padding:6px 10px 4px;font-size:10px;font-weight:700;color:var(--subtext);text-transform:uppercase;letter-spacing:.5px">Select &amp; order groupings</div>
          <label class="ms-item" id="grp-opt-region"><input type="checkbox" value="region" onchange="grpChange(this)"> Region</label>
          <label class="ms-item" id="grp-opt-tier"><input type="checkbox" value="tier" onchange="grpChange(this)"> Tier</label>
          <label class="ms-item" id="grp-opt-po"><input type="checkbox" value="po" onchange="grpChange(this)"> Portfolio Owner</label>
          <label class="ms-item" id="grp-opt-acct"><input type="checkbox" value="acct" onchange="grpChange(this)"> Account Name</label>
          <label class="ms-item" id="grp-opt-status"><input type="checkbox" value="status" onchange="grpChange(this)"> Status</label>
          <span class="ms-clear" onclick="grpClear()">Clear</span>
        </div>
      </div>
      <span id="grp-order-pills" style="display:flex;gap:3px;flex-wrap:wrap"></span>
      <span id="grp-collapse-btns" style="display:none;gap:4px">
        <button class="grp-btn" id="grp-collapse-btn" onclick="stepCollapse()">⊟ Collapse</button>
        <button class="grp-btn" id="grp-expand-btn"   onclick="stepExpand()">⊞ Expand</button>
      </span>
    </div>
    <div class="controls-row-spacer"></div>
    <button id="gm-overview-btn" onclick="openGMOverview(false)" title="AI-generated GM business overview">📊 Business Overview</button>
    <button id="help-btn" onclick="openHelp()" title="Help &amp; AI Prompts">❓ Help</button>
  </div>
</div>

<!-- GM Business Overview modal -->
<div id="gm-modal">
  <div id="gm-backdrop" onclick="closeGMOverview()"></div>
  <div id="gm-dialog">
    <div id="gm-header">
      <span id="gm-scope-label">ACC Delivery Portfolio — Business Overview</span>
      <div style="display:flex;gap:6px;align-items:center">
        <button id="gm-regen-btn" onclick="openGMOverview(true)">↻ Regenerate</button>
        <button class="btn-cancel" onclick="closeGMOverview()">✕ Close</button>
      </div>
    </div>
    <div id="gm-tabs">
      <div class="gm-tab active" id="gm-tab-overview" onclick="switchGmTab('overview')">📊 Overview</div>
      <div class="gm-tab" id="gm-tab-prompts" onclick="switchGmTab('prompts')">💡 Prompts</div>
      <div class="gm-tab" id="gm-tab-chat" onclick="switchGmTab('chat')">💬 Chat</div>
      <div class="gm-tab" id="gm-tab-library" onclick="switchGmTab('library')">📚 Library</div>
    </div>
    <!-- Overview panel -->
    <div class="gm-tab-panel active" id="gm-panel-overview" style="flex-direction:column;">
      <div id="gm-body"></div>
      <div id="gm-footer">
        <span id="gm-model-used"></span>
        <span id="gm-generated-at"></span>
      </div>
    </div>
    <!-- Prompts panel -->
    <div class="gm-tab-panel" id="gm-panel-prompts">
      <div id="gm-prompts-slack-bar" class="gm-slack-bar">
        <label class="gm-slack-chk">
          <input type="checkbox" id="gm-prompts-slack-chk">
          <span>🔍 Search Slack channels for each project <span class="gm-slack-note">(adds ~10–30s per project found)</span></span>
        </label>
      </div>
      <div id="gm-prompts-panel"></div>
    </div>
    <!-- Chat panel -->
    <div class="gm-tab-panel" id="gm-panel-chat">
      <div id="gm-chat-thread"></div>
      <div id="gm-chat-slack-bar" class="gm-slack-bar">
        <label class="gm-slack-chk">
          <input type="checkbox" id="gm-chat-slack-chk">
          <span>🔍 Include Slack channel summaries (7-day) in context <span class="gm-slack-note">(adds ~10–30s per project found)</span></span>
        </label>
      </div>
      <div id="gm-chat-input-row">
        <textarea id="gm-chat-input" rows="3" placeholder="Ask anything about the portfolio…"></textarea>
        <button id="gm-chat-save" onclick="openSavePrompt()" title="Save prompt to library">💾</button>
        <button id="gm-chat-clear" onclick="gmChatClear()" title="Clear conversation">🗑</button>
        <button id="gm-chat-send" onclick="gmChatSend()">Send</button>
      </div>
    </div>
    <!-- Library panel -->
    <div class="gm-tab-panel" id="gm-panel-library">
      <div id="gm-library-filter-bar">
        <button class="gml-chip active" data-aud="all" onclick="gmlSetFilter('all')">All</button>
        <button class="gml-chip" data-aud="portfolio_owner" onclick="gmlSetFilter('portfolio_owner')">Portfolio Owner</button>
        <button class="gml-chip" data-aud="dam_lead" onclick="gmlSetFilter('dam_lead')">DAM Lead / GM</button>
        <button class="gml-chip" data-aud="gm" onclick="gmlSetFilter('gm')">GM</button>
        <button class="gml-chip" data-aud="evp" onclick="gmlSetFilter('evp')">EVP</button>
      </div>
      <div id="gm-library-list">
        <div id="gm-library-empty" style="display:none">No saved prompts yet — save one from the Chat tab.</div>
      </div>
    </div>
  </div>
</div>

<!-- Save Prompt dialog -->
<div id="save-prompt-modal">
  <div id="save-prompt-backdrop" onclick="closeSavePrompt()"></div>
  <div id="save-prompt-dialog">
    <div id="save-prompt-header">
      💾 Save Prompt to Library
      <button class="btn-cancel" onclick="closeSavePrompt()">✕</button>
    </div>
    <div id="save-prompt-body">
      <div>
        <div class="spdlg-label">Title</div>
        <input id="save-prompt-title-inp" type="text" placeholder="e.g. Q3 GM Weekly Briefing" maxlength="120">
      </div>
      <div>
        <div class="spdlg-label">Audience / Category</div>
        <select id="save-prompt-audience">
          <option value="portfolio_owner">Portfolio Owner</option>
          <option value="dam_lead">DAM Lead / GM</option>
          <option value="gm">GM</option>
          <option value="evp">EVP</option>
        </select>
      </div>
    </div>
    <div id="save-prompt-footer">
      <button class="btn-cancel" onclick="closeSavePrompt()">Cancel</button>
      <button id="save-prompt-ok" onclick="confirmSavePrompt()" style="background:var(--blue);color:#fff;border:none;border-radius:5px;padding:6px 18px;font-size:13px;font-weight:600;cursor:pointer">Save</button>
    </div>
  </div>
</div>

<!-- Help modal -->
<div id="help-modal">
  <div id="help-backdrop" onclick="closeHelp()"></div>
  <div id="help-dialog">
    <div id="help-header">
      <span>❓ Help &amp; AI Prompts</span>
      <button class="btn-cancel" onclick="closeHelp()">✕ Close</button>
    </div>
    <div id="help-tabs">
      <div class="help-tab active" id="help-tab-guide" onclick="switchHelpTab('guide')">How to Use</div>
      <div class="help-tab" id="help-tab-prompts" onclick="switchHelpTab('prompts')">AI Prompts</div>
    </div>
    <div id="help-body">
      <!-- content injected by switchHelpTab() -->
    </div>
  </div>
</div>

<div id="summary-modal">
  <div id="summary-modal-header" class="popup-drag-handle">
    <strong id="summary-modal-title">Pulse Summary</strong>
    <button id="summary-modal-close" onclick="_hideSummary(true)">✕</button>
  </div>
  <div id="summary-modal-body"></div>
  <div id="summary-modal-footer"></div>
</div>

<div class="table-wrap">
  <table id="main-table">
    <thead>
      <tr>
        <th class="col-status"  data-col="status">Status<span class="sort-arrow"></span></th>
        <th class="col-tier"    data-col="tier">T<span class="sort-arrow"></span></th>
        <th class="col-project" data-col="name">Project / Account<span class="sort-arrow"></span></th>
        <th class="col-team"     data-col="team">Team Leadership<span class="sort-arrow"></span></th>
        <th class="col-resource" data-col="gdc_pct">Resourcing<span class="sort-arrow"></span></th>
        <th class="col-fin"      data-col="bookings">Financials<span class="sort-arrow"></span></th>
        <th class="col-rules"   data-col="rules">Portfolio Assurance<span class="sort-arrow"></span></th>
        <th class="col-bl"      data-col="baselines">Baselines<span class="sort-arrow"></span></th>
        <th class="col-summary" data-col="summary" id="th-summary">Summary</th>
      </tr>
    </thead>
    <tbody id="table-body"></tbody>
  </table>
  <div class="no-results" id="no-results" style="display:none">No projects match your filters.</div>
</div>

<script>
// Resource popup: open on mouseenter, close on mouseleave or close button
// ── Resource popup ─────────────────────────────────────────────────────────────
const _resMod   = document.getElementById('res-modal');
const _resTitle = document.getElementById('res-modal-title');
const _resBody  = document.getElementById('res-modal-body');
const _resPopup = _makePopup(_resMod, true);
_makeDraggable(_resMod, document.getElementById('res-modal-header'));
let _resSortCol = 'name', _resSortDir = 1, _resCurrentData = null;

function _renderResTable(data) {{
  const col = _resSortCol, dir = _resSortDir;
  const rows = data.resources.slice().sort((a,b) => {{
    const av = a[col] ?? '', bv = b[col] ?? '';
    return typeof av === 'number' ? dir*(av-bv) : dir*String(av).localeCompare(String(bv));
  }});
  const fmtH = h => (h != null && h !== 0) ? Math.round(h).toLocaleString('en-US') : '—';
  const totEst = data.resources.reduce((s,x)=>s+(x.est_hrs||0),0);
  const totAct = data.resources.reduce((s,x)=>s+(x.act_hrs||0),0);
  const totRem = data.resources.reduce((s,x)=>s+(x.rem_hrs||0),0);
  const totSch = data.resources.reduce((s,x)=>s+(x.sch_hrs||0),0);
  const arrow = c => c===col ? (dir===1?' ▲':' ▼') : ' ⇅';
  const th = (c,lbl,align) => `<th onclick="_resSortBy('${{c}}')" style="text-align:${{align||'left'}}">${{lbl}}${{arrow(c)}}</th>`;
  const bodyRows = rows.map((res,i) => {{
    const flag = res.region === 'GDC India' ? '🇮🇳' : '🌎';
    const remStyle = res.rem_hrs != null && res.rem_hrs < 0 ? ' style="color:var(--red)"' : '';
    const bg = i%2===0?'':'background:#f8f9fa';
    return `<tr style="${{bg}}">
      <td>${{res.name||'—'}}</td><td>${{res.role||'—'}}</td>
      <td>${{flag}} ${{res.region||'—'}}</td>
      <td style="text-align:right">${{fmtH(res.est_hrs)}}</td>
      <td style="text-align:right">${{fmtH(res.act_hrs)}}</td>
      <td style="text-align:right"${{remStyle}}>${{fmtH(res.rem_hrs)}}</td>
      <td style="text-align:right">${{fmtH(res.sch_hrs)}}</td>
      <td>${{res.start||'—'}}</td><td>${{res.end||'—'}}</td>
    </tr>`;
  }}).join('');
  const remTotStyle = totRem<0?' style="color:var(--red)"':'';
  const totRow = `<tr style="font-weight:700;border-top:2px solid var(--border);background:#f1f3f9">
    <td colspan="3" style="text-align:right;color:var(--subtext);font-size:10px;padding-right:4px">TOTAL</td>
    <td style="text-align:right">${{fmtH(totEst)}}</td><td style="text-align:right">${{fmtH(totAct)}}</td>
    <td style="text-align:right"${{remTotStyle}}>${{fmtH(totRem)}}</td>
    <td style="text-align:right">${{fmtH(totSch)}}</td><td colspan="2"></td>
  </tr>`;
  _resBody.innerHTML = `<div style="overflow-x:auto"><table>
    <thead><tr>${{th('name','Resource')}}${{th('role','Role')}}${{th('region','Region')}}
      ${{th('est_hrs','Planned','right')}}${{th('act_hrs','Actual','right')}}
      ${{th('rem_hrs','Remaining','right')}}${{th('sch_hrs','Scheduled','right')}}
      ${{th('start','Start')}}${{th('end','End')}}</tr></thead>
    <tbody>${{bodyRows}}${{totRow}}</tbody>
  </table></div>
  <div style="margin-top:6px;font-size:10px;color:var(--subtext)">🇮🇳 GDC India &nbsp; 🌎 Other regions &nbsp;·&nbsp; Click column headers to sort</div>`;
}}
function _resSortBy(col) {{
  if (_resSortCol===col) {{ _resSortDir*=-1; }} else {{ _resSortCol=col; _resSortDir=1; }}
  if (_resCurrentData) _renderResTable(_resCurrentData);
}}
function _showRes(el) {{
  _resPopup.show(() => {{
    try {{
      const data = JSON.parse(el.dataset.res.replace(/&quot;/g,'"'));
      _resTitle.innerHTML = data.name + `<span style="font-size:11px;font-weight:400;color:rgba(255,255,255,.75);margin-left:8px">Assigned Resources (${{data.resources.length}})</span>`;
      _resCurrentData = data;
      _resSortCol = 'name'; _resSortDir = 1;
      _renderResTable(data);
      _resMod.classList.add('open');
    }} catch(e) {{}}
  }});
}}
function _pinRes(el) {{ _resPopup.pin(() => _showRes(el)); }}
function _hideRes(force) {{ _resPopup.hide(force===true); }}

// Legacy stubs (referenced by old inline handlers — no-op now)
function resOpen(id) {{}}
function resClose(id) {{}}

function toggleScorecard() {{
  const sc = document.getElementById('scorecard');
  const ch = document.getElementById('sc-chevron');
  const collapsed = sc.classList.toggle('collapsed');
  ch.textContent = collapsed ? '▼' : '▲';
}}

const _INLINE = {rows_json};
let RAW = _INLINE.rows || [];

async function loadData() {{
  // Try to load fresher data from _data/ (works when platform routes it correctly)
  try {{
    const [tmt, cbs] = await Promise.all([
      fetch('./_data/acc_amer_tmt_data.json', {{credentials:'include'}}).then(r=>r.ok?r.json():null).catch(()=>null),
      fetch('./_data/acc_amer_cbs_data.json', {{credentials:'include'}}).then(r=>r.ok?r.json():null).catch(()=>null),
    ]);
    const fetched = [tmt, cbs].filter(Boolean).flatMap(d => d.rows || []);
    if (fetched.length > 0) {{
      RAW = fetched;
      const latest = [tmt, cbs].filter(Boolean).map(d => d.generated).filter(Boolean).sort().pop() || '';
      document.getElementById('data-refresh-date').textContent = latest;
    }} else {{
      document.getElementById('data-refresh-date').textContent = (_INLINE.generated || '{REPORT_DATE}');
    }}
  }} catch(e) {{
    document.getElementById('data-refresh-date').textContent = (_INLINE.generated || '{REPORT_DATE}');
  }}
  try {{ await loadAssignments(); }} catch(e) {{ console.warn('loadAssignments (non-fatal):', e); }}
  try {{ await loadPoRoster(); }} catch(e) {{ console.warn('loadPoRoster (non-fatal):', e); }}
  // Restore last region from localStorage
  try {{
    const savedRegion = localStorage.getItem(_REGION_KEY);
    if (savedRegion) {{
      filterRegion = savedRegion;
      document.querySelectorAll('.rtab').forEach(b => {{
        b.classList.toggle('active', b.textContent.trim() === savedRegion || b.getAttribute('onclick').includes("'" + savedRegion + "'"));
      }});
    }}
  }} catch(e) {{}}
  applyFilters();
}}

// ── DB credentials (embedded at publish time) ─────────────────────────────────
const _DB_CRED  = '{_db_b64}';
const _TILE_ID  = {_tile_id_h};

// ── Self-serve assignments ────────────────────────────────────────────────────
async function loadAssignments() {{
  if (!_DB_CRED) return;
  try {{
    // serverQuery hits the live server-side DB (same one db/write updates)
    const data = await TileDB.serverQuery('SELECT pid, tier, po FROM assignments');
    const rows = data.rows || [];
    rows.forEach(a => {{
      const r = RAW.find(x => x.pid === a.pid);
      if (r) {{
        if (a.tier != null) r.tier = parseInt(a.tier);
        if (a.po   != null) r.po   = a.po;
      }}
    }});
  }} catch(e) {{ console.warn('loadAssignments:', e); }}
}}

// Always saves both tier + po together so a partial upsert never wipes the other field
async function saveAssignment(pid, tier, po) {{
  if (!_DB_CRED) return;
  try {{
    const resp = await fetch('/api/tiles/' + _TILE_ID + '/db/write', {{
      method: 'POST',
      headers: {{ 'Authorization': 'Basic ' + _DB_CRED, 'Content-Type': 'application/json' }},
      body: JSON.stringify({{
        table: 'assignments', mode: 'upsert',
        key_column: 'pid', columns: ['pid', 'tier', 'po'],
        rows: [{{ pid, tier: parseInt(tier), po: po }}],
      }}),
    }});
    if (!resp.ok) {{
      const errText = await resp.text().catch(() => resp.status);
      console.warn('saveAssignment HTTP', resp.status, errText);
      return;
    }}
    const r = RAW.find(x => x.pid === pid);
    if (r) {{ r.tier = parseInt(tier); r.po = po; }}
    applyFilters();
  }} catch(e) {{ console.warn('saveAssignment:', e); }}
}}

// _poRoster: region → sorted string[]  (in-memory, loaded from DB + seeded from PO_BY_REGION)
const PO_BY_REGION = {po_by_region_js};
let _poRoster = {{}};  // populated by loadPoRoster()

async function loadPoRoster() {{
  if (!_DB_CRED) {{ _poRoster = JSON.parse(JSON.stringify(PO_BY_REGION)); return; }}
  try {{
    const data = await TileDB.serverQuery('SELECT region, name FROM po_roster ORDER BY name');
    const rows = data.rows || [];
    if (rows.length > 0) {{
      const m = {{}};
      rows.forEach(r => {{ (m[r.region] = m[r.region] || []).push(r.name); }});
      _poRoster = m;
    }} else {{
      // First run: seed DB from publish-time PO_BY_REGION
      _poRoster = JSON.parse(JSON.stringify(PO_BY_REGION));
      const seedRows = [];
      Object.entries(PO_BY_REGION).forEach(([region, names]) =>
        names.forEach(name => seedRows.push({{ region, name }}))
      );
      if (seedRows.length) {{
        await fetch('/api/tiles/' + _TILE_ID + '/db/write', {{
          method: 'POST',
          headers: {{ 'Authorization': 'Basic ' + _DB_CRED, 'Content-Type': 'application/json' }},
          body: JSON.stringify({{ table: 'po_roster', mode: 'upsert', key_column: null, columns: ['region','name'], rows: seedRows }}),
        }}).catch(e => console.warn('po_roster seed:', e));
      }}
    }}
  }} catch(e) {{ _poRoster = JSON.parse(JSON.stringify(PO_BY_REGION)); console.warn('loadPoRoster:', e); }}
}}

function _rosterForRegion(region) {{
  return (_poRoster[region] || []).slice().sort((a,b) => a.localeCompare(b));
}}

function _rebuildPoDropdown(region, currentPo) {{
  const poSel = document.getElementById('assign-po');
  poSel.innerHTML = '';
  const blank = document.createElement('option');
  blank.value = ''; blank.textContent = '— Unassigned —';
  if (!currentPo || currentPo === 'Unassigned') blank.selected = true;
  poSel.appendChild(blank);
  const roster = _rosterForRegion(region);
  // include current PO even if not in roster (legacy)
  const all = [...new Set([...(currentPo && currentPo !== 'Unassigned' ? [currentPo] : []), ...roster])].sort((a,b) => a.localeCompare(b));
  all.forEach(po => {{
    const opt = document.createElement('option');
    opt.value = po; opt.textContent = po;
    if (po === currentPo) opt.selected = true;
    poSel.appendChild(opt);
  }});
}}

function _renderRosterList(region) {{
  const list = document.getElementById('roster-list');
  list.innerHTML = '';
  const assignedPos = new Set(RAW.map(r => r.po).filter(Boolean));
  _rosterForRegion(region).forEach(name => {{
    const inUse = assignedPos.has(name);
    const item = document.createElement('div');
    item.className = 'roster-item';
    item.innerHTML = `<span title="${{name}}">${{name}}</span>
      <button class="roster-del" onclick="_delPoFromRoster('${{region}}','${{name.replace(/'/g,"\\\\'")}}')"
        title="${{inUse ? 'Assigned to projects — cannot delete' : 'Remove from list'}}"
        ${{inUse ? 'disabled' : ''}}>🗑</button>`;
    list.appendChild(item);
  }});
  if (!list.children.length) list.innerHTML = '<div style="font-size:11px;color:var(--subtext);padding:2px 0">No POs in roster yet.</div>';
}}

async function _addPoToRoster() {{
  const input = document.getElementById('roster-add-input');
  const name  = (input.value || '').trim();
  if (!name) return;
  const r = RAW.find(x => x.pid === _assignPid);
  const region = r ? r.region : '';
  if (!region) return;
  if ((_poRoster[region] || []).includes(name)) {{ input.value = ''; return; }}
  (_poRoster[region] = _poRoster[region] || []).push(name);
  input.value = '';
  _renderRosterList(region);
  _rebuildPoDropdown(region, document.getElementById('assign-po').value);
  // persist
  if (_DB_CRED) {{
    await fetch('/api/tiles/' + _TILE_ID + '/db/write', {{
      method: 'POST',
      headers: {{ 'Authorization': 'Basic ' + _DB_CRED, 'Content-Type': 'application/json' }},
      body: JSON.stringify({{ table: 'po_roster', mode: 'upsert', key_column: null, columns: ['region','name'], rows: [{{ region, name }}] }}),
    }}).catch(e => console.warn('addPo:', e));
  }}
}}

async function _delPoFromRoster(region, name) {{
  const assignedPos = new Set(RAW.map(r => r.po).filter(Boolean));
  if (assignedPos.has(name)) return; // safety guard
  _poRoster[region] = (_poRoster[region] || []).filter(n => n !== name);
  const r = RAW.find(x => x.pid === _assignPid);
  _renderRosterList(region);
  _rebuildPoDropdown(region, document.getElementById('assign-po').value);
  // persist deletion
  if (_DB_CRED) {{
    await fetch('/api/tiles/' + _TILE_ID + '/db/write', {{
      method: 'POST',
      headers: {{ 'Authorization': 'Basic ' + _DB_CRED, 'Content-Type': 'application/json' }},
      body: JSON.stringify({{ table: 'po_roster', mode: 'delete', filter: {{ region, name }} }}),
    }}).catch(e => console.warn('delPo:', e));
  }}
}}

function _toggleRoster() {{
  const panel = document.getElementById('assign-roster-panel');
  const r = RAW.find(x => x.pid === _assignPid);
  if (!panel.classList.contains('open')) {{
    document.getElementById('roster-region-label').textContent = r ? r.region : '';
    _renderRosterList(r ? r.region : '');
    document.getElementById('roster-add-input').value = '';
  }}
  panel.classList.toggle('open');
}}

let _assignPid = null;
function openAssignModal(pid) {{
  const r = RAW.find(x => x.pid === pid);
  if (!r) return;
  _assignPid = pid;
  document.getElementById('assign-acct').textContent = r.acct || r.name || pid;
  document.getElementById('assign-tier').value = String(r.tier || 2);
  document.getElementById('assign-roster-panel').classList.remove('open');
  _rebuildPoDropdown(r.region, r.po);
  document.getElementById('assign-modal').classList.add('open');
}}
function closeAssignModal() {{
  document.getElementById('assign-modal').classList.remove('open');
  document.getElementById('assign-roster-panel').classList.remove('open');
  _assignPid = null;
}}
async function commitAssign() {{
  if (!_assignPid) return;
  const tier = parseInt(document.getElementById('assign-tier').value);
  const po   = document.getElementById('assign-po').value;
  const r = RAW.find(x => x.pid === _assignPid);
  const tierChanged = r && r.tier !== tier;
  const poChanged   = r && r.po   !== po;
  if (tierChanged || poChanged) {{
    await saveAssignment(_assignPid, tier, po);
  }} else {{
    applyFilters();
  }}
  closeAssignModal();
  if (saves.length === 0) applyFilters();
}}
document.addEventListener('keydown', e => {{ if (e.key === 'Escape') closeAssignModal(); }});

const STATUS_ORDER = {{Red:0, Yellow:1, Watermelon:2, Green:3, 'No Pulse':4, 'On Hold':5}};
const TIER_LABEL   = {{1:'T1',2:'T2',3:'T3'}};

const RULE_KEY = {{
  'MARGIN_RED':        'Close Margin critically below Bid Margin (RED threshold)',
  'MARGIN_YELLOW':     'Close Margin moderately below Bid Margin (YELLOW threshold)',
  'FAR_RED_NEG':       'FAR is negative — project forecasted to overrun budget',
  'FAR_RED_UNDERUTIL': 'FAR >50% remaining with <30% schedule left — severe underutilisation',
  'FAR_YELLOW':        'FAR utilisation warning — moderate schedule/spend mismatch',
  'RR_RISK':           'Open Resource Requests with pending revenue — staffing gap risk',
  'GDC_LOW':           'GDC India assigned share ≤ 65% — more than ⅓ of the team is non-GDC',
  'OVERDUE_INV':       'Overdue invoices outstanding — cash collection risk',
  'SWE_BURNING_HOT':   'SWE burn rate alert — hours/spend tracking off-plan',
  'NO_PULSE':          'No pulse submitted — project status unknown (governance violation)',
  'NO_STEERCO':     'Next Steering Committee Date required for projects ≥ $750K. ACTION: If exempt (SEH, Advisory, etc.) set date to 01/01/2100.',
  'MISSING_PTG':       'Missing project timeline/go-live date or steerco date',
  'END_DATE_PAST':  'End date has passed — project may need extension or closure',
  'END_DATE_30':    'End date within 30 days — urgent renewal/extension needed',
  'END_DATE_60':    'End date within 60 days — renewal/extension decision needed',
  'END_DATE_90':    'End date within 90 days — begin renewal planning',
  'END_DATE_120':   'End date within 120 days — flag for upcoming renewal',
  'PIPE_CLOSE_THIS_M':  'Pipeline opportunity closing this calendar month — action needed',
  'PIPE_CLOSE_THIS_Q':  'Pipeline opportunity closing this SF fiscal quarter (Q ends Jan/Apr/Jul/Oct)',
}};

const RULE_GROUP = {{
  'MARGIN_RED':        'Margin',
  'MARGIN_YELLOW':     'Margin',
  'FAR_RED_NEG':       'FAR',
  'FAR_RED_UNDERUTIL': 'FAR',
  'FAR_YELLOW':        'FAR',
  'RR_RISK':           'Resource',
  'GDC_LOW':           'Resource',
  'OVERDUE_INV':       'Invoice',
  'SWE_BURNING_HOT':   'SWE Burn',
  'NO_PULSE':          'Governance',
  'NO_STEERCO':        'Governance',
  'MISSING_PTG':       'Governance',
  'END_DATE_PAST':  'End Date',
  'END_DATE_30':    'End Date',
  'END_DATE_60':    'End Date',
  'END_DATE_90':    'End Date',
  'END_DATE_120':   'End Date',
  'PIPE_CLOSE_THIS_M':  'Pipeline',
  'PIPE_CLOSE_THIS_Q':  'Pipeline',
}};

let sortCol = 'status', sortDir = 1;
let filterHW = false, filterSWE = false;
let filterStatuses = new Set(), filterTiers = new Set(), filterPOs = new Set(), filterRules = new Set();
let filterTypes = new Set(), filterBilling = new Set();
let groupBy = '';       // legacy — kept for backward compat, not used in new render
let groupKeys = [];     // ordered array of active grouping dimensions
const _REGION_KEY = 'acc_region_817';
let filterRegion = null;
let _gmCurrentData = [];
const TIER_NAME = {{1:'Tier 1', 2:'Tier 2', 3:'Tier 3'}};


function fmtMoney(v) {{
  if (!v && v !== 0) return '—';
  const abs = Math.abs(v);
  const sign = v < 0 ? '-' : '';
  if (abs >= 1e6) return sign + '$' + (abs/1e6).toFixed(2) + 'M';
  return sign + '$' + abs.toLocaleString('en-US', {{maximumFractionDigits:0}});
}}

function fmtPct(v) {{
  if (v === null || v === undefined) return '—';
  return v.toFixed(1) + '%';
}}

function bidMarginColor(v) {{
  if (v === null || v === undefined) return '';
  if (v > 13) return 'color:var(--green)';
  if (v >= 5)  return 'color:var(--yellow)';
  return 'color:var(--red)';
}}

function statusBadge(r) {{
  const s = r.status;
  const cls = s.replace(' ','-');
  const icons = {{Red:'🔴',Yellow:'🟡',Green:'🟢',Watermelon:'🍉','No Pulse':'⚫','On Hold':'⏸'}};
  const trendArrow = {{
    'Up':'↑', 'Improving':'↑', 'Improving Slightly':'↗',
    'Stable':'→',
    'Down':'↓', 'Declining':'↓', 'Declining Slightly':'↘', 'Worsening':'↓',
  }};
  const trendColor = {{
    'Up':'var(--green)', 'Improving':'var(--green)', 'Improving Slightly':'var(--green)',
    'Stable':'var(--subtext)',
    'Down':'var(--red)', 'Declining':'var(--red)', 'Declining Slightly':'var(--yellow)', 'Worsening':'var(--red)',
  }};
  const arrow = r.pulse_trend ? (trendArrow[r.pulse_trend] || '') : '';
  const arrowCol = r.pulse_trend ? (trendColor[r.pulse_trend] || 'var(--subtext)') : '';
  const trendHtml = arrow ? ` <span style="font-size:13px;font-weight:700;color:${{arrowCol}};vertical-align:middle" title="Trend: ${{r.pulse_trend}}">${{arrow}}</span>` : '';
  let html = `<span class="badge badge-${{cls}}">${{icons[s]||''}} ${{s}}</span>${{trendHtml}}`;
  function scoreCls(v) {{ return v >= 70 ? 'score-green' : v >= 30 ? 'score-yellow' : 'score-red'; }}
  const scores = [];
  if (r.health_risk  != null) scores.push(`<a class="score-chip ${{scoreCls(r.health_risk)}}" href="https://org62.lightning.force.com/lightning/n/ProjectScoringInformation?c__projectId=${{r.pid}}" target="_blank" title="Click to open Project Scoring in Salesforce">H&amp;R&nbsp;${{r.health_risk.toFixed(0)}}</a>`);
  if (r.data_quality != null) scores.push(`<a class="score-chip ${{scoreCls(r.data_quality)}}" href="https://org62.lightning.force.com/lightning/n/ProjectScoringInformation?c__projectId=${{r.pid}}" target="_blank" title="Click to open Project Scoring in Salesforce">DQ&nbsp;${{r.data_quality.toFixed(0)}}</a>`);
  if (r.csat_score   != null) scores.push(`<span class="score-chip ${{scoreCls(r.csat_score * 20)}}">CSAT&nbsp;${{r.csat_score.toFixed(1)}}</span>`);
  if (scores.length) html += `<div style="margin-top:4px">${{scores.join(' ')}}</div>`;
  html += pulseIcon(r);
  return html;
}}

function dimCls(v) {{
  if (!v) return 'pdim-grey';
  const lv = v.toLowerCase();
  if (lv === 'green') return 'pdim-green';
  if (lv === 'yellow') return 'pdim-yellow';
  if (lv === 'red') return 'pdim-red';
  return 'pdim-grey';
}}

// ── Shared popup behaviour: hover shows, click pins, X/Escape force-closes ─────
// hoverClose=true: modal closes on mouseleave (pulse, FAR); false: only X/Escape/backdrop (pipeline)
function _makePopup(modal, hoverClose) {{
  const _resetPos = () => {{
    modal.style.transform = '';
    modal.style.left = '';
    modal.style.top  = '';
  }};
  const obj = {{
    timer: null,
    pinned: false,
    show(openFn) {{
      clearTimeout(obj.timer);
      if (!modal.classList.contains('open')) _resetPos();
      openFn();
    }},
    pin(openFn) {{
      clearTimeout(obj.timer);
      if (!modal.classList.contains('open')) _resetPos();
      openFn();
      obj.pinned = true;
    }},
    hide(force) {{
      if (obj.pinned && !force) return;
      obj.timer = setTimeout(() => {{
        modal.classList.remove('open');
        obj.pinned = false;
        _resetPos();
      }}, 120);
    }},
  }};
  if (hoverClose) {{
    modal.addEventListener('mouseenter', () => clearTimeout(obj.timer));
    modal.addEventListener('mouseleave', () => obj.hide(false));
  }}
  return obj;
}}

// ── Draggable popup windows ─────────────────────────────────────────────────────
// Call once per popup; the header element becomes the drag handle.
// On first drag the centered transform is replaced with explicit px so the
// window can be positioned freely — just like a standard OS window.
function _makeDraggable(modal, handle) {{
  handle.addEventListener('mousedown', e => {{
    if (e.target.tagName === 'BUTTON' || e.target.tagName === 'A') return;
    e.preventDefault();
    // Resolve current visual position to explicit px
    const rect = modal.getBoundingClientRect();
    modal.style.transform = 'none';
    modal.style.left = rect.left + 'px';
    modal.style.top  = rect.top  + 'px';
    const startX = e.clientX - rect.left;
    const startY = e.clientY - rect.top;
    const onMove = ev => {{
      let nx = ev.clientX - startX;
      let ny = ev.clientY - startY;
      // Keep within viewport
      nx = Math.max(0, Math.min(nx, window.innerWidth  - modal.offsetWidth));
      ny = Math.max(0, Math.min(ny, window.innerHeight - modal.offsetHeight));
      modal.style.left = nx + 'px';
      modal.style.top  = ny + 'px';
    }};
    const onUp = () => {{
      document.removeEventListener('mousemove', onMove);
      document.removeEventListener('mouseup',   onUp);
    }};
    document.addEventListener('mousemove', onMove);
    document.addEventListener('mouseup',   onUp);
  }});
}}

// ── Pulse popup ────────────────────────────────────────────────────────────────
const _ptt      = document.getElementById('pulse-tooltip');
const _pttTitle = document.getElementById('pulse-tooltip-title');
const _pttBody  = document.getElementById('pulse-tooltip-body');
const _pttLink  = document.getElementById('pulse-tooltip-link');
const _pulsePopup = _makePopup(_ptt, true);
_makeDraggable(_ptt, document.getElementById('pulse-tooltip-header'));

function _showPulse(el) {{
  _pulsePopup.show(() => {{
    try {{
      const d = JSON.parse(el.dataset.pmodal.replace(/&quot;/g, '"'));
      _pttTitle.textContent = d.title || '';
      _pttBody.innerHTML    = d.body  || '';
      if (d.url) {{ _pttLink.href = d.url; _pttLink.style.display = 'inline'; }}
      else       {{ _pttLink.style.display = 'none'; }}
      _ptt.classList.add('open');
    }} catch(e) {{}}
  }});
}}
function _pinPulse(el) {{ _pulsePopup.pin(() => _showPulse(el)); }}
function _hidePulse(force) {{ _pulsePopup.hide(force); }}
function _movePulse() {{}}

// ── Summary (Pulse detail) popup ───────────────────────────────────────────────
const _sumMod   = document.getElementById('summary-modal');
const _sumTitle = document.getElementById('summary-modal-title');
const _sumBody  = document.getElementById('summary-modal-body');
const _sumPopup = _makePopup(_sumMod, true);
_makeDraggable(_sumMod, document.getElementById('summary-modal-header'));

function _buildSummaryHtml(r) {{
  const RAG_DOT = {{ Red:'🔴', Yellow:'🟡', Green:'🟢', Unknown:'⚫' }};
  const chip = (label, val) => {{
    const cls = ['Red','Yellow','Green'].includes(val) ? val : 'Unknown';
    return `<span class="sum-rag-chip ${{cls}}">${{RAG_DOT[cls]||'⚫'}} ${{label}}</span>`;
  }};
  const block = (label, text) => text
    ? `<div class="sum-section"><div class="sum-section-title">${{label}}</div><div class="sum-block">${{text.replace(/</g,'&lt;')}}</div></div>`
    : '';
  // RAG chips with inline P2G text when present
  const dims = [
    ['Scope',    r.pulse_scope,    r.ptg_scope],
    ['Schedule', r.pulse_sched,    r.ptg_sched],
    ['Budget',   r.pulse_budget,   r.ptg_budget],
    ['Resource', r.pulse_resource, r.ptg_resource],
    ['Customer', r.pulse_customer, r.ptg_customer],
  ];
  const ragRows = dims.map(([label, rag, ptg]) => {{
    const cls = ['Red','Yellow','Green'].includes(rag) ? rag : 'Unknown';
    const chipHtml = `<span class="sum-rag-chip ${{cls}}">${{RAG_DOT[cls]||'⚫'}} ${{label}}</span>`;
    const ptgHtml = ptg
      ? `<div style="font-size:11px;color:var(--subtext);margin:2px 0 8px 4px;padding-left:8px;border-left:2px solid var(--border)">${{ptg.replace(/</g,'&lt;')}}</div>`
      : '';
    return chipHtml + ptgHtml;
  }});
  const ragSection = `<div class="sum-section">
    <div class="sum-section-title">RAG Status &amp; Path to Green</div>
    <div>${{ragRows.join('')}}</div>
  </div>`;
  // Body: RAG → Summary → Leadership Notes → Action Needed
  const body = ragSection
    + block('Summary', r.summary)
    + block('Leadership Notes', r.leadership)
    + block('Action Needed from Leadership', r.pulse_action);
  // Footer: meta fields + pulse link (always visible, never scrolls away)
  const pulseLink = r.pulse_id
    ? `<a href="https://org62.lightning.force.com/lightning/r/Project_Health_Check__c/${{r.pulse_id}}/view"
        target="_blank" rel="noopener">↗ Open Pulse in Salesforce</a>`
    : '';
  const footer = [
    r.pulse_updated ? `<span>Last updated: <strong>${{r.pulse_updated}}</strong></span>` : '',
    r.pulse_trend   ? `<span>Trend: <strong>${{r.pulse_trend}}</strong></span>` : '',
    r.pulse_steerco ? `<span>Next SteerCo: <strong>${{r.pulse_steerco}}</strong></span>` : '',
    r.pulse_golive  ? `<span>Next Go-Live: <strong>${{r.pulse_golive}}</strong></span>` : '',
    r.high_watch    ? `<span style="color:var(--red);font-weight:700">⚑ High Watch</span>` : '',
    pulseLink,
  ].filter(Boolean).join('');
  return {{ body, footer }};
}}

function _showSummary(el) {{
  _sumPopup.show(() => {{
    _sumTitle.textContent = (el.dataset.proj || 'Pulse Summary') + ' — Pulse Detail';
    const r = JSON.parse(el.dataset.sum.replace(/&quot;/g,'"'));
    const {{ body, footer }} = _buildSummaryHtml(r);
    _sumBody.innerHTML = body;
    document.getElementById('summary-modal-footer').innerHTML = footer;
    _sumMod.classList.add('open');
  }});
}}
function _pinSummary(el) {{ _sumPopup.pin(() => _showSummary(el)); }}
function _hideSummary(force) {{ _sumPopup.hide(force === true); }}

// ── FAR popup ──────────────────────────────────────────────────────────────────
const _ftt      = document.getElementById('far-tooltip');
const _fttTitle = document.getElementById('far-tooltip-title');
const _fttBody  = document.getElementById('far-tooltip-body');
const _farPopup = _makePopup(_ftt, true);
_makeDraggable(_ftt, document.getElementById('far-tooltip-header'));

function _buildFarHtml(el) {{
  try {{
    const d = JSON.parse(el.dataset.far.replace(/&quot;/g,'"'));
    let html = '';
    if (d.reason) html += `<div><strong style="color:#888;font-size:10px">FAR REASON</strong><div style="margin-top:3px">${{d.reason}}</div></div>`;
    return html || '—';
  }} catch(e) {{ return el.dataset.far; }}
}}
function _showFar(el) {{
  _farPopup.show(() => {{
    _fttBody.innerHTML = _buildFarHtml(el);
    _ftt.classList.add('open');
  }});
}}
function _pinFar(el) {{
  _farPopup.pin(() => {{
    _fttBody.innerHTML = _buildFarHtml(el);
    _ftt.classList.add('open');
  }});
}}
function _hideFar(force) {{ _farPopup.hide(force === true); }}
function _moveFar() {{}}

// ── Path-to-Green popup ─────────────────────────────────────────────────────────
const _p2gMod   = document.getElementById('p2g-modal');
const _p2gTitle = document.getElementById('p2g-modal-title');
const _p2gBody  = document.getElementById('p2g-modal-body');
const _p2gPopup = _makePopup(_p2gMod, true);
_makeDraggable(_p2gMod, document.getElementById('p2g-modal-header'));

function _buildP2gHtml(el) {{
  try {{
    const d = JSON.parse(el.dataset.p2g.replace(/&quot;/g,'"'));
    const sec = (label, val) => val
      ? `<div style="margin-top:8px;padding-top:6px;border-top:1px solid #eee">
           <div style="font-size:10px;font-weight:700;text-transform:uppercase;letter-spacing:.4px;color:#888;margin-bottom:3px">${{label}}</div>
           <div>${{val}}</div></div>`
      : '';
    let html = '';
    html += sec('Budget',   d.ptg_budget);
    html += sec('Scope',    d.ptg_scope);
    html += sec('Schedule', d.ptg_sched);
    html += sec('Resource', d.ptg_resource);
    html += sec('Customer', d.ptg_customer);
    return html || '<div style="color:var(--subtext);font-style:italic">No Path to Green entries recorded.</div>';
  }} catch(e) {{ return '—'; }}
}}
function _showP2g(el) {{
  _p2gPopup.show(() => {{
    _p2gTitle.textContent = (el.dataset.proj || '') + ' — Path to Green';
    _p2gBody.innerHTML = _buildP2gHtml(el);
    _p2gMod.classList.add('open');
  }});
}}
function _pinP2g(el) {{
  _p2gPopup.pin(() => {{
    _p2gTitle.textContent = (el.dataset.proj || '') + ' — Path to Green';
    _p2gBody.innerHTML = _buildP2gHtml(el);
    _p2gMod.classList.add('open');
  }});
}}
function _hideP2g(force) {{ _p2gPopup.hide(force === true); }}

// ── Pipeline popup (data from /_data/ pipeline JSON, refreshed each hourly run) ──
const _pipeModal  = document.getElementById('pipe-modal');
const _pipeTitle  = document.getElementById('pipe-title');
const _pipeBody   = document.getElementById('pipe-body');
let _pipeHoverTimer, _pipeCache = null;
let _pipeSortCol = 'close', _pipeSortDir = 1;

const _pipePopup = _makePopup(_pipeModal, true);
_makeDraggable(_pipeModal, document.getElementById('pipe-header'));

function _showPipe(el) {{
  clearTimeout(_pipeHoverTimer);
  _pipeHoverTimer = setTimeout(() => _openPipe(el, false), 250);
}}
function _pinPipe(el) {{
  clearTimeout(_pipeHoverTimer);
  _openPipe(el, true);
}}
function _hidePipeHover(force) {{
  clearTimeout(_pipeHoverTimer);
  _pipePopup.hide(force === true);
}}
function _renderPipeTable(allRecs) {{
  const col = _pipeSortCol, dir = _pipeSortDir;
  const recs = allRecs.slice().sort((a, b) => {{
    if (col === 'amount') return dir * ((b.converted || b.amount || 0) - (a.converted || a.amount || 0));
    if (col === 'close')  return dir * (a.close || '').localeCompare(b.close || '');
    return 0;
  }});
  const fmt = v => !v ? '—' : (Math.abs(v) >= 1e6 ? '$' + (v/1e6).toFixed(2) + 'M' : '$' + Number(v).toLocaleString());
  const fmtPct = v => v == null ? '—' : v.toFixed(1) + '%';
  const total = allRecs.reduce((s, r) => s + (r.converted || r.amount || 0), 0);
  const fcColor = f => ({{'Commit':'color:#1a6b3c;font-weight:700','Best Case':'color:#7d5a00;font-weight:700','Pipeline':'color:var(--subtext)'}})[f] || 'color:var(--subtext)';
  const thStyle = (c, align) => {{
    const active = c === col;
    const arrow = active ? (dir === -1 ? ' ▼' : ' ▲') : ' ⇅';
    return `style="padding:7px 10px;color:#fff;font-weight:600;white-space:nowrap;cursor:pointer;user-select:none${{align ? ';text-align:' + align : ''}}" onclick="_pipeSortBy('${{c}}')"`;
  }};
  const rows = recs.map((r,i) => {{
    const sfUrl = `https://org62.lightning.force.com/lightning/r/Opportunity/${{r.id}}/view`;
    const bg = i % 2 === 0 ? '' : 'background:#f8f9fa';
    const pushBadge = r.push_count > 0 ? ` <span style="font-size:9px;color:var(--red);font-weight:700">↩${{r.push_count}}</span>` : '';
    const mgrFc = r.mgr_forecast && r.mgr_forecast !== r.forecast ? `<div style="font-size:9px;color:var(--subtext)">Mgr: ${{r.mgr_forecast}}</div>` : '';
    const margin = r.bid_margin != null ? `<div style="font-size:9px;color:var(--subtext)">${{fmtPct(r.bid_margin)}} margin</div>` : '';
    const nextSteps = r.description ? `<div style="font-size:10px;color:var(--subtext);margin-top:2px;white-space:normal">${{r.description.slice(0,200)}}${{r.description.length>200?'…':''}}</div>` : '—';
    return `<tr style="${{bg}}">
      <td style="padding:5px 10px;white-space:nowrap"><a href="${{sfUrl}}" target="_blank" style="color:var(--blue);text-decoration:none">${{r.name || '—'}}</a></td>
      <td style="padding:5px 10px;white-space:nowrap">${{r.stage || '—'}}</td>
      <td style="padding:5px 10px;text-align:right;white-space:nowrap;color:var(--lblue);font-weight:600">${{fmt(r.converted || r.amount)}}${{r.currency && r.currency !== 'USD' ? `<div style="font-size:9px;color:var(--subtext)">${{r.currency}} ${{fmt(r.amount)}}</div>` : ''}}</td>
      <td style="padding:5px 10px;white-space:nowrap">${{r.close || '—'}}${{pushBadge}}</td>
      <td style="padding:5px 10px;white-space:nowrap"><span style="${{fcColor(r.forecast)}}">${{r.forecast || '—'}}</span>${{mgrFc}}</td>
      <td style="padding:5px 10px;white-space:nowrap">${{r.owner || '—'}}${{margin}}</td>
      <td style="padding:5px 10px;max-width:280px">${{nextSteps}}</td>
    </tr>`;
  }}).join('');
  const amtArrow = _pipeSortCol==='amount' ? (_pipeSortDir===-1?' ▼':' ▲') : ' ⇅';
  const clsArrow = _pipeSortCol==='close'  ? (_pipeSortDir===-1?' ▼':' ▲') : ' ⇅';
  _pipeBody.innerHTML = `
    <div style="overflow-x:auto">
    <table style="width:100%;min-width:900px;border-collapse:collapse;font-size:11.5px;table-layout:auto">
      <thead>
        <tr style="text-align:left;background:var(--navy)">
          <th style="padding:7px 10px;color:#fff;font-weight:600;white-space:nowrap">Opportunity</th>
          <th style="padding:7px 10px;color:#fff;font-weight:600;white-space:nowrap">Stage</th>
          <th style="padding:7px 10px;color:#fff;font-weight:600;text-align:right;white-space:nowrap;cursor:pointer;user-select:none" onclick="_pipeSortBy('amount')">Amount (USD)${{amtArrow}}</th>
          <th style="padding:7px 10px;color:#fff;font-weight:600;white-space:nowrap;cursor:pointer;user-select:none" onclick="_pipeSortBy('close')">Close Date${{clsArrow}}</th>
          <th style="padding:7px 10px;color:#fff;font-weight:600;white-space:nowrap">Forecast</th>
          <th style="padding:7px 10px;color:#fff;font-weight:600;white-space:nowrap">Owner</th>
          <th style="padding:7px 10px;color:#fff;font-weight:600;white-space:nowrap">Description</th>
        </tr>
      </thead>
      <tbody>${{rows}}</tbody>
      <tfoot>
        <tr style="border-top:2px solid var(--border);font-weight:700;background:#f1f3f9">
          <td colspan="2" style="padding:6px 10px">Total (${{allRecs.length}} opportunities)</td>
          <td style="padding:6px 10px;text-align:right;color:var(--lblue)">${{fmt(total)}}</td>
          <td colspan="4"></td>
        </tr>
      </tfoot>
    </table>
    </div>`;
}}
let _pipeCurrentRecs = null;
function _pipeSortBy(col) {{
  if (_pipeSortCol === col) {{ _pipeSortDir *= -1; }} else {{ _pipeSortCol = col; _pipeSortDir = col === 'amount' ? -1 : 1; /* close & others: asc */ }}
  if (_pipeCurrentRecs) _renderPipeTable(_pipeCurrentRecs);
}}
async function _openPipe(el, pin) {{
  const acctId   = el.dataset.acctId;
  const acctName = el.dataset.acct || acctId;
  _pipeTitle.textContent = acctName + ' — Open Pipeline';
  _pipeBody.innerHTML = '<div style="color:var(--subtext);padding:20px 0;text-align:center">Loading…</div>';
  _pipeModal.classList.add('open');
  if (pin) _pipePopup.pinned = true;
  try {{
    if (_pipeCache === null) {{
      _pipeCache = _INLINE.pipe || {{}};
    }}
    const recs = (_pipeCache[acctId] || []);
    if (!recs.length) {{
      _pipeBody.innerHTML = '<div style="color:var(--subtext);padding:20px 0;text-align:center">No open pipeline opportunities found for this account.</div>';
      return;
    }}
    _pipeCurrentRecs = recs;
    _pipeSortCol = 'close'; _pipeSortDir = 1;
    _renderPipeTable(recs);
  }} catch(err) {{
    _pipeBody.innerHTML = `<div style="color:var(--red);padding:20px 0;text-align:center">Failed to load pipeline: ${{err.message}}</div>`;
  }}
}}
function _hidePipe(force) {{
  _pipePopup.hide(force === true || force === undefined);
}}

function pulseIcon(r) {{
  const pulseUrl = r.pulse_id
    ? `https://org62.lightning.force.com/lightning/r/Project_Health_Check__c/${{r.pulse_id}}/view`
    : '';
  if (!r.has_pulse) {{
    const body = [
      `<div>This project has no active pulse record in Salesforce.</div>`,
      (r.bookings||0) >= 150000 ? `<div style="margin-top:8px;color:#c0392b;font-weight:700">⚠ Governance violation — Rule 2A (Bookings ≥ $150K)</div>` : '',
      r.slack_intel ? `<div style="margin-top:10px;padding-top:8px;border-top:1px solid #eee"><span style="font-size:10px;font-weight:700;color:#5b5e6d">💬 SLACK INTEL</span><div style="margin-top:4px">${{r.slack_intel}}</div></div>` : '',
    ].join('');
    const d = JSON.stringify({{title:'No Pulse on File', body, url:''}}).replace(/"/g,'&quot;');
    return `<span class="pulse-icon" data-pmodal="${{d}}" onmouseenter="_showPulse(this)" onmouseleave="_hidePulse(false)" onclick="_pinPulse(this)">⚫</span>`;
  }}
  const dims = [
    ['Scope',    r.pulse_scope],
    ['Schedule', r.pulse_sched],
    ['Budget',   r.pulse_budget],
    ['Resource', r.pulse_resource],
    ['Customer', r.pulse_customer],
  ];
  const dimHtml = `<div class="pulse-dim">${{dims.map(([lbl,val]) =>
    `<span class="pdim ${{dimCls(val)}}">${{lbl}}: ${{val||'—'}}</span>`).join('')}}</div>`;
  const body = [
    r.pulse_updated ? `<div style="color:#888;font-size:11px;margin-bottom:6px">Updated: ${{r.pulse_updated}}</div>` : '',
    dimHtml,
    r.pulse_trend   ? `<div style="margin-top:4px"><b>Trend:</b> ${{r.pulse_trend}}</div>` : '',
    r.pulse_steerco ? `<div><b>SteerCo:</b> ${{r.pulse_steerco}}</div>` : '',
    r.pulse_golive  ? `<div><b>Next Go-Live:</b> ${{r.pulse_golive}}</div>` : '',
    r.summary       ? `<div style="margin-top:10px;padding-top:8px;border-top:1px solid #eee"><div style="font-size:10px;font-weight:700;text-transform:uppercase;letter-spacing:.4px;color:#888;margin-bottom:4px">Project Progress Summary</div><div>${{r.summary}}</div></div>` : '',
    r.pulse_action  ? `<div style="margin-top:10px;padding-top:8px;border-top:1px solid #eee"><div style="font-size:10px;font-weight:700;text-transform:uppercase;letter-spacing:.4px;color:#c0392b;margin-bottom:4px">Action Needed from Leadership</div><div style="color:#c0392b">${{r.pulse_action}}</div></div>` : '',
    r.leadership    ? `<div style="margin-top:10px;padding-top:8px;border-top:1px solid #eee"><div style="font-size:10px;font-weight:700;text-transform:uppercase;letter-spacing:.4px;color:#888;margin-bottom:4px">Leadership Notes</div><div>${{r.leadership}}</div></div>` : '',
    r.slack_intel   ? `<div style="margin-top:10px;padding-top:8px;border-top:1px solid #eee"><div style="font-size:10px;font-weight:700;color:#5b5e6d;margin-bottom:4px">💬 SLACK INTEL</div><div>${{r.slack_intel}}</div></div>` : '',
  ].join('');
  const d = JSON.stringify({{title: r.name || 'Pulse', body, url: pulseUrl}}).replace(/"/g,'&quot;');
  return `<span class="pulse-icon" data-pmodal="${{d}}" onmouseenter="_showPulse(this)" onmouseleave="_hidePulse(false)" onclick="_pinPulse(this)">📋</span>`;
}}


const ROLE_LABELS = {{
  'PM':  'Project Manager',
  'PM2': 'Program Manager',
  'AP':  'ProServ Account Partner',
  'AO':  'License Account Executive',
  'PO':  'Portfolio Owner',
  'ES':  'Executive Sponsor',
}};
function teamHtml(team) {{
  if (!team) return '';
  return '<div class="team-list">' +
    team.split(' · ').map(t => {{
      const colon = t.indexOf(':');
      const role = t.slice(0,colon);
      const name = t.slice(colon+1).trim();
      const tip  = ROLE_LABELS[role] || role;
      return `<span><b class="role-abbr">${{role}}:<span class="role-tip">${{tip}}</span></b> ${{name}}</span>`;
    }}).join('') + '</div>';
}}

function resourcingHtml(r) {{
  if (r.gdc_total === 0) {{
    return `<div style="font-size:11px;color:var(--subtext)">No assigned<br>resources</div>`;
  }}
  const isLow    = r.gdc_pct !== null && r.gdc_pct <= 65;
  const pctColor = isLow ? 'var(--red)' : 'var(--green)';
  const pctBg    = isLow ? '#FDEDEC'    : '#EAFAF1';
  const pct      = r.gdc_pct !== null ? r.gdc_pct + '%' : '?';

  const payload = JSON.stringify({{
    name: (r.name || r.acct || r.pid || ''),
    resources: r.gdc_resources || []
  }}).replace(/"/g, '&quot;');

  return `<div class="res-hover-wrap"
      data-res="${{payload}}"
      onmouseenter="_showRes(this)"
      onmouseleave="_hideRes(false)"
      onclick="_pinRes(this)">
    <div style="font-size:10px;color:var(--subtext);margin-bottom:2px">GDC India</div>
    <div style="display:inline-block;padding:2px 7px;border-radius:5px;background:${{pctBg}};color:${{pctColor}};font-size:13px;font-weight:700">${{pct}}</div>
    <div style="font-size:10px;color:var(--subtext);margin-top:2px">${{r.gdc_india}}/${{r.gdc_total}} assigned</div>
  </div>`;
}}

function finHtml(r) {{
  const neg = v => v && String(v).startsWith('-') ? ' neg' : (v && String(v).startsWith('+') ? ' pos' : '');
  const typeLabel = r.billing_type || '—';

  // FAR with hover popup for reason codes only
  const farTipData = JSON.stringify({{
    reason: [r.far_reason, r.far_subreason].filter(Boolean).join(' · '),
  }}).replace(/"/g,'&quot;');
  const hasFarTip = r.far_reason || r.far_subreason;
  const farCell = hasFarTip
    ? `<span class="far-tip-wrap" data-far="${{farTipData}}"
         onmouseenter="_showFar(this)" onmouseleave="_hideFar(false)" onclick="_pinFar(this)"
         style="cursor:pointer;border-bottom:1px dotted var(--subtext)">${{r.fmt_far||'—'}} ⓘ</span>`
    : (r.fmt_far || '—');


  // Overdue invoice — always shown; red when > 0
  const overdueStyle = r.overdue_inv > 0 ? ' style="color:var(--red)"' : '';
  const overdueVal   = r.fmt_overdue || '—';

  const fi = (k, v, cls='', sty='') =>
    `<div class="fin-item"><span class="fk">${{k}}:</span><span class="fv${{cls}}"${{sty ? ` style="${{sty}}"` : ''}}>${{v}}</span></div>`;

  // Left sub-column: revenue metrics
  const leftCol = `<div class="fin-col">
    ${{fi('Type', typeLabel)}}
    ${{fi('Bookings', r.fmt_bookings||'—')}}
    ${{fi('Billings', r.fmt_billings||'—')}}
    ${{fi('FAR', farCell, r.far < 0 ? ' neg' : '')}}
    ${{fi('Overdue Inv', overdueVal, '', r.overdue_inv > 0 ? 'color:var(--red)' : '')}}
  </div>`;

  // Right sub-column: margin metrics
  const rightCol = `<div class="fin-col">
    ${{fi('Bid Margin', r.fmt_bid||'—', '', bidMarginColor(r.bid_margin_raw))}}
    ${{fi('Margin@Close', r.fmt_close||'—', '', r.close_margin_raw !== null && r.close_margin_raw < 0 ? 'color:var(--red)' : '')}}
    ${{fi('Del - Bid', r.fmt_delivered||'—', neg(r.fmt_delivered))}}
    ${{r.fmt_eva_amt ? fi('EvA $', r.fmt_eva_amt, neg(r.fmt_eva_amt)) : ''}}
    ${{r.fmt_eva_pct ? fi('EvA %', r.fmt_eva_pct, neg(r.fmt_eva_pct)) : ''}}
  </div>`;

  // Optional extras below both columns
  const extraItems = [
    r.fmt_unsch   ? fi('Unsch Backlog', r.fmt_unsch) : '',
    r.fmt_actuals ? fi('Actuals Rem',   r.fmt_actuals) : '',
    r.fmt_pipe && r.acct_id
      ? `<div class="fin-item"><span class="fk">Open Pipe:</span><span class="fv" style="color:var(--lblue)"><span class="pipe-trigger" data-acct-id="${{r.acct_id}}" data-acct="${{(r.acct||'').replace(/"/g,'&quot;')}}" data-fmt="${{r.fmt_pipe}}" onmouseenter="_showPipe(this)" onmouseleave="_hidePipeHover(false)" onclick="_pinPipe(this)" style="cursor:pointer;border-bottom:1px dotted var(--lblue)">${{r.fmt_pipe}} ⓘ</span></span></div>`
      : (r.fmt_pipe ? fi('Open Pipe', r.fmt_pipe, '', 'color:var(--lblue)') : ''),
  ].filter(Boolean).join('');

  const extrasHtml = extraItems ? `<div class="fin-extras">${{extraItems}}</div>` : '';

  return `<div><div class="fin-2col">${{leftCol}}${{rightCol}}</div>${{extrasHtml}}</div>`;
}}

function rulesHtml(rules) {{
  if (!rules) return '';
  return rules.split(', ').filter(Boolean).map(c => {{
    const isNeg = c.includes('NEG') || c.includes('RED') || c === 'END_DATE_PAST' || c === 'END_DATE_30' || c === 'PIPE_CLOSE_THIS_M';
    const isWarn = c.includes('YELLOW') || c === 'END_DATE_60' || c === 'END_DATE_90' || c === 'END_DATE_120' || c === 'PIPE_CLOSE_THIS_Q';
    const tip = RULE_KEY[c] || c;
    const grp = RULE_GROUP[c] ? `<span style="opacity:.65;font-size:9px;margin-right:2px">${{RULE_GROUP[c]}}:</span>` : '';
    const cls = isNeg ? ' neg' : isWarn ? ' warn' : '';
    return `<span class="rule-code${{cls}}">${{grp}}${{c}}<span class="rule-tip">${{tip}}</span></span>`;
  }}).join('');
}}

function baselinesHtml(bl, r) {{
  if (!bl) return '';
  const dots = '<div class="bl-list">' +
    bl.split(' | ').map(part => {{
      const eq = part.indexOf('=');
      const lbl = part.slice(0,eq);
      const val = part.slice(eq+1) || 'None';
      const dot = `<span class="bl-dot bl-dot-${{val}}"></span>`;
      return `<span><span class="bl-lbl">${{lbl}}</span>${{dot}}</span>`;
    }}).join('<br>') + '</div>';
  if (!r) return dots;
  const hasP2g = r.ptg_budget || r.ptg_scope || r.ptg_sched || r.ptg_resource || r.ptg_customer;
  if (!hasP2g) return dots;
  const p2gData = JSON.stringify({{
    ptg_budget: r.ptg_budget || '',
    ptg_scope:  r.ptg_scope  || '',
    ptg_sched:  r.ptg_sched  || '',
    ptg_resource: r.ptg_resource || '',
    ptg_customer: r.ptg_customer || '',
  }}).replace(/"/g,'&quot;');
  const p2gLink = `<div style="margin-top:4px"><span data-p2g="${{p2gData}}" data-proj="${{(r.name||r.acct||'').replace(/"/g,'&quot;')}}"
    onmouseenter="_showP2g(this)" onmouseleave="_hideP2g(false)" onclick="_pinP2g(this)"
    style="cursor:pointer;font-size:10px;color:var(--green);font-weight:600;border-bottom:1px dotted var(--green)">P2G ⓘ</span></div>`;
  return dots + p2gLink;
}}

function summaryHtml(r, idx) {{
  const text = (r.summary ? r.summary.slice(0,140) + (r.summary.length > 140 ? '…' : '') : '')
    || (r.pulse_action ? `<span style="font-size:11px;color:var(--red)">⚠ ${{r.pulse_action.slice(0,120)}}${{r.pulse_action.length>120?'…':''}}</span>` : '')
    || (r.slack_intel  ? `<span style="font-size:11px;color:var(--subtext)">💬 ${{r.slack_intel.slice(0,120)}}${{r.slack_intel.length>120?'…':''}}</span>` : '');
  if (!r.has_pulse) return `<div class="summary-cell"><div class="summary-short">${{text}}</div></div>`;
  const sumData = JSON.stringify({{
    pulse_scope:    r.pulse_scope    || '',
    pulse_sched:    r.pulse_sched    || '',
    pulse_budget:   r.pulse_budget   || '',
    pulse_resource: r.pulse_resource || '',
    pulse_customer: r.pulse_customer || '',
    ptg_scope:      r.ptg_scope      || '',
    ptg_sched:      r.ptg_sched      || '',
    ptg_budget:     r.ptg_budget     || '',
    ptg_resource:   r.ptg_resource   || '',
    ptg_customer:   r.ptg_customer   || '',
    pulse_id:       r.pulse_id       || '',
    pulse_updated:  r.pulse_updated  || '',
    pulse_trend:    r.pulse_trend    || '',
    pulse_steerco:  r.pulse_steerco  || '',
    pulse_golive:   r.pulse_golive   || '',
    pulse_action:   r.pulse_action   || '',
    summary:        r.summary        || '',
    leadership:     r.leadership     || '',
    high_watch:     r.high_watch     || false,
  }}).replace(/"/g,'&quot;');
  const proj = (r.name || r.acct || '').replace(/"/g,'&quot;');
  return `<div class="summary-cell"
    data-sum="${{sumData}}" data-proj="${{proj}}"
    onmouseenter="_showSummary(this)" onmouseleave="_hideSummary(false)" onclick="_pinSummary(this)"
    style="cursor:pointer">
    <div class="summary-short">${{text}}</div>
  </div>`;
}}

function displayName(n) {{
  return n.replace(/[ \t-]*-[ \t-]*(CSG|CO[0-9]+).*?-Project[ \t]*$/i, '').trim();
}}

function projectRow(r, idx) {{
  const hwFlag  = r.high_watch ? '<span class="hw-flag">⚑ HW</span>' : '';
  const sweFlag = r.swe_co ? '<span class="hw-flag" style="background:#E8F5E9;color:#2E7D32">SWE</span>' : '';
  const newFlag = (r.start_dt && r.start_dt > new Date().toISOString().slice(0,10)) ? '<span class="hw-flag" style="background:#FFF8E1;color:#F57F17">🆕 NEW</span>' : '';
  return `<tr data-idx="${{idx}}" data-grp="${{(r[groupBy]||'').toString().replace(/"/g,'&quot;')}}">
    <td class="col-status">${{statusBadge(r)}}</td>
    <td class="col-tier" style="text-align:center">
      <span class="tier">${{TIER_LABEL[r.tier]||r.tier}}</span>
      ${{_DB_CRED ? `<button class="assign-btn" onclick="openAssignModal('${{r.pid}}')" title="Edit tier / PO">✏️</button>` : ''}}
    </td>
    <td class="col-project">
      <div class="proj-name">${{hwFlag}}${{sweFlag}}${{newFlag}}<a href="${{r.url}}" target="_blank" class="proj-link">${{displayName(r.name)}}</a></div>
      <div class="proj-acct">${{r.acct}}</div>
      ${{(r.start_dt || r.end_dt) ? (() => {{
        const today = new Date().toISOString().slice(0,10);
        const endColor = !r.end_dt ? 'var(--subtext)'
          : r.end_dt < today ? 'var(--red)'
          : (new Date(r.end_dt) - new Date(today)) / 86400000 <= 45 ? 'var(--yellow)'
          : 'var(--subtext)';
        return `<div style="font-size:10px;color:var(--subtext);margin-top:2px"><span>Start: ${{r.start_dt||'?'}}</span> &nbsp;·&nbsp; <span>End: </span><span style="color:${{endColor}};font-weight:${{endColor!=='var(--subtext)'?'600':'400'}}">${{r.end_dt||'?'}}</span></div>`;
      }})() : ''}}
      ${{(r.stage || r.practice) ? `<div style="font-size:10px;color:#4B5563">${{[r.stage,r.practice].filter(Boolean).join(' · ')}}</div>` : ''}}
      ${{r.pulse_golive ? `<div style="font-size:10px;margin-top:2px"><span style="color:var(--subtext)">Go-Live: </span><span style="font-weight:600;color:var(--blue)">${{r.pulse_golive}}</span></div>` : ''}}
      ${{r.ironclad_url ? `<div style="font-size:10px;margin-top:3px"><a href="${{r.ironclad_url}}" target="_blank" rel="noopener" style="color:var(--blue);text-decoration:none;font-weight:500">📄 Contract ↗</a></div>` : ''}}
    </td>
    <td class="col-team">
      ${{teamHtml((r.team||'').split(' · ').filter(s=>!s.startsWith('PO:')).join(' · '))}}
      <div style="font-size:10px;margin-top:3px;color:var(--subtext)">PO: <span style="color:var(--text);font-weight:600">${{r.po||'—'}}</span>${{_DB_CRED ? ` <button class="assign-btn" onclick="openAssignModal('${{r.pid}}')" title="Edit tier / PO">✏️</button>` : ''}}</div>
    </td>
    <td class="col-resource">${{resourcingHtml(r)}}</td>
    <td class="col-fin">${{finHtml(r)}}</td>
    <td class="col-rules">${{rulesHtml(r.rules)}}</td>
    <td class="col-bl">${{baselinesHtml(r.baselines, r)}}</td>
    <td>${{summaryHtml(r, idx)}}</td>
  </tr>`;
}}

const _STATUS_GRP_ICON = {{'Red':'🔴','Yellow':'🟡','Watermelon':'🍉','Green':'🟢','No Pulse':'⚫','On Hold':'⏸'}};
function grpLabel(dim, key) {{
  if (dim === 'tier') return TIER_NAME[parseInt(key)] || ('Tier ' + key);
  if (dim === 'status') return (_STATUS_GRP_ICON[key] || '') + ' ' + (key || 'No Pulse');
  return key || '(Unassigned)';
}}

function grpSummary(rows) {{
  const cnt  = rows.length;
  const bk   = rows.reduce((s,r)=>s+(r.bookings||0),0);
  const bil  = rows.reduce((s,r)=>s+(r.billings||0),0);
  const rem  = rows.reduce((s,r)=>s+(r.actuals_rem||0),0);
  const reds = rows.filter(r=>r.status==='Red').length;
  const yels = rows.filter(r=>r.status==='Yellow').length;
  const wms  = rows.filter(r=>r.status==='Watermelon').length;
  const grns = rows.filter(r=>r.status==='Green').length;
  const nops = rows.filter(r=>r.status==='No Pulse').length;
  let pills = '';
  if (reds) pills += `<span class="grp-pill grp-pill-red">🔴 ${{reds}}</span>`;
  if (yels) pills += `<span class="grp-pill grp-pill-yellow">🟡 ${{yels}}</span>`;
  if (wms)  pills += `<span class="grp-pill grp-pill-wm">🍉 ${{wms}}</span>`;
  if (grns) pills += `<span class="grp-pill grp-pill-green">🟢 ${{grns}}</span>`;
  if (nops) pills += `<span class="grp-pill grp-pill-nopulse">⚫ ${{nops}}</span>`;
  const finParts = [`Bkd: ${{fmtMoney(bk)}}`];
  if (bil)  finParts.push(`Bil: ${{fmtMoney(bil)}}`);
  if (rem)  finParts.push(`Rem: ${{fmtMoney(rem)}}`);
  return `<span class="grp-meta">${{cnt}} proj &nbsp;·&nbsp; ${{finParts.join(' &nbsp;·&nbsp; ')}}</span>
          <span class="grp-pills">${{pills}}</span>`;
}}

// ── N-level recursive group renderer ──────────────────────────────────────────
function _grpGetKey(r, dim) {{
  if (dim === 'tier') return String(r.tier || 9);
  if (dim === 'po')   return r.po || 'Unassigned';
  if (dim === 'acct') return r.acct || '(No Account)';
  if (dim === 'region') return r.region || 'Unknown';
  if (dim === 'status') return r.status || 'No Pulse';
  return '';
}}

function _grpSortKeys(keys, dim) {{
  if (dim === 'tier') return [...keys].sort((a,b) => (parseInt(a)||99)-(parseInt(b)||99));
  if (dim === 'status') return [...keys].sort((a,b) => (_STATUS_GRP_ORDER[a]??99) - (_STATUS_GRP_ORDER[b]??99));
  return [...keys].sort((a,b) => a.localeCompare(b));
}}

function _renderGroupLevel(rows, dims, level, ancestorGids) {{
  if (!dims.length) {{
    // leaf level — render project rows
    return rows.map((r,i) => {{
      let tr = projectRow(r, i);
      // tag with all ancestor group ids for collapse/expand targeting
      const attrs = ancestorGids.map((g,d) => `data-grp-${{d}}="${{g}}"`).join(' ');
      return tr.replace('<tr ', `<tr ${{attrs}} data-grp-level="${{level}}" `);
    }}).join('');
  }}
  const dim = dims[0];
  const rest = dims.slice(1);
  // bucket rows by this dimension
  const buckets = {{}}, order = [];
  rows.forEach(r => {{
    const k = _grpGetKey(r, dim);
    if (!buckets[k]) {{ buckets[k] = []; order.push(k); }}
    buckets[k].push(r);
  }});
  const sorted = _grpSortKeys(order, dim);
  let html = '';
  sorted.forEach(k => {{
    const bRows = buckets[k];
    const gid = 'grp-' + level + '-' + k.replace(/[^a-z0-9]/gi,'_') + '-' + ancestorGids.join('_');
    const indent = level * 14;
    const allAncestors = [...ancestorGids, gid];
    // group header
    const attrs = ancestorGids.map((g,d) => `data-grp-${{d}}="${{g}}"`).join(' ');
    html += `<tr class="group-row" data-level="${{level}}" ${{attrs}} onclick="toggleGrp('${{gid}}',this)">
      <td colspan="9" style="padding-left:${{14+indent}}px">
        <span class="grp-toggle" id="arrow-${{gid}}">▼</span>
        <strong>${{grpLabel(dim, k)}}</strong>
        ${{grpSummary(bRows)}}
      </td>
    </tr>`;
    // children
    html += _renderGroupLevel(bRows, rest, level + 1, allAncestors);
  }});
  return html;
}}

function renderRows(data) {{
  const tbody = document.getElementById('table-body');
  const noRes = document.getElementById('no-results');
  if (!data.length) {{
    tbody.innerHTML = '';
    noRes.style.display = 'block';
    return;
  }}
  noRes.style.display = 'none';
  if (!groupKeys.length) {{
    tbody.innerHTML = data.map((r,i) => projectRow(r,i)).join('');
  }} else {{
    tbody.innerHTML = _renderGroupLevel(data, groupKeys, 0, []);
  }}
  document.getElementById('cnt-all').textContent = `(${{data.length}})`;
}}

// Collapse/expand level tracking: 0=fully expanded … N=only top-level groups visible
let _collapseLevel = 0;

// ── N-level collapse/expand ────────────────────────────────────────────────────
// All group headers carry data-level="0" (outermost) … data-level="N-1" (innermost).
// Leaf rows carry data-grp-0="gid0" … data-grp-K="gidK" for every ancestor.

function _grpRowsForGid(gid, level) {{
  // Returns all descendant rows (group headers deeper than level + leaf rows) for this gid.
  return document.querySelectorAll(`tr[data-grp-${{level}}="${{gid}}"]`);
}}

function toggleGrp(gid, headerRow) {{
  const level = parseInt(headerRow.dataset.level || '0');
  const arrow = document.getElementById('arrow-' + gid);
  const descendants = _grpRowsForGid(gid, level);
  const nowHidden = descendants.length && descendants[0].classList.contains('grp-hidden');
  descendants.forEach(r => r.classList.toggle('grp-hidden', !nowHidden));
  if (arrow) arrow.textContent = nowHidden ? '▼' : '▶';
}}

function _currentCollapseLevel() {{
  // Detect actual DOM state by checking descendants, not headers themselves.
  // At collapse level cl, step cl hid all tr[data-grp-(N-cl)] descendants.
  // We scan from most-collapsed to least: first cl where ALL those rows are hidden.
  const N = groupKeys.length;
  if (N === 0) return 0;
  for (let cl = N; cl >= 1; cl--) {{
    const rows = document.querySelectorAll(`tr[data-grp-${{N - cl}}]`);
    if (rows.length === 0) continue;
    if ([...rows].every(r => r.classList.contains('grp-hidden'))) return cl;
  }}
  return 0;
}}

function stepCollapse() {{
  const N = groupKeys.length;
  if (N === 0) return;
  _collapseLevel = _currentCollapseLevel();
  if (_collapseLevel >= N) return; // already fully collapsed
  const targetLevel = N - 1 - _collapseLevel;
  document.querySelectorAll(`tr.group-row[data-level="${{targetLevel}}"]`).forEach(hdr => {{
    const gid = hdr.getAttribute('onclick')?.match(/toggleGrp\\('([^']+)'/)?.[1];
    if (!gid) return;
    _grpRowsForGid(gid, targetLevel).forEach(r => r.classList.add('grp-hidden'));
    const arrow = document.getElementById('arrow-' + gid);
    if (arrow) arrow.textContent = '▶';
  }});
  _collapseLevel++;
}}

function stepExpand() {{
  const N = groupKeys.length;
  if (N === 0) return;
  _collapseLevel = _currentCollapseLevel();
  if (_collapseLevel === 0) {{
    // Mixed state: some groups manually collapsed but level counter reads 0.
    // Expand everything to a clean fully-open state.
    const anyHidden = document.querySelector('tr.grp-hidden');
    if (!anyHidden) return;
    document.querySelectorAll('tr.grp-hidden').forEach(r => r.classList.remove('grp-hidden'));
    document.querySelectorAll('[id^="arrow-"]').forEach(a => a.textContent = '▼');
    return;
  }}
  const targetLevel = N - _collapseLevel;
  document.querySelectorAll(`tr.group-row[data-level="${{targetLevel}}"]`).forEach(hdr => {{
    const isVisible = [...Array(targetLevel).keys()].every(l => {{
      const ancestorGid = hdr.dataset['grp' + l];
      if (!ancestorGid) return true;
      const ancestorRow = document.querySelector(`tr.group-row[data-level="${{l}}"][onclick*="${{ancestorGid}}"]`);
      return !ancestorRow || !ancestorRow.classList.contains('grp-hidden');
    }});
    if (!isVisible) return;
    const gid = hdr.getAttribute('onclick')?.match(/toggleGrp\\('([^']+)'/)?.[1];
    if (!gid) return;
    hdr.classList.remove('grp-hidden');
    const arrow = document.getElementById('arrow-' + gid);
    if (arrow) arrow.textContent = '▼';
    _grpRowsForGid(gid, targetLevel).forEach(r => r.classList.remove('grp-hidden'));
  }});
  _collapseLevel--;
}}

function updateScorecard(data) {{
  const n = data.length;
  const bk  = data.reduce((s,r) => s + (r.bookings||0), 0);
  const bil = data.reduce((s,r) => s + (r.billings||0), 0);
  const far = data.reduce((s,r) => s + (r.far||0), 0);
  const ov  = data.reduce((s,r) => s + (r.overdue_inv||0), 0);
  const rr  = data.reduce((s,r) => s + (r.rr_revenue||0), 0);
  const statuses = data.map(r => r.status);
  const cnt = s => statuses.filter(v => v===s).length;
  // weighted margin (exclude SWE/ARI and zero-booking projects)
  const mProj = data.filter(r => !r.swe_co && r.bookings > 0 && r.bid_margin_raw !== null && r.close_margin_raw !== null);
  let wBid = 0, wClose = 0;
  if (mProj.length) {{
    const totBk = mProj.reduce((s,r)=>s+(r.bookings||0),0);
    if (totBk > 0) {{
      wBid   = mProj.reduce((s,r)=>s+(r.bookings||0)*(r.bid_margin_raw||0),0) / totBk;
      wClose = mProj.reduce((s,r)=>s+(r.bookings||0)*(r.close_margin_raw||0),0) / totBk;
    }}
  }}
  document.getElementById('sc-total').textContent    = n;
  document.getElementById('sc-bookings').textContent = fmtMoney(bk);
  document.getElementById('sc-avg-bk').textContent   = fmtMoney(n > 0 ? bk / n : 0);
  const _bks = data.map(r => r.bookings || 0).sort((a,b) => a - b);
  const _mid = Math.floor(_bks.length / 2);
  const _med = _bks.length === 0 ? 0 : (_bks.length % 2 ? _bks[_mid] : (_bks[_mid-1] + _bks[_mid]) / 2);
  document.getElementById('sc-median-bk').textContent = fmtMoney(_med);
  document.getElementById('sc-billings').textContent = fmtMoney(bil);
  document.getElementById('sc-backlog').textContent  = fmtMoney(bk - bil);
  document.getElementById('sc-red').textContent         = cnt('Red');
  document.getElementById('sc-yellow').textContent      = cnt('Yellow');
  document.getElementById('sc-wm').textContent          = cnt('Watermelon');
  document.getElementById('sc-green').textContent       = cnt('Green');
  document.getElementById('sc-green-total').textContent = cnt('Green') + cnt('Watermelon');
  document.getElementById('sc-nopulse').textContent     = cnt('No Pulse');
  document.getElementById('sc-far').textContent      = fmtMoney(far);
  document.getElementById('sc-overdue').textContent  = fmtMoney(ov);
  document.getElementById('sc-rr').textContent       = fmtMoney(rr);
  const bidEl = document.getElementById('sc-bid');
  bidEl.textContent = fmtPct(wBid);
  bidEl.style.color = wBid > 13 ? 'var(--green)' : wBid >= 5 ? 'var(--yellow)' : 'var(--red)';
  const closeEl = document.getElementById('sc-close');
  closeEl.textContent = fmtPct(wClose);
  closeEl.style.color = wClose < 0 ? 'var(--red)' : '';
  document.getElementById('sc-delta').textContent    = (wClose - wBid >= 0 ? '+' : '') + fmtPct(wClose - wBid);
  // Avg Health Risk, Data Quality, CSAT — recalculate from filtered data
  const hrVals   = data.map(r => r.health_risk).filter(v => v != null);
  const dqVals   = data.map(r => r.data_quality).filter(v => v != null);
  const csatVals = data.map(r => r.csat_score).filter(v => v != null);
  const _avg = arr => arr.length ? arr.reduce((s,v)=>s+v,0)/arr.length : null;
  const _scoreCls = v => v == null ? '' : v < 60 ? 'red' : v < 75 ? 'yellow' : 'green';
  const _csatCls  = v => v == null ? '' : v < 3.5 ? 'red' : v < 4.0 ? 'yellow' : 'green';
  const avgHr   = _avg(hrVals);
  const avgDq   = _avg(dqVals);
  const avgCsat = _avg(csatVals);
  const hrEl   = document.getElementById('sc-avg-hr');
  const dqEl   = document.getElementById('sc-avg-dq');
  const csatEl = document.getElementById('sc-avg-csat');
  if (hrEl)   {{ hrEl.textContent   = avgHr   != null ? avgHr.toFixed(0)   : '—'; hrEl.closest('.stat').className   = 'stat ' + _scoreCls(avgHr); }}
  if (dqEl)   {{ dqEl.textContent   = avgDq   != null ? avgDq.toFixed(0)   : '—'; dqEl.closest('.stat').className   = 'stat ' + _scoreCls(avgDq); }}
  if (csatEl) {{ csatEl.textContent = avgCsat != null ? avgCsat.toFixed(1) : '—'; csatEl.closest('.stat').className = 'stat ' + _csatCls(avgCsat); }}
  document.getElementById('header-summary').innerHTML =
    n + ' projects &nbsp;|&nbsp; ' + fmtMoney(bk) + ' bookings';
  const barEl = document.getElementById('sc-bar-summary');
  if (barEl) barEl.innerHTML =
    n + ' projects &nbsp;·&nbsp; ' + fmtMoney(bk) + ' bookings &nbsp;·&nbsp; 🟢 ' +
    (cnt('Green') + cnt('Watermelon')) + ' &nbsp;·&nbsp; 🟡 ' + cnt('Yellow') +
    ' &nbsp;·&nbsp; 🔴 ' + cnt('Red');
}}

function setRegion(region, btn) {{
  filterRegion = region;
  try {{ region ? localStorage.setItem(_REGION_KEY, region) : localStorage.removeItem(_REGION_KEY); }} catch(e) {{}}
  document.querySelectorAll('.rtab').forEach(b => b.classList.remove('active'));
  if (btn) btn.classList.add('active');
  applyFilters();
}}

function applyFilters() {{
  const q = (document.getElementById('search').value || '').toLowerCase();
  // Show Clear All button if any filter is active
  const anyActive = q || filterRegion || filterStatuses.size || filterTiers.size ||
    filterPOs.size || filterRules.size || filterHW || filterSWE ||
    filterTypes.size || filterBilling.size;
  const clearBtn = document.getElementById('clear-filters-btn');
  if (clearBtn) {{
    clearBtn.style.display = anyActive ? '' : 'none';
    clearBtn.style.background    = anyActive ? '#FFC107' : '';
    clearBtn.style.borderColor   = anyActive ? '#E0A800' : '';
    clearBtn.style.color         = anyActive ? '#1a1a1a'  : '';
    clearBtn.style.fontWeight    = anyActive ? '700'      : '';
  }}
  let data = RAW.filter(r => {{
    if (filterRegion && r.region !== filterRegion) return false;
    if (filterStatuses.size > 0 && !filterStatuses.has(r.status)) return false;
    if (filterTiers.size > 0 && !filterTiers.has(String(r.tier))) return false;
    if (filterPOs.size > 0   && !filterPOs.has(r.po))             return false;
    if (filterRules.size > 0) {{
      const rowRules = (r.rules || '').split(', ').filter(Boolean);
      if (!rowRules.some(rc => filterRules.has(rc))) return false;
    }}
    if (filterHW  && !r.high_watch) return false;
    if (filterSWE && !r.swe_co)     return false;
    if (filterTypes.size > 0) {{
      const isAri = r.swe_co && (r.name||'').toLowerCase().includes('ari');
      const todayStr = new Date().toISOString().slice(0,10);
      const isNotStarted = r.start_dt && r.start_dt > todayStr;
      const match = (filterTypes.has('hw') && r.high_watch) ||
                    (filterTypes.has('swe') && r.swe_co && !isAri) ||
                    (filterTypes.has('ari') && isAri) ||
                    (filterTypes.has('not_started') && isNotStarted);
      if (!match) return false;
    }}
    if (filterBilling.size > 0 && !filterBilling.has(r.billing_type)) return false;
    if (q && !([r.name, r.acct, r.team, r.po, r.rules, r.summary, r.baselines].join(' ').toLowerCase().includes(q))) return false;
    return true;
  }});

  data = data.slice().sort((a, b) => {{
    let av = a[sortCol], bv = b[sortCol];
    if (sortCol === 'status') {{ av = STATUS_ORDER[av]??99; bv = STATUS_ORDER[bv]??99; }}
    if (typeof av === 'string') av = av.toLowerCase();
    if (typeof bv === 'string') bv = bv.toLowerCase();
    if (av < bv) return -sortDir;
    if (av > bv) return  sortDir;
    return 0;
  }});

  // groupKeys is maintained by grpChange() — no read from DOM here
  const collapseBtns = document.getElementById('grp-collapse-btns');
  if (collapseBtns) collapseBtns.style.display = groupKeys.length ? 'inline-flex' : 'none';
  _gmCurrentData = data;
  updateScorecard(data);
  renderRows(data);
}}

function togglePill(el) {{
  const filter = el.dataset.filter;
  if (filter === 'hw') {{
    el.classList.toggle('active');
    filterHW = el.classList.contains('active');
  }} else if (filter === 'swe') {{
    el.classList.toggle('active');
    filterSWE = el.classList.contains('active');
  }}
  applyFilters();
}}

// ── Multi-select combobox helpers ──────────────────────────────────────────────
const MS_CONFIG = {{
  status:  {{ set: () => filterStatuses, btnId: 'ms-status-btn',  ddId: 'ms-status-dd',  label: 'All Statuses', plural: (n) => n + ' Status' + (n>1?'es':'') }},
  tier:    {{ set: () => filterTiers,    btnId: 'ms-tier-btn',    ddId: 'ms-tier-dd',    label: 'All Tiers',    plural: (n) => n + ' Tier' + (n>1?'s':'') }},
  po:      {{ set: () => filterPOs,      btnId: 'ms-po-btn',      ddId: 'ms-po-dd',      label: 'All Owners',   plural: (n) => n + ' Owner' + (n>1?'s':'') }},
  rule:    {{ set: () => filterRules,    btnId: 'ms-rule-btn',    ddId: 'ms-rule-dd',    label: 'All Rules',    plural: (n) => n + ' Rule' + (n>1?'s':'') }},
  type:    {{ set: () => filterTypes,    btnId: 'ms-type-btn',    ddId: 'ms-type-dd',    label: 'All Types',    plural: (n) => n + ' Type' + (n>1?'s':'') }},
  billing: {{ set: () => filterBilling,  btnId: 'ms-billing-btn', ddId: 'ms-billing-dd', label: 'All Billing',  plural: (n) => n + ' Billing' }},
}};

// ── Group-by combobox ──────────────────────────────────────────────────────────
const GRP_DIM_LABELS = {{ region: 'Region', tier: 'Tier', po: 'Portfolio Owner', acct: 'Account', status: 'Status' }};
const _STATUS_GRP_ORDER = {{'Red':0,'Yellow':1,'Watermelon':2,'No Pulse':3,'Green':4,'On Hold':5}};

function toggleGrpDd() {{
  const dd = document.getElementById('ms-grp-dd');
  // close all filter dropdowns
  Object.values(MS_CONFIG).forEach(cfg => document.getElementById(cfg.ddId).classList.remove('open'));
  dd.classList.toggle('open');
}}

function grpChange(cb) {{
  // Maintain selection order: append on check, remove on uncheck
  if (cb) {{
    if (cb.checked) {{
      if (!groupKeys.includes(cb.value)) groupKeys.push(cb.value);
    }} else {{
      groupKeys = groupKeys.filter(k => k !== cb.value);
    }}
  }} else {{
    // Called without a checkbox (e.g. programmatic reset) — sync from DOM
    const checked = [...document.querySelectorAll('#ms-grp-dd input[type=checkbox]:checked')].map(c => c.value);
    groupKeys = groupKeys.filter(k => checked.includes(k));
    checked.forEach(v => {{ if (!groupKeys.includes(v)) groupKeys.push(v); }});
  }}
  // Update button label
  const btn = document.getElementById('ms-grp-btn');
  btn.textContent = groupKeys.length ? groupKeys.map(d => GRP_DIM_LABELS[d] || d).join(' › ') : 'None';
  btn.classList.toggle('ms-active', groupKeys.length > 0);
  // Update order pills
  const pillsEl = document.getElementById('grp-order-pills');
  if (pillsEl) {{
    pillsEl.innerHTML = groupKeys.map((d,i) =>
      `<span class="grp-order-pill">${{i+1}}. ${{GRP_DIM_LABELS[d]||d}}</span>`
    ).join('');
  }}
  _collapseLevel = 0;
  applyFilters();
}}

function grpClear() {{
  document.querySelectorAll('#ms-grp-dd input[type=checkbox]').forEach(cb => cb.checked = false);
  groupKeys = [];
  const btn = document.getElementById('ms-grp-btn');
  btn.textContent = 'None';
  btn.classList.remove('ms-active');
  const pillsEl = document.getElementById('grp-order-pills');
  if (pillsEl) pillsEl.innerHTML = '';
  _collapseLevel = 0;
  document.getElementById('ms-grp-dd').classList.remove('open');
  applyFilters();
}}

function toggleMs(key) {{
  const cfg = MS_CONFIG[key];
  const dd = document.getElementById(cfg.ddId);
  // close grp dropdown + all others
  document.getElementById('ms-grp-dd').classList.remove('open');
  Object.keys(MS_CONFIG).forEach(k => {{
    if (k !== key) document.getElementById(MS_CONFIG[k].ddId).classList.remove('open');
  }});
  dd.classList.toggle('open');
}}

function msChange(key) {{
  const cfg = MS_CONFIG[key];
  const s = cfg.set();
  s.clear();
  document.querySelectorAll('#' + cfg.ddId + ' input[type=checkbox]:checked').forEach(cb => s.add(cb.value));
  const btn = document.getElementById(cfg.btnId);
  btn.textContent = s.size === 0 ? cfg.label : cfg.plural(s.size);
  btn.classList.toggle('ms-active', s.size > 0);
  applyFilters();
}}

function msClear(key, skipApply) {{
  const cfg = MS_CONFIG[key];
  cfg.set().clear();
  document.querySelectorAll('#' + cfg.ddId + ' input[type=checkbox]').forEach(cb => cb.checked = false);
  const btn = document.getElementById(cfg.btnId);
  btn.textContent = cfg.label;
  btn.classList.remove('ms-active');
  if (!skipApply) applyFilters();
}}

function clearAllFilters() {{
  // Clear all multi-select filters without triggering applyFilters each time
  Object.keys(MS_CONFIG).forEach(key => msClear(key, true));
  // Clear search
  const searchEl = document.getElementById('search');
  searchEl.value = '';
  // Clear HW / SWE pills
  filterHW = false; filterSWE = false;
  document.querySelectorAll('.pill[data-filter]').forEach(p => p.classList.remove('active'));
  // Reset all filters except region — region persists until user explicitly changes tab
  // (filterRegion and active tab remain unchanged)
  // Preserve current groupKeys and collapse level across the filter reset
  const savedGroupKeys = [...groupKeys];
  const savedCollapse = _collapseLevel;
  applyFilters();
  // Re-apply collapse if grouping is still active
  if (savedGroupKeys.length && savedCollapse > 0) {{
    _collapseLevel = 0; // reset so stepCollapse starts fresh
    for (let i = 0; i < savedCollapse; i++) stepCollapse();
  }}
}}

// close dropdowns when clicking outside
document.addEventListener('click', e => {{
  if (!e.target.closest('.ms-wrap')) {{
    Object.values(MS_CONFIG).forEach(cfg => document.getElementById(cfg.ddId).classList.remove('open'));
    const grpDd = document.getElementById('ms-grp-dd');
    if (grpDd) grpDd.classList.remove('open');
  }}
}});

function toggleExpand(btn, idx) {{
  const row = btn.closest('tr');
  row.classList.toggle('expanded');
  btn.textContent = row.classList.contains('expanded') ? '▲ less' : '▼ more';
}}

document.querySelectorAll('thead th[data-col]').forEach(th => {{
  th.addEventListener('click', () => {{
    const col = th.dataset.col;
    if (sortCol === col) {{ sortDir *= -1; }}
    else {{ sortCol = col; sortDir = 1; }}
    document.querySelectorAll('thead th').forEach(t => t.classList.remove('sorted-asc','sorted-desc'));
    th.classList.add(sortDir === 1 ? 'sorted-asc' : 'sorted-desc');
    applyFilters();
  }});
}});

document.querySelector('[data-col="status"]').classList.add('sorted-asc');

// ── GM Business Overview ───────────────────────────────────────────────────────

const _GM_SYS = `Act as a Senior Business Analyst presenting a high-level operational briefing to the Professional Services General Manager (GM).

Analyze the provided portfolio data and construct an executive overview using the exact markdown formatting and sections below. Focus heavily on financial risk, resource efficiencies, data hygiene gaps, and upcoming critical milestones.

Keep these core structural definitions and accounting mechanics in mind during your financial analysis:
1. SWE designates internal investment money (where bookings/billings are zero, meaning the goal is resource efficiency, timeline baselines, and margin preservation).
2. FAR stands for Forecasted Amount Remaining—a critical financial metric tracking expected revenue remaining in contract funds for T&M projects (FAR = Bookings - Billings - Scheduled Work - Pending Resource Requests).
   - A Negative FAR (-) variance indicates the project is actively over-staffed or over budget, requiring immediate onshore cost-containment, hour reductions, or a Change Order.
   - A Positive FAR (+) variance indicates under-staffing or under-utilization, meaning funded scope is stalling and leaving money on the table.
   - Thresholds (RAG): A FAR variance > 15% of total T&M SOW bookings is flagged as Red; 5%-14% is Yellow; < 5% is Green.

Output using exactly this section layout. Use **bold** for key numbers. Use 🔴 🟡 🟢 🍉 ⚠️ ⚑ emojis for visual scanning. Output plain markdown (### headers, ** bold, - bullets, markdown tables).

---

### [Portfolio Name / Region] Executive Briefing
**To:** Professional Services General Manager
**From:** Portfolio Business Analyst
**Date:** [Insert Date]
**Portfolio Scope:** [X] Active Projects | $[X]M Total Bookings

---

### 📊 Performance at a Glance
Provide a high-level bulleted summary of macro metrics from the data scorecard:
*   Total Billings vs. Remaining Backlog (Calculate the portfolio's financial burn percentage).
*   Margin Optimization (Compare Weighted Bid Margin % against current Weighted Margin at Close % to show the exact Margin Delta optimization).
*   Client Satisfaction (Average CSAT score, but include a data hygiene caveat highlighting the exact percentage of active projects that actually have a survey on file vs. those operating blind with missing/lapsed data).
*   Data Hygiene (Average Data Quality score).

---

### 🚦 Portfolio Health & Working Capital Risks
Analyze the data and summarize the hidden risks using these specific categories:
*   The Overdue Cash Bottleneck: Highlight total overdue invoices and name the top enterprise account culprits driving the cash freeze.
*   The "Watermelon" Transparency Gap: Count how many projects are marked "Green" overall by PMs but are classified as "Watermelon Green" due to internal red flags (e.g., budget overruns, negative FAR, 0% GDC offshore utilization, or past end dates).
*   Financial At Risk (FAR): Summarize total dollar exposure trapped across the pipeline, breaking down where projects are actively leaking funds via overburns (-) vs. where revenue velocity is stalling via under-utilization (+).

---

### 🔍 Account Ledger & High Watch Status Summary
Create a comprehensive markdown table explicitly tracking all High-Watch customers (denoted by ⚑ HW) and any projects carrying a 🔴 RED status.

For each row, include:
1. Status (Use Emoji: 🔴 RED, 🟡 YELLOW, 🍉 WATERMELON)
2. Account & Project Name (Prefix with ⚑ HW if it is a High Watch customer)
3. Leadership Team (List the PM and Account Owner)
4. Key Financials (Bookings, Billings, Overdue Invoices, and FAR variance)
5. Executive Summary, Key Business Outcomes, & Top Risks — synthesize the raw data columns and the complete summary/leadership notes into a clean professional narrative structured exactly like this:
   - **Executive Summary:** A concise overview of the project's contract type, scope parameters, and lifecycle stage.
   - **Key Business Outcome:** The core business or technical milestone the client is trying to achieve.
   - **Top Risks & Actions:** A real-time breakdown of the exact operational issues occurring (e.g., product deployability, security gaps, client funding cuts, 0% offshore skews, or approaching deadlines) and the precise tactical steps being taken to resolve them.

---

### 🛠️ GM Immediate Intervention Actions
Provide exactly 5 clear, highly actionable executive directives for the GM to execute immediately based on the risks identified above. Include specific project names, target dates, resource reallocation goals, or executive alignment escalations needed to protect portfolio health.`;


const _GM_CHAT_SYS = `You are a knowledgeable professional services portfolio analyst assistant. You have access to the user's current filtered portfolio data (projects, financials, statuses, risks). Answer questions conversationally and directly — no fixed structure, no forced sections. Use the portfolio data to give specific, grounded answers: reference real project names, account names, dollar amounts, and statuses. Be concise but thorough. Use bullet points only when listing multiple items. Use **bold** for key figures. Never fabricate data not present in the portfolio context.`;

function _fmtM(v) {{
  if (!v && v !== 0) return '—';
  const a = Math.abs(v);
  if (a >= 1e6) return (v < 0 ? '-' : '') + '$' + (a / 1e6).toFixed(2) + 'M';
  if (a >= 1e3) return (v < 0 ? '-' : '') + '$' + (a / 1e3).toFixed(0) + 'K';
  return (v < 0 ? '-' : '') + '$' + a.toFixed(0);
}}

function buildGMPrompt(data) {{
  const n = data.length;
  const sum = f => data.reduce((s, r) => s + (r[f] || 0), 0);
  const totalBk  = sum('bookings');
  const totalBil = sum('billings');
  const totalFAR = sum('far');
  const totalOv  = sum('overdue_inv');
  const totalRR  = sum('rr_revenue');
  const totalUnsch = data.reduce((s,r) => s + (r.unsch_backlog||0), 0);

  // Weighted margins
  let wBidNum = 0, wCloseNum = 0, wDen = 0;
  data.forEach(r => {{
    if (r.bookings > 0 && r.bid_margin_raw != null && r.close_margin_raw != null) {{
      wBidNum   += r.bid_margin_raw   * r.bookings;
      wCloseNum += r.close_margin_raw * r.bookings;
      wDen      += r.bookings;
    }}
  }});
  const wBid   = wDen > 0 ? (wBidNum   / wDen * 100).toFixed(1) : 'N/A';
  const wClose = wDen > 0 ? (wCloseNum / wDen * 100).toFixed(1) : 'N/A';
  const delta  = wDen > 0 ? ((wCloseNum - wBidNum) / wDen * 100).toFixed(1) : 'N/A';

  // Status counts
  const byStatus = s => data.filter(r => r.status === s).length;
  const reds  = data.filter(r => r.status === 'Red');
  const yellows = data.filter(r => r.status === 'Yellow');
  const wms   = data.filter(r => r.status === 'Watermelon');
  const noPulse = data.filter(r => r.status === 'No Pulse');
  const onHold  = data.filter(r => r.status === 'On Hold');

  // Operational leaks: active projects with $0 bookings
  const zeroActive = data.filter(r => (r.bookings||0) === 0 && r.status !== 'On Hold' && r.status !== 'No Pulse');

  // GDC coverage
  const lowGdc = data.filter(r => r.gdc_total > 0 && (r.gdc_pct||0) < 50 && r.status !== 'No Pulse' && r.status !== 'On Hold').length;

  // Data quality
  const dqScores = data.filter(r => r.data_quality != null).map(r => r.data_quality);
  const avgDQ = dqScores.length > 0 ? (dqScores.reduce((s,v) => s+v, 0) / dqScores.length).toFixed(0) : 'N/A';

  // High watch
  const hw = data.filter(r => r.high_watch);

  // Top escalations by bookings (reds + yellows + wms, up to 8)
  const escalations = [...reds, ...yellows, ...wms]
    .sort((a,b) => (b.bookings||0) - (a.bookings||0))
    .slice(0, 8)
    .map(r => `  - ${{r.acct||r.name}} (${{_fmtM(r.bookings)}}) [${{r.status}}] — rules: ${{r.rules||'none'}}`).join('\\n');

  const hwList = hw.slice(0,5).map(r => `  - ${{r.acct||r.name}} (${{_fmtM(r.bookings)}})`).join('\\n');

  const scope = filterRegion ? filterRegion : 'ACC (All Regions)';

  return `PORTFOLIO SCOPE: ${{scope}} — ${{n}} active projects

## MACRO FINANCIALS
- Total Bookings: ${{_fmtM(totalBk)}}
- Total Billings to Date: ${{_fmtM(totalBil)}}
- Backlog Remaining: ${{_fmtM(totalBk - totalBil)}}
- Weighted Bid Margin: ${{wBid}}%  |  Weighted Margin @ Close: ${{wClose}}%  |  Delivery Delta: +${{delta}}%
- Total FAR (budget remaining): ${{_fmtM(totalFAR)}}
- Overdue Invoices: ${{_fmtM(totalOv)}}
- Revenue @ Risk: ${{_fmtM(totalRR)}}
- Unscheduled Backlog: ${{_fmtM(totalUnsch)}}

## STATUS BREAKDOWN
- 🔴 Red: ${{reds.length}} projects  |  🟡 Yellow: ${{yellows.length}}  |  🍉 Watermelon: ${{wms.length}}  |  🟢 Green: ${{byStatus('Green')}}
- ⚫ No Pulse: ${{noPulse.length}} projects  |  ⏸ On Hold: ${{onHold.length}}

## TOP ESCALATIONS (Red/Yellow/Watermelon by bookings)
${{escalations || '  - None'}}

## HIGH WATCH PROJECTS
${{hwList || '  - None'}}

## OPERATIONAL LEAKS
- Active $0-Booking projects (burning resources with no contract): ${{zeroActive.length}}${{zeroActive.length > 0 ? '\\n' + zeroActive.slice(0,5).map(r=>`  - ${{r.acct||r.name}} [${{r.status}}]`).join('\\n') : ''}}
- Projects with <50% GDC coverage (onshore-heavy): ${{lowGdc}}

## DATA HYGIENE
- No-Pulse projects: ${{noPulse.length}} of ${{n}} (${{Math.round(noPulse.length/Math.max(n,1)*100)}}%)
- Avg Data Quality Score: ${{avgDQ}}%

---
Generate a GM executive briefing using the 4-section layout: (1) Performance at a Glance, (2) Portfolio Health & Working Capital Risks, (3) Account Ledger & High Watch Status Summary (markdown table with full narrative per row), (4) GM Immediate Intervention Actions (exactly 5 directives). Apply the SWE and FAR accounting rules defined in your system prompt. Be action-oriented — cite specific project names, dollar amounts, and account owners.`;
}}

function _gmRenderMarkdown(text) {{
  const lines = text.split('\\n');
  let html = '';
  let inUl = false;
  for (let i = 0; i < lines.length; i++) {{
    let line = lines[i]
      .replace(/[*][*](.+?)[*][*]/g, '<strong>$1</strong>')
      .replace(/`(.+?)`/g, '<code>$1</code>');
    if (/^[#]{{1,3}}[ ]/.test(line)) {{
      if (inUl) {{ html += '</ul>'; inUl = false; }}
      html += '<h3>' + line.replace(/^[#]+[ ]*/, '') + '</h3>';
    }} else if (/^[-•][ ]/.test(line)) {{
      if (!inUl) {{ html += '<ul>'; inUl = true; }}
      html += '<li>' + line.replace(/^[-•][ ]*/, '') + '</li>';
    }} else if (line.trim() === '') {{
      if (inUl) {{ html += '</ul>'; inUl = false; }}
    }} else {{
      if (inUl) {{ html += '</ul>'; inUl = false; }}
      html += '<p>' + line + '</p>';
    }}
  }}
  if (inUl) html += '</ul>';
  return html;
}}

function _gmSkeleton(label) {{
  return `<div class="gm-spinner-wrap">
    <div class="gm-spinner"></div>
    <div class="gm-spinner-label" id="gm-spinner-label">${{label || 'Generating overview…'}}</div>
  </div>`;
}}

async function openGMOverview(forceRegen) {{
  const modal  = document.getElementById('gm-modal');
  const body   = document.getElementById('gm-body');
  const footer = document.getElementById('gm-footer');
  const scopeEl = document.getElementById('gm-scope-label');
  const regenBtn = document.getElementById('gm-regen-btn');
  const modelEl  = document.getElementById('gm-model-used');
  const tsEl     = document.getElementById('gm-generated-at');

  // Build a cache key from the exact visible rows so any filter change busts the cache
  const visibleData = _gmCurrentData.length > 0 ? _gmCurrentData : RAW.slice();
  const _cacheHash = visibleData.map(r => r.pid||r.id||r.name).sort().join(',');
  const cacheKey = 'gm_ov_' + _cacheHash.length + '_' + (_cacheHash.slice(0,60));
  // Scope label: region + filter summary
  const scopeLabel = (filterRegion ? filterRegion + ' Region' : 'ACC Delivery Portfolio') +
    ' (' + visibleData.length + ' projects)';
  scopeEl.textContent = scopeLabel + ' — Business Overview';

  modal.classList.add('open');
  switchGmTab('overview');

  // Guard: platform injects __PROXY_TOKEN__ at page load; if missing, page must be reloaded
  if (!window.__PROXY_TOKEN__ || window.__PROXY_TOKEN__ === 'undefined') {{
    const body2 = document.getElementById('gm-body');
    body2.innerHTML = `<p style="color:var(--red)">⚠️ Session token not available.<br>
      Please <strong>reload the page</strong> and try again.<br>
      <small style="color:var(--muted)">(window.__PROXY_TOKEN__ is not set — this is injected by the platform at load time)</small></p>`;
    document.getElementById('gm-generated-at').textContent = '';
    return;
  }}

  if (!forceRegen) {{
    try {{
      const cached = sessionStorage.getItem(cacheKey);
      if (cached) {{
        const c = JSON.parse(cached);
        body.innerHTML = _gmRenderMarkdown(c.text);
        modelEl.textContent = 'Model: ' + (c.model || 'unknown');
        tsEl.textContent = 'Generated: ' + c.ts;
        return;
      }}
    }} catch(e) {{}}
  }}

  // Show spinner
  body.innerHTML = _gmSkeleton('Generating overview…');
  modelEl.textContent = '';
  tsEl.textContent = '';
  regenBtn.disabled = true;

  function _setSpinnerLabel(txt) {{
    const el = document.getElementById('gm-spinner-label');
    if (el) el.textContent = txt;
  }}

  async function _gmCall(tier) {{
    const prompt = buildGMPrompt(visibleData);
    console.log('[GM] tier=' + tier + ' prompt_len=' + prompt.length + ' data_rows=' + visibleData.length);
    const resp = await fetch('/api/proxy/llm/' + window.__UPLOAD_ID__, {{
      method: 'POST',
      headers: {{ 'Content-Type': 'application/json', 'X-Proxy-Token': window.__PROXY_TOKEN__ }},
      body: JSON.stringify({{ system: _GM_SYS, prompt, tier, maxTokens: 2000 }}),
    }});
    if (!resp.ok) {{
      const errBody = await resp.text().catch(() => '');
      console.warn('[GM] HTTP', resp.status, errBody);
      throw new Error('HTTP ' + resp.status + (errBody ? ': ' + errBody.slice(0,120) : ''));
    }}
    return resp.json();
  }}

  try {{
    let result;
    for (const tier of ['powerful', 'balanced', 'fast']) {{
      try {{
        _setSpinnerLabel(tier === 'powerful' ? 'Generating overview…' : 'Retrying with ' + tier + ' model…');
        result = await _gmCall(tier);
        break;
      }} catch(e) {{
        if (e.message.startsWith('HTTP 503') && tier !== 'fast') {{
          console.warn('[GM] ' + tier + ' tier 503, trying next');
        }} else {{
          throw e;
        }}
      }}
    }}
    const text = result.response || '';
    const model = result.model_used || result.tier || '';
    const ts = new Date().toLocaleString();

    body.innerHTML = _gmRenderMarkdown(text);
    modelEl.textContent = 'Model: ' + model;
    tsEl.textContent = 'Generated: ' + ts;

    try {{ sessionStorage.setItem(cacheKey, JSON.stringify({{ text, model, ts }})); }} catch(e) {{}}
  }} catch(e) {{
    const is401 = e.message.startsWith('HTTP 401');
    body.innerHTML = is401
      ? `<p style="color:var(--red)">⚠️ Session expired — please <strong>reload the page</strong> and try again.</p>`
      : `<p style="color:var(--red)">⚠️ Failed to generate overview: ${{e.message}}</p>
         <p><button onclick="openGMOverview(true)" style="margin-top:8px;padding:6px 14px;cursor:pointer">Try again</button></p>`;
    tsEl.textContent = '';
    console.error('[GM] final error:', e);
  }} finally {{
    regenBtn.disabled = false;
  }}
}}

function closeGMOverview() {{
  document.getElementById('gm-modal').classList.remove('open');
}}

function switchGmTab(tab) {{
  ['overview','prompts','chat','library'].forEach(t => {{
    document.getElementById('gm-tab-' + t).classList.toggle('active', t === tab);
    document.getElementById('gm-panel-' + t).classList.toggle('active', t === tab);
  }});
  document.getElementById('gm-regen-btn').style.display = tab === 'overview' ? '' : 'none';
  if (tab === 'prompts') _gmRenderPrompts();
  if (tab === 'chat') setTimeout(() => document.getElementById('gm-chat-input').focus(), 50);
  if (tab === 'library') gmlLoad();
}}

// ── Prompt Library ────────────────────────────────────────────────────────────
let _gmlFilter = 'all';
let _gmlItems  = [];   // [{{id, title, text, audience, created_at}}]
let _gmlLoaded = false;

const _GML_AUDIENCE_LABELS = {{
  portfolio_owner: 'Portfolio Owner',
  dam_lead:        'DAM Lead / GM',
  gm:              'GM',
  evp:             'EVP',
}};

async function gmlLoad(force) {{
  if (_gmlLoaded && !force) {{ gmlRender(); return; }}
  if (!_DB_CRED) {{
    document.getElementById('gm-library-list').innerHTML =
      '<div id="gm-library-empty" style="display:block">⚠️ No DB credentials — library unavailable.</div>';
    return;
  }}
  try {{
    const data = await TileDB.serverQuery('SELECT id, title, text, audience, created_at FROM prompt_library ORDER BY created_at DESC');
    _gmlItems  = data.rows || [];
    _gmlLoaded = true;
    gmlRender();
  }} catch(e) {{
    console.warn('gmlLoad:', e);
    document.getElementById('gm-library-list').innerHTML =
      '<div id="gm-library-empty" style="display:block">⚠️ Failed to load library: ' + e.message + '</div>';
  }}
}}

function gmlSetFilter(aud) {{
  _gmlFilter = aud;
  document.querySelectorAll('.gml-chip').forEach(c => c.classList.toggle('active', c.dataset.aud === aud));
  gmlRender();
}}

function gmlRender() {{
  const list = document.getElementById('gm-library-list');
  const items = _gmlFilter === 'all' ? _gmlItems : _gmlItems.filter(p => p.audience === _gmlFilter);
  if (!items.length) {{
    list.innerHTML = '<div id="gm-library-empty" style="display:block">' +
      (_gmlItems.length ? 'No prompts for this audience.' : 'No saved prompts yet — save one from the Chat tab.') +
      '</div>';
    return;
  }}
  list.innerHTML = items.map(p => {{
    const aud   = p.audience || 'portfolio_owner';
    const label = _GML_AUDIENCE_LABELS[aud] || aud;
    const prev  = (p.text || '').replace(/</g,'&lt;').replace(/>/g,'&gt;').slice(0, 240);
    return `<div class="gml-card" data-id="${{p.id}}">
      <div class="gml-card-header">
        <div class="gml-card-title">${{(p.title||'Untitled').replace(/</g,'&lt;')}}</div>
        <span class="gml-audience-badge ${{aud}}">${{label}}</span>
        <div class="gml-card-actions">
          <button class="gml-btn gml-btn-run" onclick="gmlRun('${{p.id}}')">▶ Run</button>
          <button class="gml-btn gml-btn-edit" onclick="gmlSendToChat('${{p.id}}')" title="Load into chat">✏️</button>
          <button class="gml-btn gml-btn-del" onclick="gmlDelete('${{p.id}}')" title="Delete">🗑</button>
        </div>
      </div>
      <div class="gml-preview">${{prev}}</div>
    </div>`;
  }}).join('');
}}

async function gmlSave(title, audience, text) {{
  if (!_DB_CRED) return;
  const id = 'pl_' + Date.now().toString(36) + Math.random().toString(36).slice(2,6);
  const row = {{ id, title, text, audience, created_at: new Date().toISOString() }};
  await fetch('/api/tiles/' + _TILE_ID + '/db/write', {{
    method: 'POST',
    headers: {{ 'Authorization': 'Basic ' + _DB_CRED, 'Content-Type': 'application/json' }},
    body: JSON.stringify({{
      table: 'prompt_library', mode: 'upsert',
      key_column: 'id',
      columns: ['id','title','text','audience','created_at'],
      rows: [row],
    }}),
  }});
  _gmlItems.unshift(row);
  _gmlLoaded = true;
}}

async function gmlDelete(id) {{
  if (!confirm('Delete this saved prompt?')) return;
  if (!_DB_CRED) return;
  await fetch('/api/tiles/' + _TILE_ID + '/db/write', {{
    method: 'POST',
    headers: {{ 'Authorization': 'Basic ' + _DB_CRED, 'Content-Type': 'application/json' }},
    body: JSON.stringify({{
      table: 'prompt_library', mode: 'delete',
      key_column: 'id', rows: [{{ id }}],
    }}),
  }});
  _gmlItems = _gmlItems.filter(p => p.id !== id);
  gmlRender();
}}

function gmlSendToChat(id) {{
  const p = _gmlItems.find(x => x.id === id);
  if (!p) return;
  const inp = document.getElementById('gm-chat-input');
  inp.value = p.text;
  inp.style.height = 'auto';
  inp.style.height = Math.max(68, Math.min(inp.scrollHeight, Math.floor(window.innerHeight * 0.4))) + 'px';
  switchGmTab('chat');
  setTimeout(() => {{ inp.focus(); inp.setSelectionRange(inp.value.length, inp.value.length); }}, 80);
}}

async function gmlRun(id) {{
  const p = _gmlItems.find(x => x.id === id);
  if (!p) return;
  switchGmTab('overview');
  const body    = document.getElementById('gm-body');
  const tsEl    = document.getElementById('gm-generated-at');
  const modelEl = document.getElementById('gm-model-used');
  body.innerHTML = _gmSkeleton('Running saved prompt…');
  tsEl.textContent = '';
  modelEl.textContent = '';
  document.getElementById('gm-regen-btn').disabled = true;
  const data = _gmCurrentData.length > 0 ? _gmCurrentData : RAW.slice();
  const fullPrompt = p.text + '\\n\\n---\\n\\n' + buildGMPrompt(data);
  try {{
    let result;
    for (const tier of ['powerful','balanced','fast']) {{
      try {{
        const resp = await fetch('/api/proxy/llm/' + window.__UPLOAD_ID__, {{
          method: 'POST',
          headers: {{ 'Content-Type': 'application/json', 'X-Proxy-Token': window.__PROXY_TOKEN__ }},
          body: JSON.stringify({{ system: _GM_CHAT_SYS, prompt: fullPrompt, tier, maxTokens: 2000 }}),
        }});
        if (!resp.ok) {{ const t = await resp.text().catch(()=>''); throw new Error('HTTP ' + resp.status + (t ? ': ' + t.slice(0,100) : '')); }}
        result = await resp.json();
        break;
      }} catch(e) {{
        if (e.message.startsWith('HTTP 503') && tier !== 'fast') continue;
        throw e;
      }}
    }}
    body.innerHTML = _gmRenderMarkdown(result.response || '');
    modelEl.textContent = 'Model: ' + (result.model_used || result.tier || '');
    tsEl.textContent = 'Generated: ' + new Date().toLocaleString();
  }} catch(e) {{
    body.innerHTML = `<p style="color:var(--red)">⚠️ Failed: ${{e.message}}</p>`;
  }} finally {{
    document.getElementById('gm-regen-btn').disabled = false;
  }}
}}

// ── Save Prompt dialog ────────────────────────────────────────────────────────
function openSavePrompt() {{
  const txt = (document.getElementById('gm-chat-input').value || '').trim();
  if (!txt) {{ alert('Type a prompt in the chat box first.'); return; }}
  document.getElementById('save-prompt-title-inp').value = '';
  document.getElementById('save-prompt-modal').classList.add('open');
  setTimeout(() => document.getElementById('save-prompt-title-inp').focus(), 60);
}}

function closeSavePrompt() {{
  document.getElementById('save-prompt-modal').classList.remove('open');
}}

async function confirmSavePrompt() {{
  const title = document.getElementById('save-prompt-title-inp').value.trim();
  if (!title) {{ document.getElementById('save-prompt-title-inp').focus(); return; }}
  const audience = document.getElementById('save-prompt-audience').value;
  const text     = document.getElementById('gm-chat-input').value.trim();
  const btn      = document.getElementById('save-prompt-ok');
  btn.disabled   = true;
  btn.textContent = 'Saving…';
  try {{
    await gmlSave(title, audience, text);
    closeSavePrompt();
    // brief confirmation
    const saveBtn = document.getElementById('gm-chat-save');
    saveBtn.textContent = '✓ Saved';
    saveBtn.style.color = 'var(--green)';
    setTimeout(() => {{ saveBtn.textContent = '💾'; saveBtn.style.color = ''; }}, 2200);
  }} catch(e) {{
    alert('Save failed: ' + e.message);
  }} finally {{
    btn.disabled = false;
    btn.textContent = 'Save';
  }}
}}

function _gmRenderPrompts() {{
  const panel = document.getElementById('gm-prompts-panel');
  if (panel.dataset.rendered) return;
  panel.dataset.rendered = '1';
  panel.innerHTML = _HELP_PROMPTS.map((p, i) => `
    <div class="gmp-card">
      <div class="gmp-card-header">
        <div>
          <div class="gmp-card-title">${{p.title}}</div>
          <div class="gmp-card-audience">${{p.audience}}</div>
        </div>
        <div style="display:flex;gap:6px;align-items:center">
          <button class="gmp-chat-btn" onclick="_gmPromptToChat(${{i}})" title="Copy into Chat for editing">✏️ Edit &amp; Send</button>
          <button class="gmp-run-btn" id="gmp-run-${{i}}" onclick="_gmRunPrompt(${{i}})">▶ Run</button>
        </div>
      </div>
      <div class="gmp-preview">${{p.text.replace(/</g,'&lt;').replace(/>/g,'&gt;').slice(0,220)}}</div>
    </div>`).join('');
}}

function _gmPromptToChat(i) {{
  const inp = document.getElementById('gm-chat-input');
  inp.value = _HELP_PROMPTS[i].text;
  inp.style.height = 'auto';
  inp.style.height = Math.max(68, Math.min(inp.scrollHeight, Math.floor(window.innerHeight * 0.4))) + 'px';
  switchGmTab('chat');
  setTimeout(() => {{ inp.focus(); inp.setSelectionRange(inp.value.length, inp.value.length); }}, 80);
}}

async function _gmRunPrompt(i) {{
  const btn = document.getElementById('gmp-run-' + i);
  btn.disabled = true;
  btn.textContent = '…';
  switchGmTab('overview');
  const body   = document.getElementById('gm-body');
  const tsEl   = document.getElementById('gm-generated-at');
  const modelEl = document.getElementById('gm-model-used');
  body.innerHTML = _gmSkeleton('Running prompt…');
  tsEl.textContent = '';
  modelEl.textContent = '';
  document.getElementById('gm-regen-btn').disabled = true;

  const data = _gmCurrentData.length > 0 ? _gmCurrentData : RAW.slice();
  const useSlack = document.getElementById('gm-prompts-slack-chk')?.checked;
  let slackCtx = '';
  if (useSlack) {{
    const lbl = document.getElementById('gm-spinner-label');
    slackCtx = await _buildSlackContext(data, msg => {{ if (lbl) lbl.textContent = msg; }});
    if (lbl) lbl.textContent = 'Running prompt…';
  }}
  const fullPrompt = _HELP_PROMPTS[i].text + '\\n\\n---\\n\\n' + buildGMPrompt(data) + slackCtx;
  try {{
    let result;
    for (const tier of ['powerful', 'balanced', 'fast']) {{
      try {{
        const lbl = document.getElementById('gm-spinner-label');
        if (lbl) lbl.textContent = tier === 'powerful' ? 'Running prompt…' : 'Retrying with ' + tier + ' model…';
        const resp = await fetch('/api/proxy/llm/' + window.__UPLOAD_ID__, {{
          method: 'POST',
          headers: {{ 'Content-Type': 'application/json', 'X-Proxy-Token': window.__PROXY_TOKEN__ }},
          body: JSON.stringify({{ system: _GM_CHAT_SYS, prompt: fullPrompt, tier, maxTokens: 2000 }}),
        }});
        if (!resp.ok) {{
          const errBody = await resp.text().catch(() => '');
          throw new Error('HTTP ' + resp.status + (errBody ? ': ' + errBody.slice(0,120) : ''));
        }}
        result = await resp.json();
        break;
      }} catch(e) {{
        if (e.message.startsWith('HTTP 503') && tier !== 'fast') continue;
        throw e;
      }}
    }}
    const text = result.response || '';
    body.innerHTML = _gmRenderMarkdown(text);
    modelEl.textContent = 'Model: ' + (result.model_used || result.tier || '');
    tsEl.textContent = 'Generated: ' + new Date().toLocaleString();
  }} catch(e) {{
    const is401 = e.message.startsWith('HTTP 401');
    body.innerHTML = is401
      ? `<p style="color:var(--red)">⚠️ Session expired — please <strong>reload the page</strong> and try again.</p>`
      : `<p style="color:var(--red)">⚠️ Failed: ${{e.message}}</p>`;
    tsEl.textContent = '';
  }} finally {{
    document.getElementById('gm-regen-btn').disabled = false;
    if (btn) {{ btn.disabled = false; btn.textContent = '▶ Run'; }}
  }}
}}

// ── Slack integration ──────────────────────────────────────────────────────────
async function _slackProxyFetch(url, options) {{
  const resp = await fetch('/api/proxy/fetch/' + window.__UPLOAD_ID__ + '/' + url, {{
    ...options,
    headers: {{ ...(options && options.headers), 'X-Proxy-Token': window.__PROXY_TOKEN__ }},
  }});
  if (!resp.ok) {{
    const err = await resp.json().catch(() => ({{}}));
    throw new Error(err.error || 'Slack HTTP ' + resp.status);
  }}
  return resp.json();
}}

async function _slackFindChannel(name) {{
  // Normalise: lowercase, replace non-alphanum with -, collapse multiples, trim
  const slug = name.toLowerCase().replace(/[^a-z0-9]+/g, '-').replace(/^-|-$/g, '').slice(0, 40);
  // Try exact name first, then progressive prefix (80-char Slack channel name limit)
  const candidates = [slug, slug.slice(0, 21), slug.split('-').slice(0, 3).join('-')];
  for (const q of [...new Set(candidates)]) {{
    if (!q) continue;
    try {{
      const r = await _slackProxyFetch(
        'https://slack.com/api/conversations.list?exclude_archived=true&limit=200&types=public_channel,private_channel',
        {{}}
      );
      if (!r.ok) continue;
      const match = (r.channels || []).find(c =>
        c.name === q || c.name.includes(q) || q.includes(c.name.slice(0,10))
      );
      if (match) return match;
    }} catch(e) {{ console.warn('[Slack] find error:', e); return null; }}
  }}
  return null;
}}

async function _slackChannelSummary(channelId, channelName) {{
  const oldest = Math.floor((Date.now() - 7 * 24 * 3600 * 1000) / 1000);
  try {{
    const r = await _slackProxyFetch(
      `https://slack.com/api/conversations.history?channel=${{channelId}}&oldest=${{oldest}}&limit=100`,
      {{}}
    );
    if (!r.ok || !r.messages || r.messages.length === 0) return null;
    const msgs = r.messages
      .filter(m => m.type === 'message' && !m.subtype)
      .slice(0, 40)
      .map(m => (m.text || '').slice(0, 300))
      .join('\\n');
    return `#${{channelName}} (last 7 days, ${{r.messages.length}} msgs):\\n${{msgs}}`;
  }} catch(e) {{ console.warn('[Slack] history error:', e); return null; }}
}}

async function _buildSlackContext(data, onStatus) {{
  const projects = data.slice(0, 20); // cap at 20 to stay within rate limit
  const results = [];
  for (const p of projects) {{
    const label = (p.acct || p.name || '').slice(0, 40);
    if (onStatus) onStatus(`Searching Slack for ${{label}}…`);
    const ch = await _slackFindChannel(p.acct || p.name || '');
    if (!ch) continue;
    const summary = await _slackChannelSummary(ch.id, ch.name);
    if (summary) results.push(summary);
  }}
  return results.length > 0
    ? '\\n\\n--- Slack Channel Activity (last 7 days) ---\\n' + results.join('\\n\\n')
    : '\\n\\n(No matching Slack channels found for the visible projects.)';
}}

// ── GM Chat ────────────────────────────────────────────────────────────────────
let _gmChatHistory = [];

function gmChatClear() {{
  _gmChatHistory = [];
  document.getElementById('gm-chat-thread').innerHTML = '';
}}

function _gmChatBubble(role, html) {{
  const thread = document.getElementById('gm-chat-thread');
  const div = document.createElement('div');
  div.className = 'chat-msg ' + role;
  const avatar = role === 'user' ? 'You' : 'AI';
  div.innerHTML = `<div class="chat-avatar">${{avatar}}</div><div class="chat-bubble">${{html}}</div>`;
  thread.appendChild(div);
  thread.scrollTop = thread.scrollHeight;
  return div;
}}

function _gmThinking() {{
  const thread = document.getElementById('gm-chat-thread');
  const div = document.createElement('div');
  div.className = 'chat-msg ai';
  div.id = 'gm-thinking';
  div.innerHTML = `<div class="chat-avatar">AI</div>
    <div class="chat-bubble chat-thinking"><span></span><span></span><span></span></div>`;
  thread.appendChild(div);
  thread.scrollTop = thread.scrollHeight;
}}

async function gmChatSend() {{
  const input = document.getElementById('gm-chat-input');
  const sendBtn = document.getElementById('gm-chat-send');
  const text = input.value.trim();
  if (!text) return;

  if (!window.__PROXY_TOKEN__ || window.__PROXY_TOKEN__ === 'undefined') {{
    _gmChatBubble('ai', '<p style="color:var(--red)">⚠️ Session token not available — please reload the page.</p>');
    return;
  }}

  input.value = '';
  input.style.height = 'auto';
  _gmChatBubble('user', text.replace(/</g,'&lt;').replace(/>/g,'&gt;'));
  _gmThinking();
  sendBtn.disabled = true;

  // Build context: portfolio snapshot + optional Slack + conversation history
  const data = _gmCurrentData.length > 0 ? _gmCurrentData : RAW.slice();
  const portfolioCtx = buildGMPrompt(data);
  const useSlack = document.getElementById('gm-chat-slack-chk')?.checked;
  let slackCtx = '';
  if (useSlack) {{
    let statusEl = document.getElementById('gm-chat-slack-status');
    if (!statusEl) {{
      statusEl = document.createElement('div');
      statusEl.id = 'gm-chat-slack-status';
      statusEl.className = 'gm-slack-status';
      document.getElementById('gm-chat-slack-bar').after(statusEl);
    }}
    slackCtx = await _buildSlackContext(data, msg => {{ statusEl.textContent = msg; }});
    statusEl.remove();
  }}
  const history = _gmChatHistory.map(m => m.role + ': ' + m.content).join('\\n');
  const fullPrompt = (history ? history + '\\n' : '') +
    'user: ' + text + '\\n\\n' +
    '--- Current Portfolio Data ---\\n' + portfolioCtx + slackCtx;

  _gmChatHistory.push({{ role: 'user', content: text }});

  try {{
    let result;
    for (const tier of ['powerful', 'balanced', 'fast']) {{
      try {{
        const resp = await fetch('/api/proxy/llm/' + window.__UPLOAD_ID__, {{
          method: 'POST',
          headers: {{ 'Content-Type': 'application/json', 'X-Proxy-Token': window.__PROXY_TOKEN__ }},
          body: JSON.stringify({{ system: _GM_CHAT_SYS, prompt: fullPrompt, tier, maxTokens: 2000 }}),
        }});
        if (!resp.ok) {{
          const errBody = await resp.text().catch(() => '');
          throw new Error('HTTP ' + resp.status + (errBody ? ': ' + errBody.slice(0,120) : ''));
        }}
        result = await resp.json();
        break;
      }} catch(e) {{
        if (e.message.startsWith('HTTP 503') && tier !== 'fast') continue;
        throw e;
      }}
    }}
    const reply = result.response || '';
    _gmChatHistory.push({{ role: 'assistant', content: reply }});
    document.getElementById('gm-thinking')?.remove();
    _gmChatBubble('ai', _gmRenderMarkdown(reply));
  }} catch(e) {{
    document.getElementById('gm-thinking')?.remove();
    const is401 = e.message.startsWith('HTTP 401');
    _gmChatBubble('ai', is401
      ? `<p style="color:var(--red)">⚠️ Session expired — please reload the page.</p>`
      : `<p style="color:var(--red)">⚠️ ${{e.message}}</p>`);
    _gmChatHistory.pop();
  }} finally {{
    sendBtn.disabled = false;
    input.focus();
  }}
}}

// Auto-grow textarea + Enter to send (Shift+Enter for newline)
document.addEventListener('DOMContentLoaded', () => {{
  const inp = document.getElementById('gm-chat-input');
  if (!inp) return;
  function _growInp() {{
    inp.style.height = 'auto';
    const minPx = 68;
    const maxPx = Math.floor(window.innerHeight * 0.4);
    inp.style.height = Math.max(minPx, Math.min(inp.scrollHeight, maxPx)) + 'px';
  }}
  inp.addEventListener('input', _growInp);
  inp.addEventListener('keydown', e => {{
    if (e.key === 'Enter' && !e.shiftKey) {{ e.preventDefault(); gmChatSend(); }}
  }});
}});

document.addEventListener('keydown', e => {{
  if (e.key === 'Escape') {{ closeGMOverview(); closeHelp(); _hidePulse(true); _hideFar(true); _hideP2g(true); _hidePipe(true); _hideRes(true); _hideSummary(true); }}
}});


// ── Help modal ─────────────────────────────────────────────────────────────────
const _HELP_PROMPTS = [
  {{
    title: 'License Sales EVP Briefing — High-Watch Portfolio',
    audience: 'Audience: License Sales EVP',
    text: `Act as our Professional Services Business Analyst providing an executive briefing to our License Sales EVP. Summarize our High-Watch delivery portfolio ($260.42M across 11 key accounts), specifically highlighting critical upcoming customer Go-Lives that will secure/unlock software consumption, active 'Watermelon' risks to be aware of, and expansion/whitespace opportunities where delivery success can directly grease the wheel for upcoming license renewals and Net New ARR.`,
  }},
  {{
    title: 'GM Operational Briefing — Full Executive Overview',
    audience: 'Audience: Professional Services GM',
    text: `Act as a Senior Business Analyst presenting a high-level operational briefing to the Professional Services General Manager (GM).\\n\\nAnalyze the provided portfolio data and construct an executive overview using the exact markdown formatting and sections below. Focus heavily on financial risk, resource efficiencies, data hygiene gaps, and upcoming critical milestones.\\n\\nKeep these core structural definitions and accounting mechanics in mind during your financial analysis:\\n1. SWE designates internal investment money (where bookings/billings are zero, meaning the goal is resource efficiency, timeline baselines, and margin preservation).\\n2. FAR stands for Forecasted Amount Remaining—a critical financial metric tracking expected revenue remaining in contract funds for T&M projects (FAR = Bookings - Billings - Scheduled Work - Pending Resource Requests).\\n   - A Negative FAR (-) variance indicates the project is actively over-staffed or over budget, requiring immediate onshore cost-containment, hour reductions, or a Change Order.\\n   - A Positive FAR (+) variance indicates under-staffing or under-utilization, meaning funded scope is stalling and leaving money on the table.\\n   - Thresholds (RAG): A FAR variance > 15% of total T&M SOW bookings is flagged as Red; 5%-14% is Yellow; < 5% is Green.\\n\\n---\\n\\n### [Portfolio Name / Region] Executive Briefing\\n**To:** Professional Services General Manager\\n**From:** Portfolio Business Analyst\\n**Date:** [Insert Date]\\n**Portfolio Scope:** [X] Active Projects | $[X]M Total Bookings\\n\\n---\\n\\n### 📊 Performance at a Glance\\nProvide a high-level bulleted summary of macro metrics from the data scorecard:\\n*   Total Billings vs. Remaining Backlog (Calculate the portfolio's financial burn percentage).\\n*   Margin Optimization (Compare Weighted Bid Margin % against current Weighted Margin at Close % to show the exact Margin Delta optimization).\\n*   Client Satisfaction (Average CSAT score, but include a data hygiene caveat highlighting the exact percentage of active projects that actually have a survey on file vs. those operating blind with missing/lapsed data).\\n*   Data Hygiene (Average Data Quality score).\\n\\n---\\n\\n### 🚦 Portfolio Health & Working Capital Risks\\nAnalyze the data and summarize the hidden risks using these specific categories:\\n*   The Overdue Cash Bottleneck: Highlight total overdue invoices and name the top enterprise account culprits driving the cash freeze.\\n*   The "Watermelon" Transparency Gap: Count how many projects are marked "Green" overall by PMs but are classified as "Watermelon Green" due to internal red flags (e.g., budget overruns, negative FAR, 0% GDC offshore utilization, or past end dates).\\n*   Financial At Risk (FAR): Summarize total dollar exposure trapped across the pipeline, breaking down where projects are actively leaking funds via overburns (-) vs. where revenue velocity is stalling via under-utilization (+).\\n\\n---\\n\\n### 🔍 Account Ledger & High Watch Status Summary\\nCreate a comprehensive markdown table explicitly tracking all High-Watch customers (denoted by ⚑ HW or similar indicators) and any projects carrying a 🔴 RED status.\\n\\nFor each row, include:\\n1. Status (Use Emoji: 🔴 RED, 🟡 YELLOW, 🍉 WATERMELON)\\n2. Account & Project Name (Prefix with ⚑ HW if it is a High Watch customer)\\n3. Leadership Team (List the PM and Account Owner)\\n4. Key Financials (Bookings, Billings, Overdue Invoices, and FAR variance)\\n5. Executive Summary, Key Business Outcomes, & Top Risks (Use your natural language processing to synthesize the raw data columns and the complete summary column into a clean, professional narrative structured exactly like this):\\n   - **Executive Summary:** A concise overview of the project's contract type, scope parameters, and lifecycle stage.\\n   - **Key Business Outcome:** The core business or technical milestone the client is trying to achieve.\\n   - **Top Risks & Actions:** A real-time breakdown of the exact operational issues occurring (e.g., product deployability, security gaps, client funding cuts, 0% offshore skews, or approaching deadlines) and the precise tactical steps being taken to resolve them.\\n\\n---\\n\\n### 🛠️ GM Immediate Intervention Actions\\nProvide exactly 5 clear, highly actionable executive directives for the GM to execute immediately based on the risks identified above. Include specific project names, target dates, resource reallocation goals, or executive alignment escalations needed to protect portfolio health.\\n\\n---\\n\\n[PASTE YOUR RAW PORTFOLIO DATA, LEDGER TEXT, OR SCREENSHOT DATA HERE]`,
  }},
  {{
    title: 'Portfolio Owner Battle Plan — Tactical Action Plan',
    audience: 'Audience: Portfolio Owner (self-service)',
    text: `Act as a Senior Delivery and Operational Excellence Business Analyst. I am the Portfolio Owner for this specific segment of accounts.\\n\\nAnalyze my portfolio data and construct a high-impact, tactical action plan using the exact markdown structure below. Focus heavily on identifying execution leaks, upcoming milestone freezes, data hygiene gaps, and immediate operational triage.\\n\\nKeep these two critical structural definitions in mind during your financial analysis:\\n1. SWE designates internal investment money (where bookings/billings are zero, meaning the goal is resource efficiency and margin preservation).\\n2. FAR stands for Forecasted Amount Remaining—a critical financial metric tracking expected revenue remaining in contract funds. A negative FAR indicates an overburn risk, while a high positive FAR on a stalling track indicates unconsumed pipeline funds.\\n\\n---\\n\\n### 🎯 Portfolio Owner Battle Plan: [Insert Portfolio / Region Name]\\n**Portfolio Owner:** [Insert Your Name]\\n**Data Snapshot:** As of [Insert Date] | [X] High-Watch Engagements\\n\\n---\\n\\n### 📊 Tactical Scorecard\\nExtract these precise health metrics from my portfolio view:\\n*   **Backlog Velocity:** Total Billings vs. Remaining Backlog.\\n*   **Margin Realization:** Compare our Weighted Bid Margin % against our current Weighted Margin at Close % to see if execution is bleeding or building profit.\\n*   **Leakage Metrics:** Aggregate the total dollar value of my portfolio's Overdue Invoices and check for outlying project-level FAR (Forecasted Amount Remaining) variances.\\n\\n---\\n\\n### 🚨 Top 3 Critical Delivery Leaks & Deficit Clearances\\nIdentify the 3 highest-risk projects in this specific list that require my personal operational intervention this week. For each risk, explicitly structure it as follows:\\n\\n1. **[Account Name - Project Name] | [Current Status Emoji & Color]**\\n   * **The Exposure:** (State the precise breakdown of the financial or timeline metric failing: e.g., $X overdue invoices, negative or unconsumed FAR, 0% GDC leverage, or an imminent milestone/Go-Live date).\\n   * **The Ground Reality:** (Analyze the project summary or leadership notes to explain *why* this is failing—e.g., client funding cuts, a temporary strategic pause, internal investment bounds, or missing executive governance).\\n   * **My Direct Action:** (Provide a clear tactical next step for me to execute as the Portfolio Owner to protect the account).\\n\\n---\\n\\n### 🍉 The "Watermelon" Transparency Audit\\nReview the projects reported as Green. Identify which ones are "Watermelon Green" (surface green, internal red flags). List them and detail the hidden structural or resourcing metrics (e.g., low GDC offshore mix, missing baselines, slipped budgets, or past end dates) that I need to force my Project Managers to correct.\\n\\n---\\n\\n### 📝 Strategic Ground Opportunities & Whitespace\\nScan the qualitative updates (Leadership Notes/Summaries) to find where a project pause, a pivot, or a major delivery success (including successful internal SWE proof-of-concepts) has created an opening to expand our services footprint, convert internal investments to paid tracks, or introduce new capabilities (e.g., Agentforce, Data Cloud, MuleSoft). Identify the key client stakeholder to target.\\n\\n---\\n\\n[PASTE YOUR RAW ACCOUNT LEDGER, STATUS SUMMARIES, OR SCORECARD DATA HERE]`,
  }},
];

const _HELP_GUIDE = `
<h3>Navigating the Dashboard</h3>
<ul>
  <li><strong>Region Tabs</strong> — Switch between ACC (all), AMER TMT, and AMER CBS views at the top.</li>
  <li><strong>Search</strong> — Type any account name, project name, or keyword to filter rows instantly.</li>
  <li><strong>Status / Type / Tier / PO / Rules filters</strong> — Use the dropdowns to multi-select any combination. Active filters highlight in blue.</li>
  <li><strong>Type filter</strong> — Filter by ⚑ High Watch, 🔧 SWE, 🏷 ARI, or 🗓 Not Started (projects whose start date is in the future).</li>
  <li><strong>High Watch / SWE toggles</strong> — Quick-filter to High Watch projects or SWE-only tracks.</li>
  <li><strong>Group By</strong> — Organize rows by Status, Tier, Portfolio Owner, or Region. Use ⊟/⊞ to collapse or expand groups.</li>
  <li><strong>✕ Clear Filters</strong> — Resets all filters and search in one click. Group-by and collapse/expand state are preserved.</li>
</ul>
<h3>Scorecard KPIs</h3>
<ul>
  <li>KPI tiles update live to reflect whatever filters are active — the numbers always match what you see in the table.</li>
  <li>Click any KPI tile to drill into that metric (where applicable).</li>
</ul>
<h3>Editing Tier &amp; Portfolio Owner</h3>
<ul>
  <li>Click the <strong>✎ pencil icon</strong> on any project row to open the Edit Assignment dialog.</li>
  <li>Change the Tier and/or Portfolio Owner, then click <strong>Save</strong>. Changes persist to the server immediately.</li>
</ul>
<h3>Project Detail</h3>
<ul>
  <li>Click anywhere on a project row to expand the detail panel — shows financials, GDC breakdown, rules, and the latest pulse note.</li>
  <li>The <strong>Resources</strong> column shows GDC India %. <strong>Hover</strong> to preview, <strong>click</strong> to pin the resource detail dialog (sorted by name; click column headers to re-sort). Press <strong>Escape</strong> or click ✕ to close.</li>
</ul>
<h3>Baselines</h3>
<ul>
  <li>The Baselines column shows a color-coded circle for each dimension: <span style="display:inline-block;width:9px;height:9px;border-radius:50%;background:#c0392b;vertical-align:middle"></span> Red &nbsp;<span style="display:inline-block;width:9px;height:9px;border-radius:50%;background:#e67e22;vertical-align:middle"></span> Yellow &nbsp;<span style="display:inline-block;width:9px;height:9px;border-radius:50%;background:#27ae60;vertical-align:middle"></span> Green.</li>
  <li>Dimensions: Scope · Schedule · Budget · Resource · Customer.</li>
</ul>
<h3>FAR &amp; Path to Green</h3>
<ul>
  <li>In the Financials column, <strong>hover or click the FAR value</strong> (shown with ⓘ) to see the FAR reason codes.</li>
  <li>When a project has Path to Green notes entered in Salesforce, a <strong style="color:#27ae60">P2G ⓘ</strong> link appears next to the FAR value. <strong>Hover or click</strong> to view Path to Green notes by dimension (Budget, Scope, Schedule, Resource, Customer).</li>
</ul>
<h3>Pipeline</h3>
<ul>
  <li><strong>Hover or click</strong> the pipeline amount in the Financials column to open the pipeline popup for that account.</li>
  <li>Pipeline is sorted by close date (soonest first) by default. Click <strong>Amount</strong> or <strong>Close Date</strong> column headers to re-sort.</li>
  <li>Rules <strong>PIPE_CLOSE_THIS_M</strong> (red) and <strong>PIPE_CLOSE_THIS_Q</strong> (yellow) flag accounts with pipeline closing this month or SF quarter.</li>
</ul>
<h3>Pulse &amp; Popups</h3>
<ul>
  <li>All popups (Pulse, FAR, P2G, Pipeline, Resources) share the same interaction model: <strong>hover</strong> to preview, <strong>click</strong> to pin (stays open while you read), <strong>Escape</strong> or ✕ to force-close.</li>
</ul>
<h3>Business Overview (AI)</h3>
<ul>
  <li>Click <strong>📊 Business Overview</strong> to generate an AI-written executive summary of the currently filtered portfolio.</li>
  <li>The summary reflects whatever region, filters, or search are active — narrow to a PO or tier first for a focused briefing.</li>
  <li>Use <strong>↻ Regenerate</strong> to get a fresh response. Results are cached per session so switching regions won't re-run unnecessarily.</li>
  <li>See the <strong>AI Prompts</strong> tab for copy-ready prompts to use with any external AI tool.</li>
</ul>
<h3>Portfolio Assurance Rules</h3>
<p style="font-size:12px;color:var(--subtext);margin:0 0 8px">Rules fire automatically based on Salesforce data. Each project row shows its active rule codes in the Portfolio Assurance column.</p>
<table style="width:100%;border-collapse:collapse;font-size:12px">
  <thead><tr style="background:var(--navy);color:#fff">
    <th style="padding:6px 10px;text-align:left;white-space:nowrap">Rule Code</th>
    <th style="padding:6px 10px;text-align:left">Category</th>
    <th style="padding:6px 10px;text-align:left">Description &amp; Trigger</th>
  </tr></thead>
  <tbody>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">MARGIN_RED</td><td style="padding:5px 10px;color:var(--subtext)">Margin</td><td style="padding:5px 10px">Close Margin critically below Bid Margin. Triggers when margin delta &lt; −5%.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">MARGIN_YELLOW</td><td style="padding:5px 10px;color:var(--subtext)">Margin</td><td style="padding:5px 10px">Close Margin moderately below Bid Margin. Triggers when delta is between 0% and −5%.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">FAR_RED_NEG</td><td style="padding:5px 10px;color:var(--subtext)">FAR</td><td style="padding:5px 10px">FAR is negative — project is forecasted to overrun budget.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">FAR_RED_UNDERUTIL</td><td style="padding:5px 10px;color:var(--subtext)">FAR</td><td style="padding:5px 10px">FAR &gt;50% remaining with &lt;30% schedule consumed — severe underutilisation risk.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">FAR_YELLOW</td><td style="padding:5px 10px;color:var(--subtext)">FAR</td><td style="padding:5px 10px">Moderate FAR utilisation mismatch vs schedule progress.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">RR_RISK</td><td style="padding:5px 10px;color:var(--subtext)">Resource</td><td style="padding:5px 10px">Open Resource Requests with pending revenue — staffing gap risk.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">GDC_LOW</td><td style="padding:5px 10px;color:var(--subtext)">Resource</td><td style="padding:5px 10px">GDC India assigned share ≤65% — more than ⅓ of the team is non-GDC.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">OVERDUE_INV</td><td style="padding:5px 10px;color:var(--subtext)">Invoice</td><td style="padding:5px 10px">Overdue invoices outstanding (&gt;30 days past due). Cash collection risk.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">SWE_BURNING_HOT</td><td style="padding:5px 10px;color:var(--subtext)">SWE Burn</td><td style="padding:5px 10px">SWE burn rate alert — hours/spend tracking off-plan vs schedule.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">NO_PULSE</td><td style="padding:5px 10px;color:var(--subtext)">Governance</td><td style="padding:5px 10px">No pulse submitted — project status unknown. Governance violation.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">NO_STEERCO</td><td style="padding:5px 10px;color:var(--subtext)">Governance</td><td style="padding:5px 10px">SteerCo date required for projects ≥$750K. If exempt (SEH/Advisory) set date to 01/01/2100.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">MISSING_PTG</td><td style="padding:5px 10px;color:var(--subtext)">Governance</td><td style="padding:5px 10px">Missing Path to Green for a Red/Yellow baseline dimension.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">END_DATE_PAST</td><td style="padding:5px 10px;color:var(--subtext)">End Date</td><td style="padding:5px 10px">Project end date has passed — may need extension or closure.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">END_DATE_30</td><td style="padding:5px 10px;color:var(--subtext)">End Date</td><td style="padding:5px 10px">End date within 30 days — urgent renewal/extension decision needed.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">END_DATE_60</td><td style="padding:5px 10px;color:var(--subtext)">End Date</td><td style="padding:5px 10px">End date in 31–60 days — renewal/extension decision needed now.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">END_DATE_90</td><td style="padding:5px 10px;color:var(--subtext)">End Date</td><td style="padding:5px 10px">End date in 61–90 days — begin renewal planning.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">END_DATE_120</td><td style="padding:5px 10px;color:var(--subtext)">End Date</td><td style="padding:5px 10px">End date in 91–120 days — flag for upcoming renewal.</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:var(--red);white-space:nowrap">PIPE_CLOSE_THIS_M</td><td style="padding:5px 10px;color:var(--subtext)">Pipeline</td><td style="padding:5px 10px">Pipeline opportunity closing this calendar month — action needed to book.</td></tr>
    <tr style="border-bottom:1px solid var(--border);background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">PIPE_CLOSE_THIS_Q</td><td style="padding:5px 10px;color:var(--subtext)">Pipeline</td><td style="padding:5px 10px">Pipeline closing this SF fiscal quarter (Q ends Jan/Apr/Jul/Oct).</td></tr>
    <tr style="border-bottom:1px solid var(--border)"><td style="padding:5px 10px;font-weight:700;color:#b45309;white-space:nowrap">CSAT_OVERDUE</td><td style="padding:5px 10px;color:var(--subtext)">CSAT</td><td style="padding:5px 10px">No CSAT survey sent. Triggers when bookings &gt;$150K, active &gt;90 days, stage = In Progress.</td></tr>
    <tr style="background:var(--bg)"><td style="padding:5px 10px;font-weight:700;color:var(--subtext);white-space:nowrap">CSAT_EXEMPT</td><td style="padding:5px 10px;color:var(--subtext)">CSAT</td><td style="padding:5px 10px">Project is marked Do Not Survey in Salesforce (informational, not a violation).</td></tr>
  </tbody>
</table>
<h3>Need Help?</h3>
<ul>
  <li>Questions, bugs, or feature requests? Send a Slack message to <strong>@Alex Penkrat</strong>.</li>
</ul>
`;

function openHelp() {{
  document.getElementById('help-modal').classList.add('open');
  switchHelpTab('guide');
}}

function closeHelp() {{
  document.getElementById('help-modal').classList.remove('open');
}}

function switchHelpTab(tab) {{
  document.getElementById('help-tab-guide').classList.toggle('active', tab === 'guide');
  document.getElementById('help-tab-prompts').classList.toggle('active', tab === 'prompts');
  const body = document.getElementById('help-body');
  if (tab === 'guide') {{
    body.innerHTML = _HELP_GUIDE;
  }} else {{
    body.innerHTML = _HELP_PROMPTS.map((p, i) => `
      <div class="prompt-card">
        <div class="prompt-card-header">
          <div>
            <div class="prompt-card-title">${{p.title}}</div>
            <div class="prompt-card-audience">${{p.audience}}</div>
          </div>
          <button class="prompt-copy-btn" id="copy-btn-${{i}}" onclick="_copyPrompt(${{i}})">Copy</button>
        </div>
        <div class="prompt-text">${{p.text.replace(/</g,'&lt;').replace(/>/g,'&gt;')}}</div>
      </div>`).join('');
  }}
}}

function _copyPrompt(i) {{
  navigator.clipboard.writeText(_HELP_PROMPTS[i].text).then(() => {{
    const btn = document.getElementById('copy-btn-' + i);
    btn.textContent = '✓ Copied!';
    btn.classList.add('copied');
    setTimeout(() => {{ btn.textContent = 'Copy'; btn.classList.remove('copied'); }}, 2000);
  }}).catch(() => {{
    const ta = document.createElement('textarea');
    ta.value = _HELP_PROMPTS[i].text;
    document.body.appendChild(ta);
    ta.select();
    document.execCommand('copy');
    document.body.removeChild(ta);
    const btn = document.getElementById('copy-btn-' + i);
    btn.textContent = '✓ Copied!';
    btn.classList.add('copied');
    setTimeout(() => {{ btn.textContent = 'Copy'; btn.classList.remove('copied'); }}, 2000);
  }});
}}

loadData();
</script>
</body>
</html>"""

    import json as _json_html

    # Write pipeline detail JSON
    pipe_by_acct = {}
    for opp in q5b_pipe:
        aid = opp.get('AccountId') or ''
        if not aid:
            continue
        row = {
            'id':           opp.get('Id', ''),
            'name':         opp.get('Name', ''),
            'stage':        opp.get('StageName', ''),
            'amount':       opp.get('Amount') or 0,
            'converted':    opp.get('convertedAmount') or opp.get('Amount') or 0,
            'currency':     opp.get('CurrencyIsoCode', 'USD'),
            'close':        (opp.get('CloseDate') or '')[:10],
            'created':      (opp.get('CreatedDate') or '')[:10],
            'owner':        opp.get('Owner.Name', '') or '',
            'forecast':     opp.get('ForecastCategory', ''),
            'mgr_forecast': opp.get('Manager_Forecast_Judgement__c', ''),
            'description':  opp.get('Description', '') or '',
            'bid_margin':   opp.get('Overall_Bid_Margin__c'),
            'push_count':   opp.get('PushCount') or 0,
        }
        pipe_by_acct.setdefault(aid, []).append(row)
    pipe_json_path = os.path.join(OUTPUT_DIR, f"acc_{REGION_SLUG.lower()}_pipeline.json")
    with open(pipe_json_path, 'w', encoding='utf-8') as fh:
        _json_html.dump({'generated': REPORT_DATE, 'region': REGION_LABEL, 'opps': pipe_by_acct}, fh, ensure_ascii=False)
    print(f"✅  Pipeline JSON saved: {pipe_json_path}")

    json_path = os.path.join(OUTPUT_DIR, f"acc_{REGION_SLUG.lower()}_data.json")
    with open(json_path, 'w', encoding='utf-8') as fh:
        _json_html.dump({'generated': REPORT_DATE, 'region': REGION_LABEL, 'rows': rows_data, 'pipe': pipe_by_acct}, fh, ensure_ascii=False)
    print(f"✅  JSON data saved: {json_path}")

    path = base_path + '.html'
    with open(path, 'w', encoding='utf-8') as fh:
        fh.write(html)
    print(f"✅  HTML saved: {path}")

    # Undated copy for Heroku / static hosting (always overwrites)
    latest_path = os.path.join(OUTPUT_DIR, f"{REGION_SLUG}_Audit_latest.html")
    with open(latest_path, 'w', encoding='utf-8') as fh:
        fh.write(html)
    print(f"✅  HTML (latest) saved: {latest_path}")

def write_json():
    import json as _json
    def fmt_m(v):
        if v is None: return ''
        return f'${v/1e6:.2f}M' if abs(v) >= 1_000_000 else f'${v:,.0f}'
    def fmt_pct(v):
        return f'{v:.1f}%' if v is not None else ''
    def status_label(r):
        if r['stage'] == 'On Hold': return 'On Hold'
        if r['is_watermelon']: return 'Watermelon'
        h = r['health'].lower()
        if h == 'red':    return 'Red'
        if h == 'yellow': return 'Yellow'
        if h == 'green':  return 'Green'
        return 'No Pulse'
    rows_data = []
    for r in results:
        dd = (r['close_margin'] - r['bid_margin']) if (r['bid_margin'] is not None and r['close_margin'] is not None) else None
        viol_codes = ', '.join(v[0] for v in r['violations'])
        bl = ' | '.join(f"{lbl}={v}" for lbl, v in baselines_list(r))
        team = []
        if r['pm']:        team.append(f"PM: {r['pm']}")
        if r.get('pm2'):   team.append(f"PM2: {r['pm2']}")
        if r.get('opp_owner'): team.append(f"AP: {r['opp_owner']}")
        if r.get('exec_sponsor'): team.append(f"ES: {r['exec_sponsor']}")
        if r.get('acct_owner'): team.append(f"AO: {r['acct_owner']}")
        if r.get('owner') and r.get('owner') != 'Unassigned': team.append(f"PO: {r['owner']}")
        po = r.get('owner','') or ''
        rows_data.append({
            'name':          r['name'],
            'acct':          r['acct'],
            'acct_id':       r.get('acct_id') or '',
            'status':        status_label(r),
            'tier':          r['tier'],
            'team':          ' · '.join(team),
            'po':            po,
            'pid':           r['pid'],
            'region':        r.get('region', REGION_LABEL),
            'url':           f"https://org62.lightning.force.com/lightning/r/pse__Proj__c/{r['pid']}/view",
            'rev_treat':     r.get('rev_treat') or '',
            'billing_type':  r.get('billing_type') or '',
            'fmt_pipe':      fmt_m(r.get('open_pipe')) if r.get('open_pipe') else '',
            'health_risk':   r.get('health_risk_score'),
            'data_quality':  r.get('data_quality_score'),
            'bookings':      r['bookings'] or 0,
            'billings':      r['billings'] or 0,
            'far':           r['far'] or 0,
            'overdue_inv':   r['overdue_inv'] or 0,
            'rr_revenue':    r['rr_revenue'] or 0,
            'bid_margin_raw':   r['bid_margin'],
            'close_margin_raw': r['close_margin'],
            'rules':         viol_codes,
            'baselines':     bl,
            'pulse_id':      r.get('pulse_id', ''),
            'summary':       ' '.join((r.get('overall_summary') or '').split()),
            'leadership':    ' '.join((r.get('leadership_notes') or '').split())[:400],
            'swe_co':        r['swe_co'],
            'high_watch':    bool(r.get('high_watch')),
            'has_pulse':     r.get('has_pulse', False),
            'pulse_trend':   r.get('trend') or '',
            'pulse_updated': (r.get('last_updated') or '')[:10],
            'pulse_scope':   r.get('scope_s') or '',
            'pulse_sched':   r.get('sched_s') or '',
            'pulse_budget':  r.get('budget_s') or '',
            'pulse_resource':r.get('resource_s') or '',
            'pulse_customer':r.get('customer_s') or '',
            'pulse_action':  ' '.join((r.get('action_needed') or '').split())[:300],
            'pulse_steerco': r['steerco_date'].isoformat() if r.get('steerco_date') else '',
            'pulse_golive':  r['next_golive'].isoformat() if r.get('next_golive') else '',
            'slack_intel':   r.get('slack_intel') or '',
            'start_dt':      r['start_dt'].isoformat() if r.get('start_dt') else '',
            'end_dt':        r['end_dt'].isoformat() if r.get('end_dt') else '',
            'stage':         r.get('stage') or '',
            'practice':      r.get('practice') or '',
            'unsch_backlog': r.get('unsch_backlog') or 0,
            'actuals_rem':   r.get('actuals_rem') or 0,
            'eva_amt':       r.get('eva_amt'),
            'eva_pct':       r.get('eva_pct'),
            'csat_score':    r.get('csat_score'),
            'fmt_unsch':     fmt_m(r.get('unsch_backlog')) if r.get('unsch_backlog') else '',
            'fmt_actuals':   fmt_m(r.get('actuals_rem')) if r.get('actuals_rem') else '',
            'fmt_eva_amt':   (f"+${r['eva_amt']:,.0f}" if r.get('eva_amt') is not None and r['eva_amt'] >= 0 else f"-${abs(r['eva_amt']):,.0f}" if r.get('eva_amt') is not None else ''),
            'fmt_eva_pct':   (f"{r['eva_pct']:+.1f}%" if r.get('eva_pct') is not None else ''),
            'far_reason':    r.get('far_reason') or '',
            'far_subreason': r.get('far_subreason') or '',
            'ptg_budget':    r.get('ptg_budget') or '',
            'ptg_scope':     r.get('ptg_scope') or '',
            'ptg_sched':     r.get('ptg_sched') or '',
            'ptg_resource':  r.get('ptg_resource') or '',
            'ptg_customer':  r.get('ptg_customer') or '',
            'fmt_bookings':  fmt_m(r['bookings']),
            'fmt_billings':  fmt_m(r['billings']),
            'fmt_bid':       fmt_pct(r['bid_margin']),
            'fmt_close':     fmt_pct(r['close_margin']),
            'fmt_delivered': (f'{dd:+.1f}%' if dd is not None else ''),
            'fmt_far':       fmt_m(r['far']),
            'fmt_overdue':   fmt_m(r['overdue_inv']) if r['overdue_inv'] else '',
            'fmt_rr':        fmt_m(r['rr_revenue']) if r['rr_revenue'] else '',
            'gdc_total':     r.get('gdc_total') or 0,
            'gdc_india':     r.get('gdc_india') or 0,
            'gdc_pct':       round(r['gdc_pct'] * 100, 1) if r.get('gdc_pct') is not None else None,
            'gdc_resources': r.get('gdc_resources') or [],
            'ironclad_url':  r.get('ironclad_url') or '',
        })
    json_path = os.path.join(OUTPUT_DIR, f"acc_{REGION_SLUG.lower()}_data.json")
    with open(json_path, 'w', encoding='utf-8') as fh:
        _json.dump({'generated': REPORT_DATE, 'region': REGION_LABEL, 'rows': rows_data}, fh, ensure_ascii=False)
    print(f"✅  JSON saved: {json_path}")

# ── Dispatch ──────────────────────────────────────────────────────────────────
if 'json' in OUTPUT_FORMATS: write_json()
if 'txt'  in OUTPUT_FORMATS: write_txt()
if 'docx' in OUTPUT_FORMATS: write_docx()
if 'pptx' in OUTPUT_FORMATS: write_pptx()
if 'html' in OUTPUT_FORMATS: write_html()
# end region loop

