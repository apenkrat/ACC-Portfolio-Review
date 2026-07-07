#!/usr/bin/env python3
"""
Fill the TMT Bi-Weekly PowerPoint template with live Org62 audit data.
Copies the template, updates every account card in-place.
Usage: python3 generate_template_pptx.py
"""
import subprocess, json, sys, os
from datetime import date, datetime, timedelta
from collections import defaultdict

TEMPLATE = "/Users/apenkrat/Downloads/TMT Bi-Weekly Sector Review - 07022026 .pptx"
OUTPUT_DIR = "/Users/apenkrat/Library/CloudStorage/GoogleDrive-apenkrat@salesforce.com/My Drive/TMT Reports"
TODAY = date.today()
REPORT_DATE = TODAY.isoformat()

SWE_PM_OVERRIDES = [
    'aws - afx salesforce rca - csg',
    'disney parks - ccn - r2',
    'nvidia <> slack',
    'nvidia — slack',
]

# ── Auth ───────────────────────────────────────────────────────────────────────
def get_sf():
    try:
        from simple_salesforce import Salesforce
        result = subprocess.run(
            ['sf', 'org', 'display', '--target-org', 'org62', '--json'],
            capture_output=True, text=True, check=True
        )
        d = json.loads(result.stdout)['result']
        return Salesforce(instance_url=d['instanceUrl'], session_id=d['accessToken'])
    except Exception as e:
        print(f"Auth failed: {e}"); sys.exit(1)

sf = get_sf()

def soql(q, label=''):
    print(f'  {label}...')
    recs, result = [], sf.query(q)
    recs.extend(result['records'])
    while not result['done']:
        result = sf.query_more(result['nextRecordsUrl'], True)
        recs.extend(result['records'])
    def flatten(obj, prefix=''):
        row = {}
        for k, v in obj.items():
            if k == 'attributes': continue
            key = f'{prefix}.{k}' if prefix else k
            row.update(flatten(v, key) if isinstance(v, dict) else {key: v})
        return row
    rows = [flatten(r) for r in recs]
    print(f'    -> {len(rows)} rows')
    return rows

def to_f(v):
    if v is None: return None
    try: return float(str(v).replace(',', '').replace('$', '').strip())
    except: return None

def to_d(v):
    if not v: return None
    if isinstance(v, (date, datetime)): return v.date() if isinstance(v, datetime) else v
    for fmt in ['%Y-%m-%d', '%m/%d/%Y']:
        try: return datetime.strptime(str(v)[:10], fmt).date()
        except: pass
    return None

# ── Pull ──────────────────────────────────────────────────────────────────────
print("Pulling live data from Org62...")
q1_rows = soql("""
SELECT Id, Name, pse__Stage__c, pse__Account__c, pse__Account__r.Name,
  pse__Project_Manager__r.Name, ProjectManager2Contact__r.Name,
  pse__Billing_Type__c, pse__Start_Date__c, pse__End_Date__c,
  pse__Bookings__c, pse__Billings__c, Total_Amount_Remaining__c,
  Revenue_Recognized_Comments__c, Forecasted_Amount_Remaining_Subreason__c,
  Forecasted_Amount_Remaining_Details__c,
  Overall_Bid_Margin_new__c, Margin_at_Close_Percent__c,
  Percent_Complete__c, Revenue_Treatment__c
FROM pse__Proj__c
WHERE pse__Stage__c IN ('In Progress', 'In Progress - SWE', 'On Hold')
  AND Subregion_new__c IN ('AMER TMT - 1', 'AMER TMT - 2')
  AND pse__Practice__r.Name != 'FDE'
  AND pse__Account__r.Name != 'Salesforce'
ORDER BY pse__Account__r.Name, Name
""", 'Q1 Projects')

q2_rows = soql("""
SELECT Project__c, Overall_Health__c, Trend_new__c, High_Watch_Visibility__c,
  Scope_Status__c, Schedule_Status__c, Budget_Status__c,
  Resource_Status__c, Customer_Status__c,
  Reason_for_RY_Path_to_Green_Scope__c, Reason_for_RY_Path_to_Green_Schedule__c,
  Reason_for_RY_Path_to_Green_Budget__c, Reason_for_RY_Path_to_Green_Resource__c,
  Reason_for_RY_Path_to_Green_Customer__c,
  Summary_Status__c, Leadership_Notes__c, Action_Needed_from_Leadership__c,
  Next_Steering_Committee_Date__c, Next_Go_Live_Date__c, SWE_or_CO_anticipated__c
FROM Project_Health_Check__c
WHERE Not_Primary_Pulse_Record__c = false
  AND Project__r.pse__Stage__c IN ('In Progress', 'In Progress - SWE', 'On Hold')
  AND Project__r.Subregion_new__c IN ('AMER TMT - 1', 'AMER TMT - 2')
  AND Project__r.pse__Practice__r.Name != 'FDE'
  AND Project__r.pse__Account__r.Name != 'Salesforce'
""", 'Q2 Pulse')

q4_rows = soql("""
SELECT pse__Project__c, COUNT(Id), SUM(pse__Request_Billable_Amount__c), MIN(pse__Start_Date__c)
FROM pse__Resource_Request__c
WHERE pse__Status__c IN ('Draft', 'Ready to Staff', 'Tentative', 'Hold')
  AND pse__Project__r.pse__Stage__c IN ('In Progress', 'In Progress - SWE', 'On Hold')
  AND pse__Project__r.Subregion_new__c IN ('AMER TMT - 1', 'AMER TMT - 2')
  AND pse__Project__r.pse__Practice__r.Name != 'FDE'
GROUP BY pse__Project__c LIMIT 2000
""", 'Q4 Open RRs')

# ── Build maps ────────────────────────────────────────────────────────────────
pulse_map = {r['Project__c']: r for r in q2_rows if r.get('Project__c')}
rr_map = {}
for r in q4_rows:
    pid = r.get('pse__Project__c', '')
    if not pid: continue
    rr_map[pid] = {
        'count': int(to_f(r.get('expr0')) or 0),
        'revenue': to_f(r.get('expr1')) or 0,
        'earliest': str(r.get('expr2',''))[:10],
    }

def is_swe(name, stage, swe_field):
    nl = name.lower()
    if any(o in nl for o in SWE_PM_OVERRIDES): return True
    if stage == 'In Progress - SWE': return True
    if swe_field and swe_field not in ('', 'No', None): return True
    if 'swe' in nl: return True
    return False

def rag(v): return str(v).strip().lower().capitalize() if v else ''

# ── Evaluate projects ─────────────────────────────────────────────────────────
print("Evaluating projects...")
results = []
for p in q1_rows:
    pid      = p.get('Id', '')
    name     = p.get('Name', '') or ''
    acct     = p.get('pse__Account__r.Name', '') or ''
    pm       = p.get('pse__Project_Manager__r.Name', '') or ''
    stage    = p.get('pse__Stage__c', '') or ''
    bookings = to_f(p.get('pse__Bookings__c'))
    billings = to_f(p.get('pse__Billings__c'))
    far      = to_f(p.get('Total_Amount_Remaining__c'))
    bid_m    = to_f(p.get('Overall_Bid_Margin_new__c'))
    close_m  = to_f(p.get('Margin_at_Close_Percent__c'))
    rev_treat= p.get('Revenue_Treatment__c', '') or ''
    end_dt   = to_d(p.get('pse__End_Date__c'))

    pulse     = pulse_map.get(pid, {})
    health    = pulse.get('Overall_Health__c') or 'Null'
    trend_raw = pulse.get('Trend_new__c') or ''
    swe_field = pulse.get('SWE_or_CO_anticipated__c') or ''
    resource_s= rag(pulse.get('Resource_Status__c'))
    summary   = ' '.join((pulse.get('Summary_Status__c') or '').split())
    lead_notes= ' '.join((pulse.get('Leadership_Notes__c') or '').split())
    action    = ' '.join((pulse.get('Action_Needed_from_Leadership__c') or '').split())
    scope_s   = rag(pulse.get('Scope_Status__c'))
    sched_s   = rag(pulse.get('Schedule_Status__c'))
    budget_s  = rag(pulse.get('Budget_Status__c'))
    customer_s= rag(pulse.get('Customer_Status__c'))

    ptg_parts = [pulse.get(f'Reason_for_RY_Path_to_Green_{d}__c') or ''
                 for d in ['Scope','Schedule','Budget','Resource','Customer']]
    ptg = ' '.join(p for p in ptg_parts if p)

    swe_co = is_swe(name, stage, swe_field)
    baseline_ry = any(s in ('Red','Yellow') for s in [scope_s, sched_s, budget_s, resource_s, customer_s])

    violations = []
    if not swe_co and bid_m is not None and close_m is not None:
        delta = close_m - bid_m
        if delta < -5:   violations.append('1A_RED')
        elif delta < 0:  violations.append('1A_YELLOW')
    if far is not None and far < -1:
        violations.append('1C_RED_NEG')
    rr = rr_map.get(pid, {})
    rr_rev = rr.get('revenue', 0)
    if rr_rev > 0: violations.append('1C_RR_RISK')

    fin_reds = {'1A_RED', '1C_RED_NEG', '1C_RED_UNDERUTIL', '1E_OVERDUE_INV', '1F_HOT'}
    is_watermelon = health.lower() == 'green' and (baseline_ry or bool(fin_reds & set(violations)))

    results.append({
        'pid': pid, 'name': name, 'acct': acct, 'pm': pm, 'stage': stage,
        'bookings': bookings, 'billings': billings, 'far': far,
        'bid_m': bid_m, 'close_m': close_m, 'swe_co': swe_co,
        'health': health, 'trend_raw': trend_raw, 'is_watermelon': is_watermelon,
        'violations': violations, 'summary': summary, 'lead_notes': lead_notes,
        'action': action, 'ptg': ptg, 'resource_s': resource_s,
        'scope_s': scope_s, 'sched_s': sched_s, 'budget_s': budget_s, 'customer_s': customer_s,
        'rr_count': rr.get('count',0), 'rr_rev': rr_rev,
        'rev_treat': rev_treat, 'end_dt': end_dt,
    })

# ── Account-level aggregation ─────────────────────────────────────────────────
by_acct = defaultdict(list)
for r in results:
    by_acct[r['acct']].append(r)

def acct_summary(projs):
    total_bk  = sum(r['bookings'] or 0 for r in projs)
    total_bil = sum(r['billings'] or 0 for r in projs)
    total_rr  = sum(r['rr_rev'] for r in projs)
    total_rr_count = sum(r['rr_count'] for r in projs)

    mp = [r for r in projs if not r['swe_co'] and (r['bookings'] or 0) > 0
          and r['bid_m'] is not None and r['close_m'] is not None]
    mbk = sum(r['bookings'] for r in mp)
    w_bid   = sum(r['bid_m']   * r['bookings'] for r in mp) / mbk if mbk else None
    w_close = sum(r['close_m'] * r['bookings'] for r in mp) / mbk if mbk else None

    # Health: red > yellow > watermelon > green > null
    healths = [r['health'].lower() for r in projs]
    if 'red' in healths and not all(r['is_watermelon'] for r in projs if r['health'].lower()=='red'):
        health = 'Red'
    elif 'yellow' in healths:
        health = 'Yellow'
    elif any(r['is_watermelon'] for r in projs):
        health = 'Green'  # watermelons self-report green
    elif 'green' in healths:
        health = 'Green'
    else:
        health = 'Null'

    health_icon = {'Red':'🔴','Yellow':'🟡','Green':'🟢','Null':'⚫'}.get(health,'⚫')

    # Trend: down > stable > up
    trends = [r['trend_raw'].lower() for r in projs if r['trend_raw']]
    if any('down' in t for t in trends):
        trend_arrow = '↓'
    elif any('up' in t for t in trends):
        trend_arrow = '↑'
    else:
        trend_arrow = '↔'

    # Primary project = largest bookings with a pulse
    pulsed = [r for r in projs if r['summary']]
    primary = max(pulsed, key=lambda r: r['bookings'] or 0) if pulsed else max(projs, key=lambda r: r['bookings'] or 0)

    # Status snapshot
    status_snap = (primary['summary'] or primary['lead_notes'] or '')[:300]

    # Path to green — from red/yellow projects
    ry_projs = [r for r in projs if r['health'].lower() in ('red','yellow')]
    ptg_text = ''
    if ry_projs:
        ptg_text = (ry_projs[0]['ptg'] or ry_projs[0]['lead_notes'] or '')[:250]

    # Margin reason
    all_viols = [v for r in projs for v in r['violations']]
    if '1A_RED' in all_viols:
        margin_reason = 'Margin compression > -5% (bid vs. close)'
    elif '1A_YELLOW' in all_viols:
        margin_reason = 'Minor margin compression (< 0%)'
    elif any(r['swe_co'] for r in projs):
        margin_reason = 'SWE/Investment costs excluded from bid'
    elif w_close and w_bid and w_close > w_bid:
        margin_reason = 'Above bid — on target'
    else:
        margin_reason = 'On target'

    return {
        'health': health, 'health_icon': health_icon, 'trend_arrow': trend_arrow,
        'total_bk': total_bk, 'total_bil': total_bil,
        'w_bid': w_bid, 'w_close': w_close,
        'status_snap': status_snap, 'ptg_text': ptg_text, 'margin_reason': margin_reason,
        'total_rr': total_rr, 'total_rr_count': total_rr_count,
    }

acct_data = {acct: acct_summary(projs) for acct, projs in by_acct.items()}

# ── Account name mapping: template display name → SF account name fragments ───
# Each entry: template_name -> list of substrings (any match = include)
ACCT_MAP = {
    # Tier 1
    'Amazon':         ['Amazon Web Services', 'Amazon.com', 'Amazon Devices', 'Amazon - Amazon', 'Amazon NA', 'Amazon Retail'],
    'Apple':          ['Apple Inc.'],
    'AT&T':           ['AT&T Services', 'AT&T - Agentforce', "TERM 1 -6 - AT&T"],
    'Bell Canada':    ['Bell Canada'],
    'Cisco':          ['Cisco Systems'],
    'Crowdstrike':    ['CrowdStrike', 'Crowdstrike'],
    'DirecTV':        ['DTV ', 'DirecTV', 'DIRECTV', 'Betfair Interactive US LLC'],
    'Disney':         ['Disney Parks', 'Disney Worldwide', 'DX - ShopDisney'],
    'Lumen/ Qwest':   ['Qwest Corporation'],
    'Siemens':        ['Siemens Schweiz', 'Siemens SI', 'Siemens Digital'],
    'Sirius XM':      ['Sirius XM', 'SiriusXM'],
    'TDS':            ['TDS Telecommunications'],
    # Tier 2
    'Betfair (FanDuel)': ['Betfair Interactive'],
    'Bloomberg':      ['Bloomberg'],
    'CDK Global':     ['CDK Global', 'CDK '],
    'Comcast':        ['Comcast'],
    'Cornerstone':    ['Cornerstone'],
    'Dell':           ['Dell USA'],
    'Equinix':        ['Equinix'],
    'Fortra':         ['Fortra'],
    'Google':         ['Google LLC', 'Google Fiber', 'Google - ', 'Google Cloud'],
    'Meta':           ['Meta Platforms'],
    'NetApp':         ['NetApp'],
    'Nexstar':        ['Nexstar'],
    'NVIDIA':         ['NVIDIA'],
    'PayPal':         ['PayPal'],
    'Procore':        ['Procore'],
    'ScanSource':     ['ScanSource'],
    'Sony':           ['Sony Interactive'],
    'Spotify':        ['Spotify'],
    'Stansberry':     ['Stansberry', 'Marketwise'],
    'T-Mobile':       ['T-Mobile'],
    'Teradata':       ['Teradata'],
    'The Nielsen Company': ['Nielsen', 'The Nielsen'],
    'UST':            ['UST Global'],
    'Veeam':          ['Veeam'],
    'Verizon':        ['Verizon'],
    'Workday':        ['Workday'],
}

def get_acct_data_for(display_name):
    fragments = ACCT_MAP.get(display_name, [display_name])
    # Collect all projects whose account name matches any fragment
    matched_projs = []
    for acct_name, projs in by_acct.items():
        if any(frag.lower() in acct_name.lower() for frag in fragments):
            matched_projs.extend(projs)
    if not matched_projs:
        return None
    return acct_summary(matched_projs)

# ── PPTX update helpers ───────────────────────────────────────────────────────
def set_run_text(run, text):
    run.text = text

def update_card_content(content_box, data):
    """Update the content textbox of a card in-place with live data."""
    paras = content_box.text_frame.paragraphs
    if len(paras) < 8:
        return

    # Para 0: Health/Trend: [icon] [Health] [arrow]
    p0 = paras[0]
    runs0 = p0.runs
    health_str = f'{data["health_icon"]} {data["health"]} {data["trend_arrow"]}'
    # Find the run with the health/trend value (last meaningful run)
    if len(runs0) >= 3:
        runs0[-1].text = health_str
        for r in runs0[1:-1]:
            r.text = ' '
    elif len(runs0) == 2:
        runs0[-1].text = health_str
    elif len(runs0) == 1:
        # Label + value in one run — keep label
        if ':' in runs0[0].text:
            runs0[0].text = 'Health/Trend: ' + health_str
        else:
            runs0[0].text = health_str

    # Para 1: Status Snapshot
    p1 = paras[1]
    runs1 = p1.runs
    snap = data['status_snap'][:280] if data['status_snap'] else 'Status current as of ' + REPORT_DATE
    if len(runs1) >= 2:
        runs1[-1].text = snap
        for r in runs1[1:-1]: r.text = ''
    elif runs1:
        # Preserve label if present
        if 'Snapshot' in runs1[0].text or len(runs1[0].text) < 25:
            runs1[0].text = 'Status Snapshot: ' + snap
        else:
            runs1[0].text = snap

    # Para 2: Path to Green
    p2 = paras[2]
    runs2 = p2.runs
    ptg = data['ptg_text'][:230] if data['ptg_text'] else 'No red/yellow projects — sustain current delivery.'
    if len(runs2) >= 2:
        runs2[-1].text = ptg
        for r in runs2[1:-1]: r.text = ''
    elif runs2:
        if 'Green' in runs2[0].text or len(runs2[0].text) < 20:
            runs2[0].text = 'Path to Green: ' + ptg
        else:
            runs2[0].text = ptg

    # Para 3: Financials: — label only, leave as-is

    # Para 4: Bookings
    bk_str = f'${data["total_bk"]/1e6:.1f}M' if data['total_bk'] >= 1e6 else f'${data["total_bk"]/1e3:.0f}K'
    if paras[4].runs:
        paras[4].runs[0].text = f'Bookings: {bk_str}'

    # Para 5: Bid Margin
    if paras[5].runs:
        bid_str = f'{data["w_bid"]:.1f}%' if data['w_bid'] is not None else 'N/A'
        paras[5].runs[0].text = f'Bid Margin: {bid_str}'

    # Para 6: Delivered (EAC)
    if paras[6].runs:
        del_str = f'{data["w_close"]:.1f}%' if data['w_close'] is not None else 'N/A'
        paras[6].runs[0].text = f'Delivered (EAC): {del_str}'

    # Para 7: Margin Reason
    if paras[7].runs:
        paras[7].runs[0].text = data['margin_reason']

    # Paras 8+ (Defined Customer Outcomes, Growth) — leave as-is (already in template)

# ── Also update scorecard stats on slide 2 ───────────────────────────────────
def replace_text_in_shape(shape, old, new):
    if not shape.has_text_frame: return
    for para in shape.text_frame.paragraphs:
        full = ''.join(r.text for r in para.runs)
        if old in full:
            replacement = full.replace(old, new)
            if para.runs:
                para.runs[0].text = replacement
                for r in para.runs[1:]: r.text = ''

# ── Load template and fill ────────────────────────────────────────────────────
from pptx import Presentation

print("\nLoading template...")
prs = Presentation(TEMPLATE)

# ── Slide 2: Portfolio Health scorecard ──────────────────────────────────────
slide2 = prs.slides[1]

total_bk  = sum(r['bookings'] or 0 for r in results)
total_bil = sum(r['billings'] or 0 for r in results)
watermelons = [r for r in results if r['is_watermelon']]
reds        = [r for r in results if r['health'].lower()=='red' and not r['is_watermelon']]
yellows     = [r for r in results if r['health'].lower()=='yellow']
clean_green = [r for r in results if r['health'].lower()=='green' and not r['is_watermelon']]
no_pulse    = [r for r in results if not pulse_map.get(r['pid'])]

mp = [r for r in results if not r['swe_co'] and (r['bookings'] or 0)>0
      and r['bid_m'] is not None and r['close_m'] is not None]
mbk = sum(r['bookings'] for r in mp)
w_bid   = sum(r['bid_m']   * r['bookings'] for r in mp) / mbk if mbk else 0
w_close = sum(r['close_m'] * r['bookings'] for r in mp) / mbk if mbk else 0

SCORECARD_REPLACEMENTS = {
    '168': str(len(results)),
    '20.9%': f'{w_close:.1f}%',
    '9.3%':  f'{w_bid:.1f}%',
    'Yellow Risk Highlights (21 Projects)': f'Yellow Risk Highlights ({len(yellows)} Projects)',
    'Green Highlights (69 Projects)': f'Green Highlights ({len(clean_green) + len(watermelons)} Projects)',
    'Green': 'Green',  # keep as-is placeholder
}

for shape in slide2.shapes:
    for old, new in [
        ('168', str(len(results))),
        ('20.9%', f'{w_close:.1f}%'),
        ('9.3%', f'{w_bid:.1f}%'),
        ('Yellow Risk Highlights (21 Projects)', f'Yellow Risk Highlights ({len(yellows)} Projects)'),
        ('Green Highlights (69 Projects)', f'Green Highlights ({len(clean_green)+len(watermelons)} Projects)'),
    ]:
        replace_text_in_shape(shape, old, new)

# ── Slides 3 (Region Snapshot) ────────────────────────────────────────────────
slide3 = prs.slides[2]
t1_projs = [r for r in results if sum(r2['bookings'] or 0 for r2 in by_acct[r['acct']]) >= 7_000_000]
t2_projs = [r for r in results if 1_500_000 <= sum(r2['bookings'] or 0 for r2 in by_acct[r['acct']]) < 7_000_000]
t3_projs = [r for r in results if sum(r2['bookings'] or 0 for r2 in by_acct[r['acct']]) < 1_500_000]

def tier_counts(projs):
    return {
        'count': len(projs),
        'bk': sum(r['bookings'] or 0 for r in projs),
        'bil': sum(r['billings'] or 0 for r in projs),
        'red': sum(1 for r in projs if r['health'].lower()=='red' and not r['is_watermelon']),
        'yel': sum(1 for r in projs if r['health'].lower()=='yellow'),
        'grn': sum(1 for r in projs if r['health'].lower()=='green' and not r['is_watermelon']),
    }

t1 = tier_counts(t1_projs)
t2 = tier_counts(t2_projs)
t3 = tier_counts(t3_projs)

S3_REPL = [
    ('Project Count: 47',  f'Project Count: {t1["count"]}'),
    ('Bookings: $304.6M',  f'Bookings: ${t1["bk"]/1e6:.1f}M'),
    ('Billings: $183.8M',  f'Billings: ${t1["bil"]/1e6:.1f}M'),
    ('🔴 Red: 2\n',         f'🔴 Red: {t1["red"]}\n'),
    ('🟡 Yellow: 10\n',     f'🟡 Yellow: {t1["yel"]}\n'),
    ('🟢 Green: 22\n',      f'🟢 Green: {t1["grn"]}\n'),
    ('Project Count: 51',  f'Project Count: {t2["count"]}'),
    ('Bookings: $51.4M',   f'Bookings: ${t2["bk"]/1e6:.1f}M'),
    ('Billings: $22.8M',   f'Billings: ${t2["bil"]/1e6:.1f}M'),
    ('Project Count: 67',  f'Project Count: {t3["count"]}'),
    ('Bookings: $7.2M',    f'Bookings: ${t3["bk"]/1e6:.1f}M'),
    ('Billings: $4.3M',    f'Billings: ${t3["bil"]/1e6:.1f}M'),
]
for shape in slide3.shapes:
    if shape.shape_type == 6:
        for child in shape.shapes:
            for old, new in S3_REPL:
                replace_text_in_shape(child, old, new)
    else:
        for old, new in S3_REPL:
            replace_text_in_shape(shape, old, new)

# ── Update account cards (slides 7-23) ───────────────────────────────────────
print("Updating account cards...")
cards_updated = 0

for slide_idx in range(6, min(24, len(prs.slides))):
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.shape_type != 6:
            continue
        # Find title box (short text, no Health/Trend) and content box
        title_box = None
        content_box = None
        for child in shape.shapes:
            if child.has_text_frame:
                txt = child.text_frame.text.strip()
                if 'Health/Trend' in txt:
                    content_box = child
                elif txt and len(txt) < 50 and '\n' not in txt and 'Health' not in txt:
                    title_box = child
        if not title_box or not content_box:
            continue

        display_name = title_box.text_frame.text.strip()
        data = get_acct_data_for(display_name)
        if not data:
            print(f'  ⚠ No SF data for card: {display_name!r}')
            continue

        update_card_content(content_box, data)
        cards_updated += 1
        print(f'  ✓ {display_name} → {data["health_icon"]} {data["health"]} {data["trend_arrow"]}  ${data["total_bk"]/1e6:.1f}M')

print(f'\nUpdated {cards_updated} cards.')

# ── Save ──────────────────────────────────────────────────────────────────────
os.makedirs(OUTPUT_DIR, exist_ok=True)
output_path = f"{OUTPUT_DIR}/TMT_Bi-Weekly_{REPORT_DATE}.pptx"
prs.save(output_path)
print(f'\n✅  Saved: {output_path}')
